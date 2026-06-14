"""COM-022 — PowerPoint export service.

Generates a .pptx file from a Solution Architecture Blueprint:
  - Cover slide: solution title, architect, date, overall completeness
  - One slide per blueprint section: title, narrative text, completeness badge
"""

import io
import logging
from datetime import datetime
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt

logger = logging.getLogger(__name__)

# Brand colours
_BLUE = RGBColor(0x1E, 0x40, 0xAF)
_NEAR_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GREEN = RGBColor(0x16, 0xA3, 0x4A)
_AMBER = RGBColor(0xD9, 0x77, 0x06)
_RED = RGBColor(0xDC, 0x26, 0x26)
_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)

# Blueprint section ordered list (matches BlueprintCompletenessService.SECTION_TITLES)
SECTION_TITLES = {
    "executive_summary": "Executive Summary",
    "vision_motivation": "Vision & Motivation",
    "value_stream_map": "Value Stream Map",
    "business_process_view": "Business Process View",
    "application_cooperation": "Application Co-operation",
    "data_information": "Data & Information",
    "deployment_view": "Deployment View",
    "network_communication": "Network & Communication",
    "gap_analysis": "Gap Analysis",
    "transition_roadmap": "Transition Roadmap",
    "work_packages": "Work Packages",
    "security_viewpoint": "Security Viewpoint",
    "nfr_satisfaction": "NFR Satisfaction",
    "requirements_traceability": "Requirements Traceability",
    "erp_fit_gap": "ERP Fit-Gap Analysis",
    "integration_architecture": "Integration Architecture",
}

_MAX_NARRATIVE_WORDS = 200


def _truncate_to_words(text: str, max_words: int) -> str:
    """Return text truncated to at most max_words words."""
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]) + " …"


def _num_score(score):
    """Coerce a score (sometimes a dict of metrics) to a number, or None."""
    if isinstance(score, dict):
        for _k in ("completeness", "overall_score", "score", "value", "percentage"):
            _v = score.get(_k)
            if isinstance(_v, (int, float)):
                return _v
        return None
    return score if isinstance(score, (int, float)) else None


def _badge_colour(score: Optional[float]) -> RGBColor:
    score = _num_score(score)
    if score is None:
        return _AMBER
    if score >= 80:
        return _GREEN
    if score >= 60:
        return _AMBER
    return _RED


def _badge_label(score: Optional[float]) -> str:
    score = _num_score(score)
    if score is None:
        return "N/A"
    return f"{int(round(score))}%"


class PowerPointExportService:
    """Generates a PPTX export for a solution's architecture blueprint."""

    @staticmethod
    def generate(solution_id: int) -> bytes:
        """Build a PPTX document and return its raw bytes.

        Raises:
            ValueError: if solution_id is not found.
        """
        from app.models.solution_models import Solution
        from app.models.user import User

        solution = Solution.query.get(solution_id)
        if solution is None:
            raise ValueError(f"Solution {solution_id} not found")

        architect_name = "Unknown"
        if solution.created_by_id:
            user = User.query.get(solution.created_by_id)
            if user:
                architect_name = (user.full_name() if (user and callable(getattr(user, "full_name", None))) else None) or getattr(user, "name", None) or user.email or "Unknown"

        section_narratives: dict = solution.section_narratives or {}
        section_scores: dict = solution.section_scores or {}

        # Compute overall completeness as average of available section scores
        score_values = [v for v in section_scores.values() if isinstance(v, (int, float))]
        overall_score: Optional[float] = (sum(score_values) / len(score_values)) if score_values else None

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        PowerPointExportService._add_cover_slide(
            prs,
            title=solution.name or "Untitled Solution",
            architect=architect_name,
            created_at=solution.created_at,
            overall_score=overall_score,
        )

        for section_key, section_title in SECTION_TITLES.items():
            narrative = section_narratives.get(section_key, "")
            score = section_scores.get(section_key)
            PowerPointExportService._add_section_slide(
                prs,
                section_title=section_title,
                narrative=narrative or "",
                score=score,
            )

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------
    # Private slide builders
    # ------------------------------------------------------------------

    @staticmethod
    def _add_cover_slide(
        prs: Presentation,
        title: str,
        architect: str,
        created_at: Optional[datetime],
        overall_score: Optional[float],
    ) -> None:
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Full-bleed background rectangle
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = _BLUE
        bg.line.fill.background()

        # Title text box
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(11.0), Inches(1.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = _WHITE

        # Subtitle line: architect + date
        date_str = created_at.strftime("%d %b %Y") if created_at else "—"
        sub_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.4), Inches(11.0), Inches(0.6)
        )
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        srun = sp.add_run()
        srun.text = f"Architect: {architect}   |   Date: {date_str}"
        srun.font.size = Pt(18)
        srun.font.color.rgb = _NEAR_WHITE

        # Completeness badge
        badge_colour = _badge_colour(overall_score)
        badge_label = _badge_label(overall_score)

        badge = slide.shapes.add_shape(
            1,  # RECTANGLE
            Inches(0.8), Inches(4.4), Inches(2.4), Inches(0.8)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_colour
        badge.line.fill.background()
        btf = badge.text_frame
        btf.margin_top = Pt(4)
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        brun = bp.add_run()
        brun.text = f"Completeness: {badge_label}"
        brun.font.size = Pt(16)
        brun.font.bold = True
        brun.font.color.rgb = _WHITE

        # Footer label
        label_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.8), Inches(11.0), Inches(0.4)
        )
        ltf = label_box.text_frame
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.LEFT
        lrun = lp.add_run()
        lrun.text = "Solution Architecture Blueprint — A.R.C.H.I.E. Platform"
        lrun.font.size = Pt(11)
        lrun.font.color.rgb = _NEAR_WHITE

    @staticmethod
    def _add_section_slide(
        prs: Presentation,
        section_title: str,
        narrative: str,
        score: Optional[float],
    ) -> None:
        blank_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(blank_layout)

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Header strip (top bar)
        header = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0), slide_w, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _BLUE
        header.line.fill.background()

        # Section title inside header
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15), Inches(10.5), Inches(0.9)
        )
        ttf = title_box.text_frame
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        trun = tp.add_run()
        trun.text = section_title
        trun.font.size = Pt(26)
        trun.font.bold = True
        trun.font.color.rgb = _WHITE

        # Completeness badge (top-right corner)
        badge_colour = _badge_colour(score)
        badge_label = _badge_label(score)

        badge = slide.shapes.add_shape(
            1,
            Inches(11.2), Inches(0.2), Inches(1.8), Inches(0.7)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_colour
        badge.line.fill.background()
        btf = badge.text_frame
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        brun = bp.add_run()
        brun.text = badge_label
        brun.font.size = Pt(14)
        brun.font.bold = True
        brun.font.color.rgb = _WHITE

        # Narrative content area
        body_top = Inches(1.4)
        body_box = slide.shapes.add_textbox(
            Inches(0.5), body_top, Inches(12.3), slide_h - body_top - Inches(0.5)
        )
        btf2 = body_box.text_frame
        btf2.word_wrap = True

        display_text = _truncate_to_words(narrative, _MAX_NARRATIVE_WORDS) if narrative else "(No narrative recorded for this section.)"

        para = btf2.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = display_text
        run.font.size = Pt(14)
        run.font.color.rgb = _DARK_TEXT
