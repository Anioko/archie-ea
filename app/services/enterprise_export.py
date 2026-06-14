"""Enterprise Export Service - PDF, PowerPoint, Excel, Word, ArchiMate XML

Provides Fortune 500-grade export formats for architecture artifacts.

Exports:
- PDF: Solution Architecture Document (SAD) with diagrams
- PowerPoint: Executive presentation deck
- Excel: TCO models, cost analysis, portfolio data
- Word: Governance submission templates
- ArchiMate XML: Open Exchange Format (Sparx EA/BiZZdesign compatible)

Usage:
    from app.services.enterprise_export import EnterpriseExportService
    
    # PDF export
    pdf_bytes = EnterpriseExportService.export_solution_pdf(solution_id)
    
    # PowerPoint export
    pptx_bytes = EnterpriseExportService.export_solution_pptx(solution_id)
    
    # ArchiMate XML
    xml_str = EnterpriseExportService.export_archimate_xml(element_ids=[])
"""

from io import BytesIO
from typing import List, Optional
import logging

logger = logging.getLogger(__name__)


class EnterpriseExportService:
    """Enterprise-grade export formats for architecture artifacts."""
    
    @staticmethod
    def export_solution_pdf(solution_id: int, include_diagrams: bool = True) -> bytes:
        """Export Solution Architecture Document as PDF.
        
        Includes:
        - Executive summary
        - TOGAF ADM phases (A-H)
        - ArchiMate diagrams
        - Cost models, risks, metrics
        - ARB approval status
        
        Uses: weasyprint or wkhtmltopdf for HTML→PDF conversion
        """
        from app.models import Solution
        from flask import render_template

        from app.services.sad_document_builder import SADDocumentBuilder

        solution = Solution.query.get_or_404(solution_id)

        # OPS-003: full partner-grade SAD — cover page, document control, TOC,
        # executive summary, all 8 TOGAF ADM phases, risk register, appendix.
        # All content gathered from live data via SADDocumentBuilder; the old
        # template rendered 4 empty tables because the route never passed the
        # section variables.
        context = SADDocumentBuilder.build(solution)
        context["include_diagrams"] = include_diagrams
        context["for_pdf"] = True
        html = render_template(
            "exports/solution_architecture_document.html", **context
        )

        try:
            # Try weasyprint first (better CSS support; the template carries its
            # own @page rules, cover gradient, and table styling).
            from weasyprint import HTML
            return HTML(string=html).write_pdf()
            
        except ImportError:
            # Fallback to pdfkit (wkhtmltopdf)
            try:
                import pdfkit
                options = {
                    'page-size': 'A4',
                    'margin-top': '20mm',
                    'margin-right': '20mm',
                    'margin-bottom': '20mm',
                    'margin-left': '20mm',
                    'encoding': 'UTF-8',
                    'no-outline': None,
                    'enable-local-file-access': None
                }
                pdf_bytes = pdfkit.from_string(html, False, options=options)
                return pdf_bytes
            except Exception as e:
                logger.error(f"PDF export failed: {e}")
                raise ValueError("PDF export libraries not available. Install weasyprint or wkhtmltopdf.")
    
    @staticmethod
    def export_solution_pptx(solution_id: int) -> bytes:
        """Export Solution as PowerPoint presentation.
        
        Slides:
        1. Title slide (solution name, architect, date)
        2. Executive summary
        3. Business drivers & goals
        4. Capability map
        5. Application landscape
        6. Technology architecture
        7. Cost breakdown
        8. Risk assessment
        9. Implementation roadmap
        10. Recommendation
        
        Uses: python-pptx
        """
        from app.models import Solution
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        solution = Solution.query.get_or_404(solution_id)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = solution.name
        subtitle.text = f"Solution Architecture Document\n{solution.created_at.strftime('%B %d, %Y')}"
        
        # Slide 2: Executive Summary
        bullet_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = "Executive Summary"
        body = slide.placeholders[1].text_frame
        body.text = solution.description or "Solution overview"
        
        # Slide 3-9: Add content slides (abbreviated for initial implementation)
        # TODO: Add capability maps, diagrams, cost tables
        
        # Save to bytes
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io.read()
    
    @staticmethod
    def export_archimate_xml(element_ids: Optional[List[int]] = None) -> str:
        """Export ArchiMate elements as Open Exchange Format XML.
        
        Compatible with:
        - Sparx Enterprise Architect (.xml import)
        - BiZZdesign Enterprise Studio
        - Archi (open source)
        - Any ArchiMate 3.2 compliant tool
        
        Args:
            element_ids: List of element IDs to export. If None, exports all.
        
        Returns:
            ArchiMate XML string (Open Exchange Format)
        """
        from app.models import ArchiMateElement, ArchiMateRelationship
        from datetime import datetime
        
        if element_ids:
            elements = ArchiMateElement.query.filter(ArchiMateElement.id.in_(element_ids)).all()
        else:
            elements = ArchiMateElement.query.all()
        
        # Build XML structure (ArchiMate 3.2 Open Exchange Format)
        xml_lines = [
            '<?xml version="1.0" encoding="UTF-8"?>',
            '<model xmlns="http://www.opengroup.org/xsd/archimate/3.0/"',
            '       xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"',
            '       xsi:schemaLocation="http://www.opengroup.org/xsd/archimate/3.0/ http://www.opengroup.org/xsd/archimate/3_1/archimate3_Diagram.xsd"',
            f'       identifier="archie-export-{datetime.utcnow().strftime("%Y%m%d-%H%M%S")}"',
            '       version="3.2">',
            '  <name>A.R.C.H.I.E. Architecture Export</name>',
            '  <elements>'
        ]
        
        # Export elements
        for elem in elements:
            xml_lines.append(f'    <element identifier="{elem.id}" xsi:type="{elem.layer}:{elem.element_type}">')
            xml_lines.append(f'      <name>{_escape_xml(elem.name)}</name>')
            if elem.description:
                xml_lines.append(f'      <documentation>{_escape_xml(elem.description)}</documentation>')
            xml_lines.append('    </element>')
        
        xml_lines.extend(['  </elements>', '  <relationships>'])
        
        # Export relationships
        if element_ids:
            relationships = ArchiMateRelationship.query.filter(
                ArchiMateRelationship.source_id.in_(element_ids)
            ).all()
        else:
            relationships = ArchiMateRelationship.query.all()
        
        for rel in relationships:
            xml_lines.append(f'    <relationship identifier="{rel.id}" source="{rel.source_id}" target="{rel.target_id}" xsi:type="{rel.relationship_type}">')
            if rel.name:
                xml_lines.append(f'      <name>{_escape_xml(rel.name)}</name>')
            xml_lines.append('    </relationship>')
        
        xml_lines.extend(['  </relationships>', '</model>'])
        
        return '\n'.join(xml_lines)
    
    @staticmethod
    def export_portfolio_excel(filters: Optional[dict] = None) -> bytes:
        """Export application portfolio to Excel with multiple sheets.
        
        Sheets:
        1. Applications - Full inventory
        2. Vendors - Vendor list with risk scores
        3. Costs - TCO analysis
        4. Rationalization - TAME scores
        5. Capabilities - Coverage matrix
        
        Uses: openpyxl
        """
        from app.models import Application, Vendor
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        
        wb = Workbook()
        
        # Sheet 1: Applications
        ws_apps = wb.active
        ws_apps.title = "Applications"
        
        # Header row
        headers = ['Name', 'Vendor', 'Lifecycle Status', 'Annual Cost', 'Technical Owner', 'Business Owner', 'Criticality']
        ws_apps.append(headers)
        
        # Style header
        header_fill = PatternFill(start_color="1E293B", end_color="1E293B", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)
        for cell in ws_apps[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="left", vertical="center")
        
        # Data rows
        apps = Application.query.all()
        for app in apps:
            ws_apps.append([
                app.name,
                app.vendor.name if app.vendor else '',
                app.lifecycle_status or '',
                app.annual_cost or 0,
                app.technical_owner or '',
                app.business_owner or '',
                app.criticality or ''
            ])
        
        # Auto-size columns
        for column in ws_apps.columns:
            max_length = 0
            column = [cell for cell in column]
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(cell.value)
                except Exception as exc:
                    logger.debug("suppressed error in EnterpriseExportService.export_portfolio_excel (app/services/enterprise_export.py): %s", exc)
            adjusted_width = min(max_length + 2, 50)
            ws_apps.column_dimensions[column[0].column_letter].width = adjusted_width
        
        # Save to bytes
        excel_io = BytesIO()
        wb.save(excel_io)
        excel_io.seek(0)
        return excel_io.read()


def _escape_xml(text: str) -> str:
    """Escape special XML characters."""
    if not text:
        return ""
    return (text
            .replace('&', '&amp;')
            .replace('<', '&lt;')
            .replace('>', '&gt;')
            .replace('"', '&quot;')
            .replace("'", '&apos;'))
