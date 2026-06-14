"""
Document Text Extractor for ArchiMate Generation

Extracts text content from various document formats for architecture generation.
Supports: PDF, TXT, MD, DOCX, HTML, CSV, XLSX, XLS
"""

import logging
import os
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


def extract_text_from_file(file_path: str, file_type: str) -> str:
    """
    Extract text content from uploaded document.

    Args:
        file_path: Path to the uploaded file
        file_type: Type of file ('text', 'document', 'image', 'spreadsheet')

    Returns:
        Extracted text content
    """
    try:
        # Spreadsheet files (CSV, XLSX, XLS)
        if file_type == "spreadsheet":
            return extract_from_spreadsheet(file_path)

        # Text files (TXT, MD, HTML)
        if file_type == "text":
            return extract_from_text_file(file_path)

        # Check file extension for spreadsheets even if file_type is different
        lower_path = file_path.lower()
        if lower_path.endswith(".csv"):
            return extract_from_csv(file_path)
        elif lower_path.endswith((".xlsx", ".xls")):
            return extract_from_excel(file_path)

        # PDF documents
        elif file_type == "document" and lower_path.endswith(".pdf"):
            return extract_from_pdf(file_path)

        # Word documents (DOCX)
        elif file_type == "document" and lower_path.endswith(".docx"):
            return extract_from_docx(file_path)

        # DOC files (older Word format)
        elif file_type == "document" and lower_path.endswith(".doc"):
            return extract_from_doc(file_path)

        # PowerPoint files
        elif file_type == "document" and lower_path.endswith((".pptx", ".ppt")):
            return extract_from_pptx(file_path)

        # Fallback: try to read as text
        else:
            logger.warning(f"Unknown file type: {file_type}, attempting text extraction")
            return extract_from_text_file(file_path)

    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return f"Error extracting text from document: {str(e)}"


def extract_from_spreadsheet(file_path: str) -> str:
    """
    Extract text from spreadsheet files (CSV, XLSX, XLS).
    Auto-detects file type based on extension.
    """
    lower_path = file_path.lower()

    if lower_path.endswith(".csv"):
        return extract_from_csv(file_path)
    elif lower_path.endswith(".xlsx"):
        return extract_from_excel(file_path)
    elif lower_path.endswith(".xls"):
        return extract_from_excel_legacy(file_path)
    else:
        return "Error: Unsupported spreadsheet format. Please use CSV, XLSX, or XLS."


def extract_from_csv(file_path: str) -> str:
    """
    Extract structured text from CSV files.
    Preserves column headers and provides summary statistics.
    """
    try:
        import csv

        content_lines = []
        data_rows = []
        headers = []

        # Try different encodings
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]
        file_content = None

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding, newline="") as f:
                    file_content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if file_content is None:
            return "Error: Could not decode CSV file. Please check the file encoding."

        # Parse CSV
        reader = csv.reader(file_content.splitlines())
        rows = list(reader)

        if not rows:
            return "Error: CSV file is empty."

        # Extract headers (first row)
        headers = rows[0] if rows else []
        data_rows = rows[1:] if len(rows) > 1 else []

        # Build structured output
        content_lines.append("=" * 60)
        content_lines.append("CSV DATA EXTRACTION")
        content_lines.append("=" * 60)
        content_lines.append("")

        # Summary statistics
        content_lines.append("SUMMARY:")
        content_lines.append(f"  - Total Columns: {len(headers)}")
        content_lines.append(f"  - Total Data Rows: {len(data_rows)}")
        content_lines.append(f"  - Column Names: {', '.join(headers)}")
        content_lines.append("")

        # Column headers with indices
        content_lines.append("COLUMN STRUCTURE:")
        for i, header in enumerate(headers):
            # Get sample values from this column
            sample_values = [row[i] for row in data_rows[:5] if i < len(row) and row[i].strip()]
            unique_count = len(set(row[i] for row in data_rows if i < len(row)))
            content_lines.append(f"  [{i}] {header}")
            content_lines.append(f"      Unique values: ~{unique_count}")
            if sample_values:
                content_lines.append(f"      Sample: {', '.join(sample_values[:3])}")
        content_lines.append("")

        # Include full data (limited to prevent token overflow)
        content_lines.append("DATA RECORDS:")
        content_lines.append("-" * 40)

        # Format as table-like structure
        max_rows = min(100, len(data_rows))  # Limit to 100 rows
        for i, row in enumerate(data_rows[:max_rows]):
            record_parts = []
            for j, (header, value) in enumerate(zip(headers, row)):
                if value.strip():
                    record_parts.append(f"{header}: {value}")
            content_lines.append(f"Record {i + 1}: {' | '.join(record_parts)}")

        if len(data_rows) > max_rows:
            content_lines.append(f"... and {len(data_rows) - max_rows} more rows")

        content_lines.append("")
        content_lines.append("=" * 60)

        return "\n".join(content_lines)

    except Exception as e:
        logger.error(f"Error extracting CSV: {e}")
        return f"Error extracting CSV content: {str(e)}"


def extract_from_excel(file_path: str) -> str:
    """
    Extract structured text from Excel files (XLSX format).
    Handles multiple sheets and preserves structure.
    """
    try:
        import openpyxl

        workbook = openpyxl.load_workbook(file_path, data_only=True)
        content_lines = []

        content_lines.append("=" * 60)
        content_lines.append("EXCEL DATA EXTRACTION")
        content_lines.append("=" * 60)
        content_lines.append("")

        # Summary
        content_lines.append("WORKBOOK SUMMARY:")
        content_lines.append(f"  - Total Sheets: {len(workbook.sheetnames)}")
        content_lines.append(f"  - Sheet Names: {', '.join(workbook.sheetnames)}")
        content_lines.append("")

        # Process each sheet
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            content_lines.append("-" * 40)
            content_lines.append(f"SHEET: {sheet_name}")
            content_lines.append(f"  Dimensions: {sheet.dimensions}")
            content_lines.append(f"  Max Row: {sheet.max_row}, Max Column: {sheet.max_column}")
            content_lines.append("")

            if sheet.max_row == 0 or sheet.max_column == 0:
                content_lines.append("  (Empty sheet)")
                continue

            # Get headers (first row)
            headers = []
            for col in range(1, min(sheet.max_column + 1, 50)):  # Limit to 50 columns
                cell = sheet.cell(row=1, column=col)
                headers.append(str(cell.value) if cell.value else f"Column{col}")

            content_lines.append(f"  COLUMNS: {', '.join(headers)}")
            content_lines.append("")

            # Column analysis
            content_lines.append("  COLUMN STRUCTURE:")
            for col_idx, header in enumerate(headers, 1):
                # Sample values
                sample_values = []
                for row in range(2, min(7, sheet.max_row + 1)):
                    cell_val = sheet.cell(row=row, column=col_idx).value
                    if cell_val is not None:
                        sample_values.append(str(cell_val)[:50])
                content_lines.append(f"    [{col_idx}] {header}")
                if sample_values:
                    content_lines.append(f"        Sample: {', '.join(sample_values[:3])}")
            content_lines.append("")

            # Data rows (limited)
            content_lines.append("  DATA RECORDS:")
            max_rows = min(50, sheet.max_row - 1)  # Limit rows per sheet
            for row_idx in range(2, max_rows + 2):
                record_parts = []
                for col_idx, header in enumerate(headers, 1):
                    cell_val = sheet.cell(row=row_idx, column=col_idx).value
                    if cell_val is not None:
                        record_parts.append(f"{header}: {str(cell_val)[:100]}")
                if record_parts:
                    content_lines.append(f"    Record {row_idx - 1}: {' | '.join(record_parts)}")

            if sheet.max_row - 1 > max_rows:
                content_lines.append(f"    ... and {sheet.max_row - 1 - max_rows} more rows")

            content_lines.append("")

        workbook.close()
        content_lines.append("=" * 60)

        return "\n".join(content_lines)

    except ImportError:
        logger.error("openpyxl not installed. Install with: pip install openpyxl")
        return "Error: Excel processing library not available. Please install openpyxl or upload a CSV file."
    except Exception as e:
        logger.error(f"Error extracting Excel: {e}")
        # Try fallback: read-only mode without styles
        try:
            import openpyxl

            logger.info(
                f"Attempting fallback Excel extraction (read-only, no styles) for {file_path}"
            )
            workbook = openpyxl.load_workbook(
                file_path, data_only=True, read_only=True, keep_links=False
            )
            content_lines = []
            content_lines.append("=" * 60)
            content_lines.append("EXCEL DATA EXTRACTION (FALLBACK MODE)")
            content_lines.append("=" * 60)
            content_lines.append(
                "Note: File had formatting issues, extracted data in read-only mode"
            )
            content_lines.append("")

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_lines.append(f"SHEET: {sheet_name}")
                # Extract basic data without formatting
                for row in sheet.iter_rows(min_row=1, max_row=100, values_only=True):
                    if any(cell is not None for cell in row):
                        content_lines.append(
                            " | ".join(str(cell)[:50] if cell is not None else "" for cell in row)
                        )
                content_lines.append("")

            workbook.close()
            return "\n".join(content_lines)
        except Exception as fallback_error:
            logger.error(f"Fallback Excel extraction also failed: {fallback_error}")
            return f"Error extracting Excel content: {str(e)}\n\nFallback attempt also failed: {str(fallback_error)}\n\nPlease try:\n1. Re-saving the Excel file\n2. Exporting to CSV format\n3. Checking if the file is corrupted"


def extract_from_excel_legacy(file_path: str) -> str:
    """
    Extract text from older Excel files (XLS format).
    """
    try:
        import xlrd

        workbook = xlrd.open_workbook(file_path)
        content_lines = []

        content_lines.append("=" * 60)
        content_lines.append("EXCEL (XLS) DATA EXTRACTION")
        content_lines.append("=" * 60)
        content_lines.append("")

        content_lines.append("WORKBOOK SUMMARY:")
        content_lines.append(f"  - Total Sheets: {workbook.nsheets}")
        content_lines.append(f"  - Sheet Names: {', '.join(workbook.sheet_names())}")
        content_lines.append("")

        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)

            content_lines.append("-" * 40)
            content_lines.append(f"SHEET: {sheet.name}")
            content_lines.append(f"  Rows: {sheet.nrows}, Columns: {sheet.ncols}")
            content_lines.append("")

            if sheet.nrows == 0 or sheet.ncols == 0:
                content_lines.append("  (Empty sheet)")
                continue

            # Get headers
            headers = [
                str(sheet.cell_value(0, col) or f"Column{col}")
                for col in range(min(sheet.ncols, 50))
            ]
            content_lines.append(f"  COLUMNS: {', '.join(headers)}")
            content_lines.append("")

            # Data rows
            content_lines.append("  DATA RECORDS:")
            max_rows = min(50, sheet.nrows - 1)
            for row_idx in range(1, max_rows + 1):
                record_parts = []
                for col_idx, header in enumerate(headers):
                    cell_val = sheet.cell_value(row_idx, col_idx)
                    if cell_val:
                        record_parts.append(f"{header}: {str(cell_val)[:100]}")
                if record_parts:
                    content_lines.append(f"    Record {row_idx}: {' | '.join(record_parts)}")

            if sheet.nrows - 1 > max_rows:
                content_lines.append(f"    ... and {sheet.nrows - 1 - max_rows} more rows")

            content_lines.append("")

        content_lines.append("=" * 60)
        return "\n".join(content_lines)

    except ImportError:
        logger.error("xlrd not installed. Install with: pip install xlrd")
        return (
            "Error: Legacy Excel processing library not available. Please convert to XLSX or CSV."
        )
    except Exception as e:
        logger.error(f"Error extracting XLS: {e}")
        return f"Error extracting legacy Excel content: {str(e)}"


def extract_from_text_file(file_path: str) -> str:
    """Extract text from plain text files (TXT, MD, HTML, CSV)"""
    encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1"]

    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
        except Exception as e:
            logger.error(f"Error reading text file with {encoding}: {e}")
            continue

    # If all encodings fail, return error message
    return "Error: Could not decode text file. Please ensure it's a valid text document."


def extract_from_pdf(file_path: str) -> str:
    """Extract text from PDF files"""
    try:
        import PyPDF2

        text_content = []
        with open(file_path, "rb") as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text:
                    text_content.append(text)

        extracted_text = "\n\n".join(text_content)

        if not extracted_text.strip():
            return "Error: PDF appears to be empty or contains only images. Please use a text-based PDF."

        return extracted_text

    except ImportError:
        logger.error("PyPDF2 not installed. Install with: pip install PyPDF2")
        return "Error: PDF processing library not available. Please upload a text file instead."
    except Exception as e:
        logger.error(f"Error extracting PDF: {e}")
        return f"Error extracting PDF content: {str(e)}"


def extract_from_docx(file_path: str) -> str:
    """Extract text from DOCX files"""
    try:
        import docx

        doc = docx.Document(file_path)
        text_content = []

        # Extract from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)

        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_content.append(row_text)

        extracted_text = "\n\n".join(text_content)

        if not extracted_text.strip():
            return "Error: Word document appears to be empty."

        return extracted_text

    except ImportError:
        logger.error("python-docx not installed. Install with: pip install python-docx")
        return "Error: Word document processing library not available. Please upload a text or PDF file instead."
    except Exception as e:
        logger.error(f"Error extracting DOCX: {e}")
        return f"Error extracting Word document content: {str(e)}"


def extract_from_doc(file_path: str) -> str:
    """Extract text from older DOC files"""
    try:
        # Try using antiword (if available on system)
        import subprocess

        result = subprocess.run(["antiword", file_path], capture_output=True, text=True, timeout=30)

        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
        else:
            return "Error: Could not extract text from .DOC file. Please convert to DOCX or PDF."

    except FileNotFoundError:
        logger.warning("antiword not installed for DOC extraction")
        return "Error: Old Word (.doc) format not supported. Please convert to DOCX or PDF."
    except subprocess.TimeoutExpired:
        return "Error: Document processing timed out. Please try a smaller file."
    except Exception as e:
        logger.error(f"Error extracting DOC: {e}")
        return f"Error extracting old Word document: {str(e)}"


def extract_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint files"""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"--- Slide {slide_num} ---")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Extract from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_text.append(row_text)

            if len(slide_text) > 1:  # More than just the slide header
                text_content.append("\n".join(slide_text))

        extracted_text = "\n\n".join(text_content)

        if not extracted_text.strip():
            return "Error: PowerPoint appears to be empty or contains only images."

        return extracted_text

    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        return "Error: PowerPoint processing library not available. Please upload a text or PDF file instead."
    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
        return f"Error extracting PowerPoint content: {str(e)}"


def get_document_summary(text: str, max_chars: int = 500) -> str:
    """
    Get a summary/preview of document text.

    Args:
        text: Full document text
        max_chars: Maximum characters for summary

    Returns:
        Truncated text with ellipsis if needed
    """
    if len(text) <= max_chars:
        return text

    return text[:max_chars] + "..."


def parse_spreadsheet_to_records(file_path: str) -> Dict[str, Any]:
    """
    Parse spreadsheet into structured records for direct database import.

    Args:
        file_path: Path to spreadsheet file

    Returns:
        Dictionary with headers, records, and metadata
    """
    lower_path = file_path.lower()

    if lower_path.endswith(".csv"):
        return _parse_csv_to_records(file_path)
    elif lower_path.endswith(".xlsx"):
        return _parse_xlsx_to_records(file_path)
    elif lower_path.endswith(".xls"):
        return _parse_xls_to_records(file_path)
    else:
        return {"error": "Unsupported spreadsheet format", "records": []}


def _parse_csv_to_records(file_path: str) -> Dict[str, Any]:
    """Parse CSV into structured records."""
    try:
        import csv

        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
        file_content = None

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding, newline="") as f:
                    file_content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if file_content is None:
            return {"error": "Could not decode CSV file", "records": []}

        reader = csv.DictReader(file_content.splitlines())
        records = list(reader)
        headers = reader.fieldnames or []

        return {
            "success": True,
            "format": "csv",
            "headers": headers,
            "records": records,
            "record_count": len(records),
            "column_count": len(headers),
        }

    except Exception as e:
        return {"error": str(e), "records": []}


def _parse_xlsx_to_records(file_path: str) -> Dict[str, Any]:
    """Parse XLSX into structured records."""
    try:
        import openpyxl

        workbook = openpyxl.load_workbook(file_path, data_only=True)
        all_records = {}

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            if sheet.max_row < 1:
                continue

            # Get headers from first row
            headers = []
            for col in range(1, sheet.max_column + 1):
                cell = sheet.cell(row=1, column=col)
                headers.append(str(cell.value) if cell.value else f"column_{col}")

            # Get records
            records = []
            for row in range(2, sheet.max_row + 1):
                record = {}
                for col, header in enumerate(headers, 1):
                    cell_val = sheet.cell(row=row, column=col).value
                    record[header] = cell_val
                records.append(record)

            all_records[sheet_name] = {
                "headers": headers,
                "records": records,
                "record_count": len(records),
            }

        workbook.close()

        return {
            "success": True,
            "format": "xlsx",
            "sheets": all_records,
            "sheet_count": len(all_records),
        }

    except ImportError:
        return {"error": "openpyxl not installed", "records": []}
    except Exception as e:
        logger.warning(f"Primary XLSX parsing failed: {e}. Attempting fallback (read-only mode)...")
        # Try fallback: read-only mode without styles
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(
                file_path, data_only=True, read_only=True, keep_links=False
            )
            all_records = {}

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                if sheet.max_row < 1:
                    continue

                # Get headers from first row
                headers = []
                first_row = next(sheet.iter_rows(min_row=1, max_row=1, values_only=True), None)
                if first_row:
                    headers = [
                        str(cell)[:50] if cell is not None else f"column_{i + 1}"
                        for i, cell in enumerate(first_row)
                    ]

                # Get records
                records = []
                for row in sheet.iter_rows(min_row=2, values_only=True):
                    if any(cell is not None for cell in row):
                        record = {}
                        for col, header in enumerate(headers):
                            if col < len(row):
                                record[header] = row[col]
                        records.append(record)

                all_records[sheet_name] = {
                    "headers": headers,
                    "records": records,
                    "record_count": len(records),
                }

            workbook.close()

            # SUCCESS: Fallback worked, return success without scary error messages
            logger.info(
                f"✅ Fallback XLSX parsing succeeded: {len(all_records)} sheets, {sum(s['record_count'] for s in all_records.values())} total records"
            )
            return {
                "success": True,
                "format": "xlsx",
                "sheets": all_records,
                "sheet_count": len(all_records),
                "fallback_mode": True,
                "fallback_note": "File had formatting issues but data was successfully extracted in read-only mode",
            }
        except Exception as fallback_error:
            logger.error(f"Fallback XLSX parsing also failed: {fallback_error}")
            return {
                "error": f"Excel parsing failed: {str(e)}. Fallback also failed: {str(fallback_error)}",
                "records": [],
            }


def _parse_xls_to_records(file_path: str) -> Dict[str, Any]:
    """Parse XLS into structured records."""
    try:
        import xlrd

        workbook = xlrd.open_workbook(file_path)
        all_records = {}

        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)

            if sheet.nrows < 1:
                continue

            # Get headers
            headers = [
                str(sheet.cell_value(0, col) or f"column_{col}") for col in range(sheet.ncols)
            ]

            # Get records
            records = []
            for row in range(1, sheet.nrows):
                record = {}
                for col, header in enumerate(headers):
                    record[header] = sheet.cell_value(row, col)
                records.append(record)

            all_records[sheet.name] = {
                "headers": headers,
                "records": records,
                "record_count": len(records),
            }

        return {
            "success": True,
            "format": "xls",
            "sheets": all_records,
            "sheet_count": len(all_records),
        }

    except ImportError:
        return {"error": "xlrd not installed", "records": []}
    except Exception as e:
        return {"error": str(e), "records": []}
