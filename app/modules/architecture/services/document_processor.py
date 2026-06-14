"""
Document Processor for ArchiMate Engine

Orchestrates multi-stage document processing:
1. File upload and validation
2. Content extraction (text, images, diagrams)
3. Multi-modal AI analysis
4. ArchiMate element mapping
5. Validation and merging
"""

import asyncio
import inspect
import json
import logging
import os
import uuid
from datetime import datetime
from typing import Dict, List, Optional, Tuple

from app import db
from app.models import (
    APISettings,
    ArchiMateElement,
    ArchiMateRelationship,
    ArchitectureModel,
    LLMInteraction,
    Requirement,
    User,
)
from app.services.llm_service import LLMService
from app.services.progress_service import get_progress_tracker

from app.services.archimate.archimate_validator import ArchiMateValidator
from app.services.archimate.document_upload_service import DocumentUploadService
from app.services.archimate.gemini_file_search_service import GeminiExtractionResult, GeminiFileSearchService
from app.services.archimate.multi_modal_llm_service import MultiModalLLMService
from app.services.archimate.tabular_data_extractor import TabularDataExtractor


class DocumentProcessor:
    """
    Orchestrates end-to-end document processing for ArchiMate generation
    """

    def __init__(self, upload_folder: str):
        """
        Initialize document processor

        Args:
            upload_folder: Base directory for uploads
        """
        self.upload_service = DocumentUploadService(upload_folder)
        self.multi_modal_service = MultiModalLLMService()
        self.validator = ArchiMateValidator()
        self.tabular_extractor = TabularDataExtractor()
        self._gemini_service: Optional[GeminiFileSearchService] = None
        self._logger = logging.getLogger(__name__)

    def _get_gemini_service(self) -> GeminiFileSearchService:
        """Lazily instantiate the Gemini File Search service."""
        if self._gemini_service is None:
            self._gemini_service = GeminiFileSearchService()
        return self._gemini_service

    async def _resolve_mapping(self, candidate, stage_name: str) -> Dict:
        """Ensure extracted data is a dictionary, awaiting awaitables as needed."""
        value = candidate

        if inspect.isawaitable(value):
            value = await value

        if isinstance(value, dict):
            return value

        if hasattr(value, "to_dict"):
            try:
                return value.to_dict()  # type: ignore[call-arg]
            except Exception as exc:  # pragma: no cover - defensive logging
                self._logger.warning(
                    "Failed to convert %s stage data via to_dict: %s", stage_name, exc
                )

        self._logger.warning(
            "Unexpected data type for %s stage: %s. Falling back to empty mapping.",
            stage_name,
            type(value).__name__,
        )
        return {}

    @staticmethod
    def _infer_document_type_hint(file_path: str) -> str:
        """Infer a loose document type string from the file extension."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext in {".pdf"}:
            return "requirements"
        if ext in {".doc", ".docx"}:
            return "business_case"
        if ext in {".ppt", ".pptx"}:
            return "architecture_briefing"
        return "enterprise_document"

    @staticmethod
    def _resolve_gemini_model() -> str:
        """Resolve the configured Gemini model, falling back to the service default."""
        settings = APISettings.query.filter_by(provider="gemini", enabled=True).first()
        if settings and settings.default_model:
            return settings.default_model
        return GeminiFileSearchService.DEFAULT_MODEL

    async def process_uploaded_document(
        self,
        file_path: str,
        file_type: str,
        model_name: str,
        user_id: int,
        additional_context: Optional[str] = None,
        provider: str = "claude",
        task_id: Optional[str] = None,
        enable_progress: bool = False,
    ) -> Tuple[ArchitectureModel, List[LLMInteraction], Dict]:
        """
        Process uploaded document and generate ArchiMate architecture

        Args:
            file_path: Path to uploaded file
            file_type: 'image', 'document', 'text', or 'spreadsheet'
            model_name: Name for the architecture model
            user_id: User ID
            additional_context: Optional context to guide generation
            provider: 'claude', 'openai', or 'gemini'
            task_id: Optional task ID for progress tracking (auto-generated if enable_progress=True)
            enable_progress: Enable real-time progress tracking

        Returns:
            (architecture_model, llm_interactions, processing_report)
        """
        processing_report = {
            "file_path": file_path,
            "file_type": file_type,
            "stages": [],
            "errors": [],
            "warnings": [],
        }

        interactions = []

        # Initialize progress tracking if enabled
        progress_tracker = None
        if enable_progress:
            progress_tracker = get_progress_tracker()
            if not task_id:
                task_id = str(uuid.uuid4())

            # Determine total stages based on file type
            total_stages = 4  # extraction, enrichment (optional), model creation, validation
            if file_type == "spreadsheet":
                total_stages = 2  # extraction, validation
            elif not additional_context:
                total_stages = 3  # extraction, model creation, validation

            progress_tracker.create_task(
                task_id=task_id,
                total_stages=total_stages,
                description=f"Processing {file_type} document: {model_name}",
            )
            processing_report["task_id"] = task_id

        # Stage 1: Extract content based on file type
        if file_type == "image":
            if progress_tracker:
                progress_tracker.update_stage(task_id, 1, "Analyzing image", 0)

            extracted_data, interaction = await self._process_image(file_path, provider)
            extracted_data = await self._resolve_mapping(extracted_data, "image_analysis")
            if interaction:
                interactions.append(interaction)

            stage_result = {
                "stage": "image_analysis",
                "status": "completed",
                "elements_extracted": len(extracted_data.get("elements", [])),
                "relationships_extracted": len(extracted_data.get("relationships", [])),
            }
            processing_report["stages"].append(stage_result)

            if progress_tracker:
                progress_tracker.complete_stage(task_id, "Image Analysis", stage_result)
                progress_tracker.update_stage(task_id, 1, "Image analysis complete", 100)

        elif file_type == "document":
            if progress_tracker:
                progress_tracker.update_stage(task_id, 1, "Analyzing document", 0)

            extracted_data, doc_interactions = await self._process_document(
                file_path, provider, additional_context
            )
            extracted_data = await self._resolve_mapping(extracted_data, "document_analysis")
            interactions.extend(doc_interactions)

            stage_result = {
                "stage": "document_analysis",
                "status": "completed",
                "elements_extracted": len(extracted_data.get("elements", [])),
                "relationships_extracted": len(extracted_data.get("relationships", [])),
            }
            processing_report["stages"].append(stage_result)

            if progress_tracker:
                progress_tracker.complete_stage(task_id, "Document Analysis", stage_result)
                progress_tracker.update_stage(task_id, 1, "Document analysis complete", 100)

        elif file_type == "text":
            if progress_tracker:
                progress_tracker.update_stage(task_id, 1, "Analyzing text", 0)

            extracted_data, interaction = await self._process_text_file(
                file_path, provider, additional_context
            )
            extracted_data = await self._resolve_mapping(extracted_data, "text_analysis")
            if interaction:
                interactions.append(interaction)

            stage_result = {
                "stage": "text_analysis",
                "status": "completed",
                "elements_extracted": len(extracted_data.get("elements", [])),
                "relationships_extracted": len(extracted_data.get("relationships", [])),
            }
            processing_report["stages"].append(stage_result)

            if progress_tracker:
                progress_tracker.complete_stage(task_id, "Text Analysis", stage_result)
                progress_tracker.update_stage(task_id, 1, "Text analysis complete", 100)

        elif file_type == "spreadsheet":
            if progress_tracker:
                progress_tracker.update_stage(task_id, 1, "Extracting data from spreadsheet", 0)

            # Create a temporary architecture to get an ID for the extractor
            temp_arch = ArchitectureModel(
                name=model_name,
                description=f"Auto-generated from spreadsheet upload",
                created_by_id=user_id,
                created_at=datetime.utcnow(),
            )
            db.session.add(temp_arch)
            db.session.flush()

            # Process spreadsheet
            extraction_results = self.tabular_extractor.extract_elements(
                file_path=file_path,
                architecture_id=temp_arch.id,
                portfolio_type=None,  # Auto-detect
                user_id=user_id,
            )

            stage_result = {
                "stage": "spreadsheet_extraction",
                "status": "completed",
                "portfolio_type": extraction_results["portfolio_type"],
                "elements_created": extraction_results["elements_created"],
                "relationships_created": extraction_results["relationships_created"],
                "errors": extraction_results["errors"],
                "warnings": extraction_results["warnings"],
            }
            processing_report["stages"].append(stage_result)

            if progress_tracker:
                progress_tracker.complete_stage(task_id, "Spreadsheet Extraction", stage_result)
                progress_tracker.update_stage(task_id, 1, "Spreadsheet extraction complete", 100)

            # Return early for spreadsheets since extraction already created elements
            db.session.commit()

            # Fetch the created architecture
            architecture_model = db.session.get(ArchitectureModel, temp_arch.id)

            # Validate the model
            if progress_tracker:
                progress_tracker.update_stage(task_id, 2, "Validating architecture", 0)

            validation_results = self.validator.validate_model(architecture_model)
            validation_stage = {
                "stage": "validation",
                "status": "completed",
                "is_valid": validation_results["is_valid"],
                "errors": len(validation_results["element_errors"])
                + len(validation_results["relationship_errors"]),
                "warnings": len(validation_results["warnings"]),
            }
            processing_report["stages"].append(validation_stage)

            processing_report.update(
                {
                    "validation_errors": validation_results["element_errors"]
                    + validation_results["relationship_errors"],
                    "validation_warnings": validation_results["warnings"],
                }
            )

            if progress_tracker:
                progress_tracker.complete_stage(task_id, "Validation", validation_stage)
                progress_tracker.complete_task(task_id, "completed", validation_stage)

            return architecture_model, interactions, processing_report

        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        # Stage 2: Enrich with additional context if provided
        current_stage = 2
        if additional_context:
            if progress_tracker:
                progress_tracker.update_stage(task_id, current_stage, "Enriching with context", 0)

            enriched_data, interaction = await self._enrich_with_context(
                extracted_data, additional_context, provider
            )
            enriched_data = await self._resolve_mapping(enriched_data, "context_enrichment")
            if interaction:
                interactions.append(interaction)

            enrichment_result = {
                "stage": "context_enrichment",
                "status": "completed",
                "elements_added": len(enriched_data.get("elements", []))
                - len(extracted_data.get("elements", [])),
                "relationships_added": len(enriched_data.get("relationships", []))
                - len(extracted_data.get("relationships", [])),
            }
            processing_report["stages"].append(enrichment_result)

            if progress_tracker:
                progress_tracker.complete_stage(task_id, "Context Enrichment", enrichment_result)
                progress_tracker.update_stage(
                    task_id, current_stage, "Context enrichment complete", 100
                )

            current_stage += 1
        else:
            enriched_data = extracted_data

        # Stage 3: Create ArchiMate model in database
        if progress_tracker:
            progress_tracker.update_stage(task_id, current_stage, "Creating architecture model", 0)

        enriched_data = await self._resolve_mapping(enriched_data, "model_creation")

        architecture_model = await self._create_architecture_model(
            enriched_data, model_name, user_id
        )

        model_result = {
            "stage": "model_creation",
            "status": "completed",
            "model_id": architecture_model.id,
            "elements_created": len(architecture_model.elements),
            "relationships_created": len(architecture_model.relationships),
        }
        processing_report["stages"].append(model_result)

        if progress_tracker:
            progress_tracker.complete_stage(task_id, "Model Creation", model_result)
            progress_tracker.update_stage(task_id, current_stage, "Model creation complete", 100)

        current_stage += 1

        # Stage 4: Validate generated architecture
        if progress_tracker:
            progress_tracker.update_stage(task_id, current_stage, "Validating architecture", 0)

        validation_results = self.validator.validate_model(architecture_model)

        validation_result = {
            "stage": "validation",
            "status": "completed",
            "is_valid": validation_results["is_valid"],
            "errors": len(validation_results["element_errors"])
            + len(validation_results["relationship_errors"]),
            "warnings": len(validation_results["warnings"]),
        }
        processing_report["stages"].append(validation_result)

        if not validation_results["is_valid"]:
            processing_report["warnings"].append(
                f"Architecture has {len(validation_results['element_errors']) + len(validation_results['relationship_errors'])} validation errors"
            )

        if progress_tracker:
            progress_tracker.complete_stage(task_id, "Validation", validation_result)
            progress_tracker.complete_task(task_id, "completed", validation_result)

        return architecture_model, interactions, processing_report

    async def _process_image(
        self, image_path: str, provider: str
    ) -> Tuple[Dict, Optional[LLMInteraction]]:
        """
        Process image file with vision AI

        Returns:
            (extracted_data, llm_interaction)
        """
        extracted_data, interaction = await self.multi_modal_service.extract_archimate_from_diagram(
            image_path, provider
        )

        return extracted_data, interaction

    async def _process_text_file(
        self, file_path: str, provider: str, additional_context: Optional[str] = None
    ) -> Tuple[Dict, Optional[LLMInteraction]]:
        """
        Process text file (TXT, MD, HTML)

        Returns:
            (extracted_data, llm_interaction)
        """
        # Read text content
        with open(file_path, "r", encoding="utf-8") as f:
            text_content = f.read()

        # Determine document type from content
        document_type = "requirements"  # Default
        if "requirement" in text_content.lower()[:500]:
            document_type = "requirements"
        elif "technical" in text_content.lower()[:500]:
            document_type = "technical_spec"
        elif "business case" in text_content.lower()[:500]:
            document_type = "business_case"

        if provider == "gemini":
            return await self._process_with_gemini_file_search(
                file_path=file_path,
                document_type=document_type,
                additional_context=additional_context,
            )

        extracted_data, interaction = await self.multi_modal_service.analyze_document_text(
            text_content, document_type, provider
        )

        return extracted_data, interaction

    async def _process_document(
        self, file_path: str, provider: str, additional_context: Optional[str] = None
    ) -> Tuple[Dict, List[LLMInteraction]]:
        """
        Process document file (PDF, DOC, PPT)

        For now, this is a simplified implementation.
        Future: Integrate proper PDF parsing, OCR, diagram extraction

        Returns:
            (extracted_data, list_of_interactions)
        """
        interactions = []

        if provider == "gemini":
            document_type = self._infer_document_type_hint(file_path)
            return await self._process_with_gemini_file_search(
                file_path=file_path,
                document_type=document_type,
                additional_context=additional_context,
            )

        # Simplified: Treat as text extraction for now
        # In production, use libraries like PyPDF2, python-pptx, python-docx
        # and extract both text and embedded images

        try:
            # Placeholder: Extract text (you would use proper libraries here)
            if file_path.endswith(".pdf"):
                text_content = self._extract_text_from_pdf(file_path)
            elif file_path.endswith((".doc", ".docx")):
                text_content = self._extract_text_from_word(file_path)
            elif file_path.endswith((".ppt", ".pptx")):
                text_content = self._extract_text_from_powerpoint(file_path)
            else:
                text_content = ""

            if text_content:
                extracted_data, interaction = await self.multi_modal_service.analyze_document_text(
                    text_content, "requirements", provider
                )
                if interaction:
                    interactions.append(interaction)
            else:
                extracted_data = {"elements": [], "relationships": [], "metadata": {}}

            return extracted_data, interactions

        except Exception as e:
            # Return empty results on error
            return {
                "elements": [],
                "relationships": [],
                "metadata": {"error": str(e)},
            }, interactions

    async def _process_with_gemini_file_search(
        self, file_path: str, document_type: str, additional_context: Optional[str]
    ) -> Tuple[Dict, List[LLMInteraction]]:
        """Process any file via Gemini File Search."""

        service = self._get_gemini_service()
        model_name = self._resolve_gemini_model()

        result: GeminiExtractionResult = await asyncio.to_thread(
            service.extract_archimate_elements,
            file_path,
            document_type,
            additional_context,
            model_name,
        )

        parsed = self._parse_archimate_json(result.response_text)

        metadata = parsed.setdefault("metadata", {})
        if result.grounding_metadata:
            metadata["gemini_grounding"] = result.grounding_metadata
        if result.citation_metadata:
            metadata["gemini_citations"] = result.citation_metadata

        usage = result.usage or {}
        prompt_tokens = usage.get("promptTokenCount") or 0
        response_tokens = usage.get("candidatesTokenCount") or 0

        if prompt_tokens == 0 and response_tokens == 0:
            prompt_tokens = len(result.prompt) // 4
            response_tokens = len(result.response_text) // 4

        interaction = LLMInteraction(
            model_name=model_name,
            provider="gemini",
            prompt=result.prompt[:1000],
            response=result.response_text[:2000],
            token_count_input=prompt_tokens,
            token_count_output=response_tokens,
            cost=LLMService.estimate_cost(prompt_tokens * 4, response_tokens * 4, model=model_name),
        )

        return parsed, [interaction]

    @staticmethod
    def _parse_archimate_json(response_text: str) -> Dict:
        """Parse JSON payload from a model response, capturing errors."""
        try:
            if "```json" in response_text:
                json_start = response_text.find("```json") + 7
                json_end = response_text.find("```", json_start)
                json_text = response_text[json_start:json_end].strip()
            elif "```" in response_text:
                json_start = response_text.find("```") + 3
                json_end = response_text.find("```", json_start)
                json_text = response_text[json_start:json_end].strip()
            else:
                json_text = response_text

            parsed = json.loads(json_text)
            if isinstance(parsed, dict):
                return parsed
        except json.JSONDecodeError as exc:
            return {
                "elements": [],
                "relationships": [],
                "metadata": {
                    "error": f"JSON parsing failed: {exc}",
                    "raw_response": response_text[:2000],
                },
            }

        return {
            "elements": [],
            "relationships": [],
            "metadata": {
                "error": "Unexpected response format from Gemini",
                "raw_response": response_text[:2000],
            },
        }

    def _extract_text_from_pdf(self, file_path: str) -> str:
        """
        Extract text from PDF file using pdfplumber (primary) with PyPDF2 fallback

        Args:
            file_path: Path to PDF file

        Returns:
            Extracted text content
        """
        import logging

        logger = logging.getLogger(__name__)

        text_parts = []

        try:
            # Try pdfplumber first (better quality)
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(f"\n--- Page {page_num} ---\n")
                        text_parts.append(page_text)

            if text_parts:
                logger.info(f"Extracted {len(pdf.pages)} pages from PDF using pdfplumber")
                return "\n\n".join(text_parts)

        except ImportError:
            logger.warning("pdfplumber not installed, falling back to PyPDF2")
        except Exception as e:
            logger.warning(f"pdfplumber extraction failed: {e}, falling back to PyPDF2")

        # Fallback to PyPDF2
        try:
            import PyPDF2

            text_parts = []
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(f"\n--- Page {page_num} ---\n")
                        text_parts.append(page_text)

            logger.info(f"Extracted {len(reader.pages)} pages from PDF using PyPDF2")
            return "\n\n".join(text_parts)

        except ImportError:
            logger.error(
                "Neither pdfplumber nor PyPDF2 is installed. Install with: pip install pdfplumber PyPDF2"
            )
            return ""
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return ""

    def _extract_text_from_word(self, file_path: str) -> str:
        """
        Extract text from Word document (.docx)

        Args:
            file_path: Path to DOCX file

        Returns:
            Extracted text content including tables
        """
        import logging

        logger = logging.getLogger(__name__)

        try:
            from docx import Document

            doc = Document(file_path)
            text_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            for table_num, table in enumerate(doc.tables, 1):
                text_parts.append(f"\n--- Table {table_num} ---")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)

            logger.info(
                f"Extracted {len(doc.paragraphs)} paragraphs and {len(doc.tables)} tables from DOCX"
            )
            return "\n\n".join(text_parts)

        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return ""
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return ""

    def _extract_text_from_powerpoint(self, file_path: str) -> str:
        """
        Extract text from PowerPoint (.pptx)

        Args:
            file_path: Path to PPTX file

        Returns:
            Extracted text content from all slides
        """
        import logging

        logger = logging.getLogger(__name__)

        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n--- Slide {slide_num} ---")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                    # Extract from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                text_parts.append(row_text)

            logger.info(f"Extracted {len(prs.slides)} slides from PPTX")
            return "\n\n".join(text_parts)

        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return ""

    async def _enrich_with_context(
        self, extracted_data: Dict, context: str, provider: str
    ) -> Tuple[Dict, Optional[LLMInteraction]]:
        """
        Enrich extracted architecture with additional context

        Args:
            extracted_data: Initial extracted elements/relationships
            context: Additional context (existing systems, constraints, etc.)
            provider: 'claude', 'openai', or 'gemini'

        Returns:
            (enriched_data, llm_interaction)
        """
        from .archimate_prompts import ARCHIMATE_SYSTEM_PROMPT

        prompt = f"""
You have extracted the following ArchiMate architecture from a document:

EXTRACTED ELEMENTS:
{json.dumps(extracted_data.get('elements', []), indent=2)}

EXTRACTED RELATIONSHIPS:
{json.dumps(extracted_data.get('relationships', []), indent=2)}

ADDITIONAL CONTEXT:
{context}

Based on this additional context, enhance the architecture by:
1. Adding missing elements that are implied by the context
2. Adding missing relationships
3. Enriching element descriptions with context information
4. Ensuring all layers are properly represented

Return the COMPLETE enhanced architecture in the same JSON format:
{{
  "elements": [...all elements including new ones...],
  "relationships": [...all relationships including new ones...],
  "metadata": {{
    "enrichments_made": "description of what was added"
  }}
}}
"""

        # Get model from database APISettings (NO HARDCODED MODELS per policy)
        from app.models import APISettings

        settings = APISettings.query.filter_by(provider=provider, enabled=True).first()
        if not settings or not settings.default_model:
            raise ValueError(
                f"No configured model found for provider '{provider}'. "
                f"Please configure via /admin/api-settings"
            )

        model = settings.default_model

        from .multi_modal_llm_service import MultiModalLLMService

        service = MultiModalLLMService()

        response_text, interaction = await service.llm_service.generate_with_prompt(
            prompt=prompt,
            system_prompt=ARCHIMATE_SYSTEM_PROMPT,
            model=model,
            max_tokens=8000,
            temperature=0.2,
        )

        # Parse JSON
        try:
            if "```json" in response_text:
                json_start = response_text.find("```json") + 7
                json_end = response_text.find("```", json_start)
                json_text = response_text[json_start:json_end].strip()
            elif "```" in response_text:
                json_start = response_text.find("```") + 3
                json_end = response_text.find("```", json_start)
                json_text = response_text[json_start:json_end].strip()
            else:
                json_text = response_text

            enriched_data = json.loads(json_text)
            return enriched_data, interaction

        except json.JSONDecodeError:
            # Return original if enrichment fails
            return extracted_data, interaction

    async def _create_architecture_model(
        self, data: Dict, model_name: str, user_id: int
    ) -> ArchitectureModel:
        """
        Create ArchiMate model in database from extracted data

        Args:
            data: Dictionary with 'elements' and 'relationships'
            model_name: Name for the model
            user_id: User ID

        Returns:
            Created ArchitectureModel instance
        """
        # Create architecture model
        architecture = ArchitectureModel(
            name=model_name,
            description=data.get("metadata", {}).get("notes", ""),
            created_by_id=user_id,
            created_at=datetime.utcnow(),
        )
        db.session.add(architecture)
        db.session.flush()

        # Create elements
        element_map = {}  # Maps element name to database ID
        requirement_elements = []  # Track requirement elements for later processing

        for elem_data in data.get("elements", []):
            element = ArchiMateElement(
                architecture_id=architecture.id,
                name=elem_data["name"],
                type=elem_data["type"],
                layer=elem_data.get("layer", "business"),
                description=elem_data.get("description", ""),
                properties=elem_data.get("properties", {}),
            )
            db.session.add(element)
            db.session.flush()

            element_map[elem_data["name"]] = element.id

            # Track requirement elements to create Requirement records
            if elem_data["type"].lower() == "requirement":
                requirement_elements.append((element, elem_data))

        # Create relationships
        for rel_data in data.get("relationships", []):
            source_name = rel_data["source"]
            target_name = rel_data["target"]

            if source_name in element_map and target_name in element_map:
                relationship = ArchiMateRelationship(
                    architecture_id=architecture.id,
                    source_id=element_map[source_name],
                    target_id=element_map[target_name],
                    type=rel_data["type"],
                    description=rel_data.get("description", ""),
                )
                db.session.add(relationship)

        # Create Requirement records for all Requirement-type ArchiMate elements
        # This ensures they appear on the Requirements dashboard with cards and tables
        # Use AI-powered MotivationLayerService for intelligent requirement enrichment
        from app.services.archimate import MotivationLayerService

        motivation_service = MotivationLayerService()

        for element, elem_data in requirement_elements:
            # Extract properties from the element data
            properties = elem_data.get("properties", {})

            # Create basic requirement record linked to ArchiMate element
            requirement = Requirement(
                title=element.name,
                description=element.description,
                category=properties.get("category", "Functional"),
                priority=properties.get("priority", "medium"),
                architecture_id=architecture.id,
                archimate_element_id=element.id,  # Link to ArchiMate element (Basecoat pattern)
                compliance_status="draft",
            )
            db.session.add(requirement)
            db.session.flush()  # Get requirement ID

            # Use AI to enrich the requirement with intelligent analysis
            enriched_req = motivation_service.enrich_requirement(
                requirement_id=requirement.id,
                context={
                    "source": "document_upload",
                    "properties": properties,
                    "element_data": elem_data,
                },
            )

            # AI determines type, priority, and relationships from description
            if enriched_req:
                requirement.type = enriched_req.get("type", requirement.type)
                requirement.priority = enriched_req.get("priority", requirement.priority)
                requirement.category = enriched_req.get("category", requirement.category)
                requirement.verification_method = enriched_req.get("verification_method")
                requirement.rationale = enriched_req.get("rationale")

        db.session.commit()

        # CRITICAL: Auto-link requirements to Motivation Layer elements
        # This establishes proper ArchiMate 3.2 relationships between:
        # - Requirements -> Stakeholders (influence)
        # - Requirements -> Drivers (association)
        # - Requirements -> Goals (realization)
        # - Requirements -> Parent Requirements (hierarchical decomposition)
        logger.info(f"[INFO] Auto-linking requirements to Motivation Layer elements...")
        linking_results = motivation_service.auto_link_requirements_to_motivation_elements(
            architecture_id=architecture.id, use_ai=True  # Use AI for intelligent semantic matching
        )

        logger.info(f"[SUCCESS] ArchiMate 3.2 linking complete:")
        logger.info(f"  - Requirements processed: {linking_results['requirements_processed']}")
        logger.info(f"  - Stakeholder links: {linking_results['stakeholder_links_created']}")
        logger.info(f"  - Driver links: {linking_results['driver_links_created']}")
        logger.info(f"  - Goal links: {linking_results['goal_links_created']}")
        logger.info(f"  - Parent requirement links: {linking_results['parent_links_created']}")
        logger.info(
            f"  - ArchiMate relationships: {linking_results['archimate_relationships_created']}"
        )

        # ========================================================================
        # LAYER ENRICHMENT PIPELINE - Enhance extracted elements with AI-powered layer services
        # ========================================================================
        logger.info(f"[INFO] Starting layer enrichment pipeline...")

        from app.services.archimate import (
            ApplicationLayerService,
            BusinessLayerService,
            RelationshipService,
            TechnologyLayerService,
        )

        # Initialize services
        business_service = BusinessLayerService()
        application_service = ApplicationLayerService()
        technology_service = TechnologyLayerService()
        relationship_service = RelationshipService()

        # ---- Business Layer Enrichment ----
        logger.info(f"[INFO] Enriching Business Layer...")
        business_elements = ArchiMateElement.query.filter_by(
            architecture_id=architecture.id, layer="business"
        ).all()

        if business_elements:
            # Enrich business processes with flow modeling
            business_processes = [e for e in business_elements if e.type == "BusinessProcess"]
            for process in business_processes:
                try:
                    # Model process flow (steps, decisions, exceptions)
                    flow_data = business_service.model_process_flow(
                        process.id, process.description or ""
                    )
                    # Store flow data in element properties
                    props = process.properties or {}
                    props["process_flow"] = flow_data
                    process.properties = props
                    logger.info(f"  ✓ Modeled flow for: {process.name}")
                except Exception as e:
                    logger.info(f"  ✗ Flow modeling failed for {process.name}: {str(e)}")

        # ---- Application Layer Enrichment ----
        logger.info(f"[INFO] Enriching Application Layer...")
        application_elements = ArchiMateElement.query.filter_by(
            architecture_id=architecture.id, layer="application"
        ).all()

        if application_elements:
            # Analyze application dependencies
            app_components = [e for e in application_elements if e.type == "ApplicationComponent"]
            for component in app_components:
                try:
                    # Analyze dependencies (upstream/downstream)
                    deps = application_service.analyze_application_dependencies(
                        component.id, technical_context=component.description
                    )
                    logger.info(f"  ✓ Analyzed dependencies for: {component.name}")
                except Exception as e:
                    logger.info(f"  ✗ Dependency analysis failed for {component.name}: {str(e)}")

            # Map application to business process relationships
            business_processes = ArchiMateElement.query.filter_by(
                architecture_id=architecture.id, type="BusinessProcess"
            ).all()

            for app_component in app_components:
                for process in business_processes:
                    # Check if app serves process based on naming/description similarity
                    if (
                        process.name.lower() in app_component.name.lower()
                        or process.name.lower() in (app_component.description or "").lower()
                    ):
                        try:
                            application_service.map_application_to_business_process(
                                app_component.id,
                                process.id,
                                usage_description=f"{app_component.name} supports {process.name}",
                            )
                            logger.info(f"  ✓ Linked {app_component.name} → {process.name}")
                        except Exception as e:
                            logger.info(f"  ✗ Linking failed: {str(e)}")

        # ---- Technology Layer Enrichment ----
        logger.info(f"[INFO] Enriching Technology Layer...")
        technology_elements = ArchiMateElement.query.filter_by(
            architecture_id=architecture.id, layer="technology"
        ).all()

        if technology_elements:
            # Analyze technology dependencies
            nodes = [e for e in technology_elements if e.type == "Node"]
            for node in nodes:
                try:
                    # Analyze infrastructure dependencies
                    deps = technology_service.analyze_technology_dependencies(
                        node.id, infrastructure_context=node.description
                    )
                    logger.info(f"  ✓ Analyzed tech dependencies for: {node.name}")
                except Exception as e:
                    logger.info(f"  ✗ Tech dependency analysis failed for {node.name}: {str(e)}")

            # Map deployment architecture (app → nodes)
            app_components = ArchiMateElement.query.filter_by(
                architecture_id=architecture.id, type="ApplicationComponent"
            ).all()

            for app in app_components:
                for node in nodes:
                    # Check if app deploys to node based on naming/description
                    if (
                        node.name.lower() in app.name.lower()
                        or node.name.lower() in (app.description or "").lower()
                        or "server" in node.name.lower()
                        or "cluster" in node.name.lower()
                    ):
                        try:
                            technology_service.model_deployment_architecture(
                                app.id,
                                [node.id],
                                deployment_description=f"{app.name} deployed on {node.name}",
                            )
                            logger.info(f"  ✓ Deployed {app.name} → {node.name}")
                        except Exception as e:
                            logger.info(f"  ✗ Deployment mapping failed: {str(e)}")

        # ---- Cross-Layer Relationship Discovery ----
        logger.info(f"[INFO] Discovering cross-layer relationships...")
        try:
            # Discover realization relationships across layers
            context = data.get("model_description", "") + "\n" + data.get("rationale", "")
            discovered_rels = relationship_service.discover_realization_relationships(
                architecture.id, context
            )
            logger.info(f"  ✓ Discovered {len(discovered_rels)} realization relationships")
        except Exception as e:
            logger.info(f"  ✗ Relationship discovery failed: {str(e)}")

        try:
            # Discover flow relationships (data/control flows)
            flow_rels = relationship_service.discover_flow_relationships(architecture.id, context)
            logger.info(f"  ✓ Discovered {len(flow_rels)} flow relationships")
        except Exception as e:
            logger.info(f"  ✗ Flow discovery failed: {str(e)}")

        # ---- Validate Relationship Consistency ----
        logger.info(f"[INFO] Validating relationship consistency...")
        try:
            validation_results = relationship_service.validate_relationship_consistency(
                architecture.id
            )
            if validation_results["valid"]:
                logger.info(f"  ✓ All relationships valid")
            else:
                logger.info(f"  ⚠ Validation warnings:")
                for error in validation_results["errors"][:5]:  # Show first 5
                    logger.info(f"    - {error}")
                for warning in validation_results["warnings"][:5]:
                    logger.info(f"    - {warning}")
        except Exception as e:
            logger.info(f"  ✗ Validation failed: {str(e)}")

        db.session.commit()

        logger.info(f"[SUCCESS] Layer enrichment pipeline complete")

        return architecture
