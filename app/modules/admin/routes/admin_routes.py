"""
Admin Routes (migrated).

Migrated from: app/admin/views.py
Blueprint name preserved: "admin"
URL prefix preserved: /admin (applied via register() in __init__.py)

User management routes delegate to AdminUserService.
Complex routes (API settings, feature flags, abacus, seed) retain inline logic
for Phase 0.5 parity; service extraction planned for Phase 2.
"""

import hashlib
import hmac
import json
import logging
import os
from datetime import datetime

from flask import (
    Blueprint,
    flash,
    jsonify,
    redirect,
    render_template,
    request,
    url_for,
)
from flask_login import current_user, login_required
from sqlalchemy.orm import aliased, joinedload

try:
    from flask_rq import get_queue

    HAS_RQ = True
except ImportError:
    HAS_RQ = False
    get_queue = None

from app.extensions import csrf, db
from ..forms.admin_forms import (
    APISettingsForm,
    ChangeAccountTypeForm,
    ChangeUserEmailForm,
    FeatureFlagForm,
    InviteUserForm,
    NewUserForm,
)
from app.modules.account.forms.account_forms import CreatePasswordForm
from app.decorators import admin_required, audit_log
from app.models import APISettings, EditableHTML, Role, User
from app.models.feature_flags import FeatureFlag, FeatureState, FeatureType
from app.services.llm_service import test_api_key
from app.services.rbac_service import rbac_service
from app.utils.sidebar_parser import SidebarSubmenu, parse_sidebar_template
from ..services.admin_user_service import AdminUserService

admin_bp = Blueprint("admin", __name__)
logger = logging.getLogger(__name__)

_svc = AdminUserService


# ============================================================================
# Dashboard & Index
# ============================================================================


@admin_bp.route("/")
@login_required
@admin_required
def index():
    """Admin dashboard page."""
    return render_template("admin/index.html")


@admin_bp.route("/dashboard-test")
@login_required
@admin_required
def dashboard_test():
    """Admin dashboard test page for dropdown testing."""
    page = request.args.get("page", 1, type=int)
    per_page = request.args.get("per_page", 10, type=int)
    search_query = request.args.get("search", "")

    pagination = _svc.get_paginated_users(page, per_page, search_query)
    users = pagination.items

    return render_template(
        "admin/index.html",
        users=users,
        pagination=pagination,
        search_query=search_query,
        per_page=per_page,
    )


@admin_bp.route("/dashboard")
@admin_bp.route("/dashboard/overview")  # Issue 1 fix: Support both routes
@login_required
@admin_required
def dashboard():
    """Admin dashboard with stats and overview."""
    page = request.args.get("page", 1, type=int)
    per_page = request.args.get("per_page", 10, type=int)
    search_query = request.args.get("search", "")

    pagination = _svc.get_paginated_users(page, per_page, search_query)
    users = pagination.items
    roles = _svc.get_all_roles()

    return render_template(
        "admin/index.html",
        users=users,
        roles=roles,
        pagination=pagination,
        search_query=search_query,
        per_page=per_page,
    )


# ============================================================================
# User Management (thin routes -> AdminUserService)
# ============================================================================


@admin_bp.route("/new-user", methods=["GET", "POST"])
@login_required
@rbac_service.require_role("org_admin")
@admin_required
@audit_log("admin_user_create")
def new_user():
    """Create a new user."""
    form = NewUserForm()
    if form.validate_on_submit():
        # Enforce seat limit before creating the user
        org_id = getattr(current_user, "organization_id", None)
        if org_id is not None:
            from app.services.usage_metering_service import UsageMeteringService
            if not UsageMeteringService.check_seat_limit(org_id):
                return jsonify({
                    "error": "seat_limit_exceeded",
                    "message": "Upgrade your plan to add more users.",
                }), 402

        user = _svc.create_user(
            first_name=form.first_name.data,
            last_name=form.last_name.data,
            email=form.email.data,
            password=form.password.data,
            role=form.role.data,
        )
        flash("User {} successfully created".format(user.full_name()), "form-success")
    return render_template("admin/new_user.html", form=form)


@admin_bp.route("/invite-user", methods=["GET", "POST"])
@login_required
@rbac_service.require_role("org_admin")
@admin_required
@audit_log("admin_user_invite")
def invite_user():
    """Invites a new user to create an account and set their own password."""
    form = InviteUserForm()
    if form.validate_on_submit():
        # Enforce seat limit before inviting the user
        org_id = getattr(current_user, "organization_id", None)
        if org_id is not None:
            from app.services.usage_metering_service import UsageMeteringService
            if not UsageMeteringService.check_seat_limit(org_id):
                return jsonify({
                    "error": "seat_limit_exceeded",
                    "message": "Upgrade your plan to add more users.",
                }), 402

        user = _svc.invite_user(
            first_name=form.first_name.data,
            last_name=form.last_name.data,
            email=form.email.data,
            role=form.role.data,
        )
        flash("User {} successfully invited".format(user.full_name()), "form-success")
    return render_template("admin/new_user.html", form=form)


@admin_bp.route("/manage-users")
def manage_users_redirect():
    """Redirect legacy /admin/manage-users to canonical /admin/users."""
    return redirect(url_for("admin.registered_users"), code=301)


@admin_bp.route("/users")
@login_required
@admin_required
def registered_users():
    """View all registered users."""
    users = User.query.options(joinedload(User.role)).order_by(User.id.desc()).all()
    roles = Role.query.order_by(Role.name).all()
    return render_template("admin/registered_users.html", users=users, roles=roles)


@admin_bp.route("/user/<int:user_id>")
@admin_bp.route("/user/<int:user_id>/info")
@login_required
@admin_required
def user_info(user_id):
    """View a user's profile."""
    user = _svc.get_user_or_404(user_id)
    return render_template("admin/manage_user.html", user=user)


@admin_bp.route("/user/<int:user_id>/change-email", methods=["GET", "POST"])
@login_required
@admin_required
@audit_log("admin_user_email_change")
def change_user_email(user_id):
    """Change a user's email."""
    user = _svc.get_user_or_404(user_id)
    form = ChangeUserEmailForm()
    if form.validate_on_submit():
        _svc.change_user_email(user, form.email.data)
        flash(
            "Email for user {} successfully changed to {}.".format(
                user.full_name(), user.email
            ),
            "form-success",
        )
    return render_template("admin/manage_user.html", user=user, form=form)


@admin_bp.route("/user/<int:user_id>/change-account-type", methods=["GET", "POST"])
@login_required
@admin_required
@audit_log("admin_user_role_change")
def change_account_type(user_id):
    """Change a user's account type."""
    if current_user.id == user_id:
        flash(
            "You cannot change the type of your own account. Please ask "
            "another administrator to do this.",
            "error",
        )
        return redirect(url_for("admin.user_info", user_id=user_id))

    user = _svc.get_user_or_404(user_id)
    form = ChangeAccountTypeForm()
    if form.validate_on_submit():
        _svc.change_user_role(user, form.role.data)
        role_name = user.role.name if user.role else "No Role"
        flash(
            "Role for user {} successfully changed to {}.".format(
                user.full_name(), role_name
            ),
            "form-success",
        )
    return render_template("admin/manage_user.html", user=user, form=form)


@admin_bp.route("/user/<int:user_id>/set-password", methods=["GET", "POST"])
@login_required
@admin_required
@audit_log("admin_user_password_set")
def set_user_password(user_id):
    """Set or reset a user's password."""
    user = _svc.get_user_or_404(user_id)
    form = CreatePasswordForm()
    if form.validate_on_submit():
        _svc.set_user_password(user, form.password.data)
        flash(
            "Password for user {} successfully updated. The account is now login-ready.".format(
                user.full_name()
            ),
            "form-success",
        )
        return redirect(url_for("admin.user_info", user_id=user_id))
    return render_template("admin/manage_user.html", user=user, form=form)


@admin_bp.route("/user/<int:user_id>/delete")
@login_required
@admin_required
def delete_user_request(user_id):
    """Request deletion of a user's account."""
    user = _svc.get_user_or_404(user_id)
    return render_template("admin/manage_user.html", user=user)


@admin_bp.route("/user/<int:user_id>/_delete")
@login_required
@admin_required
@audit_log("admin_user_delete")
def delete_user(user_id):
    """Delete a user's account."""
    if current_user.id == user_id:
        flash(
            "You cannot delete your own account. Please ask another "
            "administrator to do this.",
            "error",
        )
    else:
        user = _svc.get_user_or_404(user_id)
        success, message = _svc.delete_user(user)
        flash(message, "success")
    return redirect(url_for("admin.registered_users"))


# ============================================================================
# Editor Contents
# ============================================================================


@admin_bp.route("/_update_editor_contents", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_editor_update")
def update_editor_contents():
    """Update the contents of an editor."""
    edit_data = request.form.get("edit_data")
    editor_name = request.form.get("editor_name")

    editor_contents = EditableHTML.query.filter_by(editor_name=editor_name).first()
    if editor_contents is None:
        editor_contents = EditableHTML(editor_name=editor_name)
    editor_contents.value = edit_data

    db.session.add(editor_contents)
    db.session.commit()
    return "OK", 200


# ============================================================================
# API Settings
# ============================================================================


@admin_bp.route("/api-settings", methods=["GET", "POST"])
@login_required
@rbac_service.require_role("org_admin")
@admin_required
@audit_log("admin_api_settings_save")
def api_settings():
    """Manage API settings for LLM providers."""
    all_settings = APISettings.query.order_by(APISettings.provider).all()
    settings_dict = {s.provider: s for s in all_settings}
    edit_provider = request.args.get("edit")
    form = APISettingsForm()

    # Known key-format prefixes (provider → required prefix(es))
    _KEY_PREFIXES = {
        "openai":      ("sk-",),
        "anthropic":   ("sk-ant-",),
        "huggingface": ("hf_",),
        "gemini":      ("AIza",),
        "deepseek":    ("sk-",),
        "openrouter":  ("sk-or-",),
        "azure":       (),          # Azure uses a 32-char GUID-style key — no reliable prefix
        "jira":        (),          # JIRA tokens vary; skip prefix check
    }

    if form.validate_on_submit():
        provider = form.provider.data
        raw_key = (form.api_key.data or "").strip()

        # If the field contains the masked placeholder (e.g. "sk-a***…***4321"), the
        # admin didn't change the key — treat it as empty so we don't overwrite.
        if raw_key and "***" in raw_key:
            raw_key = ""

        # Validate key format before touching the DB
        if raw_key:
            allowed_prefixes = _KEY_PREFIXES.get(provider, ())
            if allowed_prefixes and not any(raw_key.startswith(p) for p in allowed_prefixes):
                expected = " or ".join(f'"{p}…"' for p in allowed_prefixes)
                flash(
                    f"Invalid API key format for {provider}. "
                    f"Keys for this provider must start with {expected}. "
                    f"Check for extra spaces, incorrect provider selection, or a copy-paste error.",
                    "error",
                )
                return render_template(
                    "admin/api_settings.html",
                    form=form,
                    all_settings=all_settings,
                    edit_provider=edit_provider,
                )

        existing = settings_dict.get(provider)

        if existing:
            if raw_key:
                existing.api_key = raw_key
            existing.enabled = form.enabled.data
            existing.default_model = form.default_model.data or None
            existing.max_tokens = form.max_tokens.data or 4000
            existing.temperature = form.temperature.data or 0.7
            existing.jira_url = form.jira_url.data or None
            existing.jira_email = form.jira_email.data or None
            existing.hf_model_id = form.hf_model_id.data or None
            existing.hf_endpoint_url = form.hf_endpoint_url.data or None
            existing.custom_endpoint_url = form.custom_endpoint_url.data or None
            existing.custom_auth_method = form.custom_auth_method.data or "bearer"
            existing.custom_headers = form.custom_headers.data or None
            existing.updated_at = datetime.utcnow()
            existing.updated_by_id = current_user.id
            settings = existing
        else:
            if not raw_key:
                flash("API key is required when creating new settings.", "error")
                return render_template(
                    "admin/api_settings.html",
                    form=form,
                    all_settings=all_settings,
                    edit_provider=edit_provider,
                )

            settings = APISettings(
                provider=provider,
                api_key=raw_key,
                enabled=form.enabled.data,
                default_model=form.default_model.data or None,
                max_tokens=form.max_tokens.data or 4000,
                temperature=form.temperature.data or 0.7,
                jira_url=form.jira_url.data or None,
                jira_email=form.jira_email.data or None,
                hf_model_id=form.hf_model_id.data or None,
                hf_endpoint_url=form.hf_endpoint_url.data or None,
                custom_endpoint_url=form.custom_endpoint_url.data or None,
                custom_auth_method=form.custom_auth_method.data or "bearer",
                custom_headers=form.custom_headers.data or None,
                updated_by_id=current_user.id,
            )
            db.session.add(settings)

        # Use the just-validated raw_key for the connection test so we're testing
        # what was submitted, not whatever was already in the DB.
        api_key_to_test = raw_key if raw_key else settings.api_key
        if form.test.data:
            if not api_key_to_test:
                flash("No API key available to test.", "error")
                return redirect(url_for("admin.api_settings"))
            test_result = test_api_key(provider, api_key_to_test, model=form.default_model.data or None)
            settings.test_status = "success" if test_result["success"] else "failed"
            settings.test_message = test_result.get("message", "")
            settings.last_tested_at = datetime.utcnow()

            if test_result["success"]:
                flash(
                    f"Connection test successful: {test_result['message']}", "success"
                )
            else:
                flash(
                    f"Connection test failed: {test_result.get('message', 'Unknown error')}",
                    "error",
                )
        else:
            flash(f"API settings for {provider} saved successfully.", "success")

        db.session.commit()
        return redirect(url_for("admin.api_settings"))

    if edit_provider and edit_provider in settings_dict:
        settings = settings_dict[edit_provider]
        form.provider.data = settings.provider
        form.api_key.data = settings.get_masked_key()  # never send full key to browser
        form.enabled.data = settings.enabled
        form.default_model.data = settings.default_model or ""
        form.max_tokens.data = settings.max_tokens or 4000
        form.temperature.data = settings.temperature or 0.7
        form.jira_url.data = settings.jira_url or ""
        form.jira_email.data = settings.jira_email or ""
        form.hf_model_id.data = settings.hf_model_id or ""
        form.hf_endpoint_url.data = settings.hf_endpoint_url or ""
        form.custom_endpoint_url.data = settings.custom_endpoint_url or ""
        form.custom_auth_method.data = settings.custom_auth_method or "bearer"
        form.custom_headers.data = settings.custom_headers or ""

    return render_template(
        "admin/api_settings.html",
        form=form,
        all_settings=all_settings,
        edit_provider=edit_provider,
    )


@admin_bp.route("/api-settings/<int:settings_id>/delete", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_api_settings_delete")
def delete_api_settings(settings_id):
    """Delete API settings."""
    settings = APISettings.query.get_or_404(settings_id)
    provider = settings.provider
    db.session.delete(settings)
    db.session.commit()
    flash(f"API settings for {provider} deleted successfully.", "success")
    return redirect(url_for("admin.api_settings"))


@admin_bp.route("/api-settings/<int:settings_id>/test", methods=["POST"])
@login_required
@admin_required
@audit_log("test_api_settings")
def test_api_settings(settings_id):
    """Test API settings."""
    settings = APISettings.query.get_or_404(settings_id)

    try:
        result = test_api_key(settings.provider, settings.api_key, model=settings.default_model or None)
        flash(f"API test successful: {result}", "success")
    except Exception as e:
        flash("API test failed. Please try again.", "error")

    return redirect(url_for("admin.api_settings"))


@admin_bp.route("/api-settings/env-keys", methods=["GET"])
@login_required
@admin_required
def preview_env_keys():
    """Preview API keys found in environment variables but not yet in the database."""
    import os

    env_key_map = {
        "OPENAI_API_KEY": "openai",
        "ANTHROPIC_API_KEY": "anthropic",
        "GOOGLE_API_KEY": "gemini",
        "GEMINI_API_KEY": "gemini",
        "DEEPSEEK_API_KEY": "deepseek",
        "HUGGINGFACE_API_KEY": "huggingface",
        "HF_TOKEN": "huggingface",
        "AZURE_OPENAI_API_KEY": "azure",
        "OPENROUTER_API_KEY": "openrouter",
    }

    default_models = {
        "openai": "gpt-4o",
        "anthropic": "claude-3-5-sonnet-20241022",
        "gemini": "gemini-2.0-flash-exp",
        "deepseek": "deepseek-chat",
        "huggingface": "meta-llama/Llama-3.1-8B-Instruct",
        "azure": "gpt-4",
        "openrouter": "google/gemini-2.5-flash-preview:free",
    }

    existing = {s.provider for s in APISettings.query.all()}

    found_keys = []
    for env_var, provider in env_key_map.items():
        value = os.environ.get(env_var, "")
        if value and value.strip():
            masked = value[:8] + "..." + value[-4:] if len(value) > 12 else "****"
            found_keys.append(
                {
                    "env_var": env_var,
                    "provider": provider,
                    "masked_key": masked,
                    "in_database": provider in existing,
                    "default_model": default_models.get(provider, ""),
                }
            )

    return jsonify(
        {
            "success": True,
            "keys": found_keys,
            "count": len(found_keys),
            "new_count": sum(1 for k in found_keys if not k["in_database"]),
        }
    )


@admin_bp.route("/api-settings/update-model", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_api_model_update")
def update_provider_model():
    """Quick-update the default model for a provider without the full form."""
    data = request.get_json()
    if not data or "provider" not in data or "model" not in data:
        return jsonify(
            {"success": False, "error": "provider and model are required"}
        ), 400

    provider = data["provider"]
    model = data["model"].strip()

    settings = APISettings.query.filter_by(provider=provider).first()
    if not settings:
        return jsonify(
            {"success": False, "error": f"No settings found for {provider}"}
        ), 404

    settings.default_model = model
    settings.updated_at = datetime.utcnow()
    settings.updated_by_id = current_user.id
    db.session.commit()

    return jsonify({"success": True, "provider": provider, "model": model})


@admin_bp.route("/api-settings/load-env", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_api_keys_load_env")
def load_env_keys():
    """Import selected API keys from environment variables into the database."""
    import os

    data = request.get_json()
    if not data or "keys" not in data:
        return jsonify({"success": False, "error": "No keys specified"}), 400

    env_key_map = {
        "OPENAI_API_KEY": "openai",
        "ANTHROPIC_API_KEY": "anthropic",
        "GOOGLE_API_KEY": "gemini",
        "GEMINI_API_KEY": "gemini",
        "DEEPSEEK_API_KEY": "deepseek",
        "HUGGINGFACE_API_KEY": "huggingface",
        "HF_TOKEN": "huggingface",
        "AZURE_OPENAI_API_KEY": "azure",
        "OPENROUTER_API_KEY": "openrouter",
    }

    default_models = {
        "openai": "gpt-4o",
        "anthropic": "claude-3-5-sonnet-20241022",
        "gemini": "gemini-2.0-flash-exp",
        "deepseek": "deepseek-chat",
        "huggingface": "meta-llama/Llama-3.1-8B-Instruct",
        "azure": "gpt-4",
        "openrouter": "google/gemini-2.5-flash-preview:free",
    }

    requested_keys = data["keys"]
    update_existing = data.get("update_existing", False)

    imported = []
    skipped = []
    errors = []

    existing_map = {s.provider: s for s in APISettings.query.all()}

    for env_var in requested_keys:
        provider = env_key_map.get(env_var)
        if not provider:
            errors.append(f"Unknown env var: {env_var}")
            continue

        api_key = os.environ.get(env_var, "").strip()
        if not api_key:
            errors.append(f"{env_var} is empty")
            continue

        if provider in existing_map and not update_existing:
            skipped.append(f"{provider} already in database (skipped)")
            continue

        if provider in existing_map:
            existing_map[provider].api_key = api_key
            existing_map[provider].enabled = True
            existing_map[provider].updated_at = datetime.utcnow()
            existing_map[provider].updated_by_id = current_user.id
            imported.append(f"{provider} (updated)")
        else:
            settings = APISettings(
                provider=provider,
                api_key=api_key,
                enabled=True,
                default_model=default_models.get(provider, ""),
                max_tokens=4000,
                temperature=0.7,
                updated_by_id=current_user.id,
            )
            db.session.add(settings)
            imported.append(f"{provider} (created)")

    db.session.commit()

    return jsonify(
        {
            "success": True,
            "imported": imported,
            "skipped": skipped,
            "errors": errors,
            "message": f"Imported {len(imported)} key(s)",
        }
    )


# ============================================================================
# Consolidation Status
# ============================================================================


@admin_bp.route("/consolidation")
@login_required
@admin_required
def consolidation_status():
    """Blueprint Consolidation Status Dashboard"""
    from app.api_registry import (
        API_REGISTRY,
        CONSOLIDATION_TARGETS,
        BlueprintStatus,
        get_consolidation_status,
    )

    status = get_consolidation_status()

    blueprints = []
    for name, bp_info in API_REGISTRY.items():
        bp_data = {
            "name": name,
            "status": bp_info.status.value,
            "url_prefix": bp_info.url_prefix,
            "endpoints_count": bp_info.endpoints_count,
            "consolidate_into": bp_info.consolidate_into,
            "notes": bp_info.notes or "",
            "is_active": bp_info.status == BlueprintStatus.ACTIVE,
            "is_deprecated": bp_info.status == BlueprintStatus.DEPRECATED,
            "is_consolidated": bp_info.status == BlueprintStatus.CONSOLIDATED,
        }
        blueprints.append(bp_data)

    status_order = {"deprecated": 0, "active": 1, "consolidated": 2, "legacy": 3}
    blueprints.sort(key=lambda x: (status_order.get(x["status"], 4), x["name"]))

    targets = []
    for target_name, target_info in CONSOLIDATION_TARGETS.items():
        keep_bp = target_info["keep"]
        consolidate_list = target_info["consolidate"]
        priority = target_info["priority"]

        completed_count = 0
        total_count = len(consolidate_list)

        for bp_name in consolidate_list:
            if bp_name in API_REGISTRY:
                bp_info = API_REGISTRY[bp_name]
                if bp_info.status in [
                    BlueprintStatus.CONSOLIDATED,
                    BlueprintStatus.DEPRECATED,
                ]:
                    completed_count += 1

        targets.append(
            {
                "name": target_name,
                "keep_blueprint": keep_bp,
                "consolidate_blueprints": consolidate_list,
                "priority": priority,
                "completed_count": completed_count,
                "total_count": total_count,
                "progress_percentage": (completed_count / total_count * 100)
                if total_count > 0
                else 100,
                "is_complete": completed_count == total_count,
                "estimated_effort": target_info.get("estimated_effort", "unknown"),
            }
        )

    priority_order = {"HIGH": 0, "MEDIUM": 1, "LOW": 2}
    targets.sort(
        key=lambda x: (priority_order.get(x["priority"], 3), -x["progress_percentage"])
    )

    return render_template(
        "admin/consolidation_status.html",
        status=status,
        blueprints=blueprints,
        targets=targets,
        BlueprintStatus=BlueprintStatus,
    )


# ============================================================================
# Feature Flags Management
# ============================================================================


@admin_bp.route("/feature-flags")
@login_required
@admin_required
def feature_flags():
    """Feature flags management page with pagination."""
    page = request.args.get("page", 1, type=int)
    per_page = request.args.get("per_page", 50, type=int)
    search = request.args.get("search", "")
    filter_type = request.args.get("type", "")
    filter_state = request.args.get("state", "")
    filter_status = request.args.get("status", "")
    filter_hierarchy = request.args.get("hierarchy", "")

    query = FeatureFlag.query

    if search:
        search_term = f"%{search}%"
        query = query.filter(
            db.or_(
                FeatureFlag.name.ilike(search_term),
                FeatureFlag.key.ilike(search_term),
                FeatureFlag.description.ilike(search_term),
            )
        )

    if filter_type:
        query = query.filter(FeatureFlag.feature_type == filter_type)

    if filter_state:
        query = query.filter(FeatureFlag.state == filter_state)

    if filter_status:
        if filter_status == "enabled":
            query = query.filter(FeatureFlag.enabled == True)
        elif filter_status == "disabled":
            query = query.filter(FeatureFlag.enabled == False)

    if filter_hierarchy:
        if filter_hierarchy == "parent":
            query = query.filter(FeatureFlag.parent_id == None)
        elif filter_hierarchy == "child":
            parent_table = aliased(FeatureFlag)
            query = query.join(
                parent_table, FeatureFlag.parent_id == parent_table.id
            ).filter(FeatureFlag.parent_id != None, parent_table.parent_id == None)
        elif filter_hierarchy == "grandchild":
            parent_table = aliased(FeatureFlag)
            query = query.join(
                parent_table, FeatureFlag.parent_id == parent_table.id
            ).filter(FeatureFlag.parent_id != None, parent_table.parent_id != None)

    query = query.order_by(FeatureFlag.sort_order, FeatureFlag.name)
    pagination = query.paginate(page=page, per_page=per_page, error_out=False)

    if pagination.total == 0 and not search:
        flash(
            "No features found. Run 'flask discover-features' to auto-discover routes, "
            "or create features manually.",
            "info",
        )

    return render_template(
        "admin/feature_flags.html",
        features=pagination.items,
        pagination=pagination,
        search=search,
        filter_type=filter_type,
        filter_state=filter_state,
        filter_status=filter_status,
        filter_hierarchy=filter_hierarchy,
        FeatureType=FeatureType,
        FeatureState=FeatureState,
    )


@admin_bp.route("/feature-flags/new", methods=["GET", "POST"])
@login_required
@admin_required
@audit_log("admin_feature_flag_create")
def feature_flag_new():
    """Create new feature flag."""
    form = FeatureFlagForm()

    if form.validate_on_submit():
        try:
            routes_data = None
            if form.routes.data and form.routes.data.strip():
                try:
                    routes_data = json.loads(form.routes.data)
                except json.JSONDecodeError:
                    flash("Invalid JSON format for routes", "error")
                    return render_template(
                        "admin/feature_flag_form.html", form=form, action="New"
                    )

            feature = FeatureFlag(
                key=form.key.data,
                name=form.name.data,
                description=form.description.data,
                feature_type=FeatureType(form.feature_type.data),
                state=FeatureState(form.state.data),
                enabled=form.enabled.data,
                sidebar_label=form.sidebar_label.data,
                sidebar_icon=form.sidebar_icon.data,
                routes=routes_data,
                parent_id=form.parent_id.data or None,
                sort_order=form.sort_order.data or 0,
                last_modified_by=current_user.id,
            )

            db.session.add(feature)
            db.session.commit()

            # Clear cache for this feature
            FeatureFlag.clear_cache(feature.key)

            flash(f"Feature flag '{feature.name}' created successfully", "success")
            return redirect(url_for("admin.feature_flags"))

        except Exception as e:
            db.session.rollback()
            flash("Error creating feature flag. Please try again.", "error")

    return render_template("admin/feature_flag_form.html", form=form, action="New")


@admin_bp.route("/feature-flags/<int:id>/edit", methods=["GET", "POST"])
@login_required
@admin_required
@audit_log("admin_feature_flag_edit")
def feature_flag_edit(id):
    """Edit feature flag."""
    feature = FeatureFlag.query.get_or_404(id)
    form = FeatureFlagForm(obj=feature)

    if form.validate_on_submit():
        try:
            routes_data = None
            if form.routes.data and form.routes.data.strip():
                try:
                    routes_data = json.loads(form.routes.data)
                except json.JSONDecodeError:
                    flash("Invalid JSON format for routes", "error")
                    return render_template(
                        "admin/feature_flag_form.html",
                        form=form,
                        feature=feature,
                        action="Edit",
                    )

            feature.key = form.key.data
            feature.name = form.name.data
            feature.description = form.description.data
            feature.feature_type = FeatureType(form.feature_type.data)
            feature.state = FeatureState(form.state.data)
            feature.enabled = form.enabled.data
            feature.sidebar_label = form.sidebar_label.data
            feature.sidebar_icon = form.sidebar_icon.data
            feature.routes = routes_data
            feature.parent_id = form.parent_id.data or None
            feature.sort_order = form.sort_order.data or 0
            feature.last_modified_by = current_user.id

            db.session.commit()

            # Clear cache for this feature
            FeatureFlag.clear_cache(feature.key)

            flash(f"Feature flag '{feature.name}' updated successfully", "success")
            return redirect(url_for("admin.feature_flags"))

        except Exception as e:
            db.session.rollback()
            flash("Error updating feature flag. Please try again.", "error")

    if feature.routes and not form.routes.data:
        form.routes.data = json.dumps(feature.routes, indent=2)
    elif not feature.routes:
        form.routes.data = ""

    return render_template(
        "admin/feature_flag_form.html", form=form, feature=feature, action="Edit"
    )


@admin_bp.route("/feature-flags/<int:id>/toggle", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_feature_flag_toggle")
def feature_flag_toggle(id):
    """Quick toggle feature enabled/disabled."""
    feature = FeatureFlag.query.get_or_404(id)

    try:
        feature.enabled = not feature.enabled
        feature.last_modified_by = current_user.id
        db.session.commit()

        # Clear cache for this feature
        FeatureFlag.clear_cache(feature.key)

        status = "enabled" if feature.enabled else "disabled"
        return jsonify(
            {
                "success": True,
                "enabled": feature.enabled,
                "message": f"Feature {status}",
            }
        )
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "error": "An internal error occurred"}), 500


@admin_bp.route("/feature-flags/<int:id>/delete", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_feature_flag_delete")
def feature_flag_delete(id):
    """Delete feature flag."""
    feature = FeatureFlag.query.get_or_404(id)
    feature_key = feature.key  # Store key before deletion

    try:
        db.session.delete(feature)
        db.session.commit()

        # Clear cache for this feature
        FeatureFlag.clear_cache(feature_key)

        flash(f"Feature flag '{feature.name}' deleted successfully", "success")
    except Exception as e:
        db.session.rollback()
        flash("Error deleting feature flag. Please try again.", "error")

    return redirect(url_for("admin.feature_flags"))


@admin_bp.route("/feature-flags/discover-sidebar")
@login_required
@admin_required
def feature_flags_discover_sidebar():
    """Discover sidebar menu items for feature flagging."""
    try:
        parser = parse_sidebar_template()
        existing_flags = {ff.key: ff for ff in FeatureFlag.query.all()}

        all_items = []
        for section_title, submenu_title, link in parser.get_all_links():
            feature_key = link.feature_key
            if submenu_title:
                feature_key = f"{SidebarSubmenu(submenu_title, None, section_title, []).feature_key}_{link.feature_key}"

            all_items.append(
                {
                    "section": section_title,
                    "submenu": submenu_title,
                    "text": link.text,
                    "href": link.href,
                    "icon": link.icon,
                    "endpoint": link.endpoint,
                    "feature_key": feature_key,
                    "already_flagged": feature_key in existing_flags,
                    "flag_id": existing_flags[feature_key].id
                    if feature_key in existing_flags
                    else None,
                }
            )

        return render_template(
            "admin/discover_sidebar.html",
            items=all_items,
            sections=parser.sections,
            parser_data=parser.to_dict(),
        )

    except Exception as e:
        flash("Error parsing sidebar. Please try again.", "error")
        return redirect(url_for("admin.feature_flags"))


@admin_bp.route("/feature-flags/discover-sidebar/create", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_feature_flags_bulk_create")
def feature_flags_create_from_sidebar():
    """Create feature flags from selected sidebar items."""
    try:
        selected_keys = request.form.getlist("selected_items")

        if not selected_keys:
            flash("No items selected", "warning")
            return redirect(url_for("admin.feature_flags_discover_sidebar"))

        parser = parse_sidebar_template()

        items_map = {}
        for section_title, submenu_title, link in parser.get_all_links():
            feature_key = link.feature_key
            if submenu_title:
                submenu_obj = SidebarSubmenu(submenu_title, None, section_title, [])
                feature_key = f"{submenu_obj.feature_key}_{link.feature_key}"

            items_map[feature_key] = {
                "section": section_title,
                "submenu": submenu_title,
                "link": link,
            }

        created_count = 0
        skipped_count = 0

        for key in selected_keys:
            existing = FeatureFlag.query.filter_by(key=key).first()
            if existing:
                skipped_count += 1
                continue

            item_data = items_map.get(key)
            if not item_data:
                continue

            link = item_data["link"]
            submenu = item_data["submenu"]
            section = item_data["section"]

            feature = FeatureFlag(
                key=key,
                name=link.text,
                description=f"Controls visibility of '{link.text}' in {section}"
                + (f" > {submenu}" if submenu else ""),
                feature_type=FeatureType.SIDEBAR_LINK,
                state=FeatureState.BETA,
                enabled=True,
                sidebar_label=link.text,
                sidebar_icon=link.icon,
                routes=[link.endpoint] if link.endpoint else [],
                last_modified_by=current_user.id,
            )

            db.session.add(feature)
            created_count += 1

        db.session.commit()

        message = f"Created {created_count} feature flag(s)"
        if skipped_count > 0:
            message += f", skipped {skipped_count} existing"

        flash(message, "success")
        return redirect(url_for("admin.feature_flags"))

    except Exception as e:
        db.session.rollback()
        flash("Error creating feature flags. Please try again.", "error")
        return redirect(url_for("admin.feature_flags_discover_sidebar"))


# ============================================================================
# Abacus Integration Management
# ============================================================================


@admin_bp.route("/abacus-settings", methods=["GET", "POST"])
@login_required
@admin_required
@audit_log("admin_abacus_settings_save")
def abacus_settings():
    """Manage Abacus connector configuration."""
    from flask_wtf import FlaskForm
    from wtforms import BooleanField, IntegerField, PasswordField, StringField
    from wtforms.validators import URL, DataRequired

    from app.models.models import ExternalSystem

    class AbacusSettingsForm(FlaskForm):
        base_url = StringField("Base URL", validators=[DataRequired(), URL()])
        client_id = StringField("Client ID", validators=[DataRequired()])
        client_secret = PasswordField("Client Secret")
        enabled = BooleanField("Enable Integration", default=False)
        sync_enabled = BooleanField("Enable Auto-Sync", default=False)
        sync_interval_minutes = IntegerField("Sync Interval (minutes)", default=1440)
        filter_countries = StringField("Country Filter", default="")

    abacus_config = ExternalSystem.query.filter_by(
        system_name="abacus", system_type="ea_tool"
    ).first()

    form = AbacusSettingsForm()

    if form.validate_on_submit():
        import os

        credentials = {}
        if form.client_secret.data and form.client_secret.data.strip():
            credentials = {
                "client_id": form.client_id.data,
                "client_secret": form.client_secret.data,
            }
        elif abacus_config and abacus_config.credentials:
            try:
                existing_creds = json.loads(abacus_config.credentials)
                credentials = {
                    "client_id": form.client_id.data,
                    "client_secret": existing_creds.get("client_secret", ""),
                }
            except (json.JSONDecodeError, KeyError, TypeError):
                logger.exception("Failed to JSON parsing")
                pass

        if abacus_config:
            abacus_config.base_url = form.base_url.data
            abacus_config.api_endpoint = f"{form.base_url.data}/api"
            abacus_config.auth_type = "oauth2"
            if credentials:
                abacus_config.credentials = json.dumps(credentials)

            filter_countries = (form.filter_countries.data or "").strip()
            config_data = {"filter_countries": filter_countries}
            abacus_config.config_json = json.dumps(config_data)
            os.environ["ABACUS_FILTER_COUNTRIES"] = filter_countries

            abacus_config.enabled = form.enabled.data
            abacus_config.sync_enabled = form.sync_enabled.data
            abacus_config.sync_interval_minutes = form.sync_interval_minutes.data
            abacus_config.updated_by_id = current_user.id
            abacus_config.updated_at = datetime.utcnow()

            flash("Abacus settings updated successfully.", "success")
        else:
            if not credentials or not credentials.get("client_secret"):
                flash("Client secret is required for new configuration.", "error")
                return render_template(
                    "admin/abacus_settings.html", form=form, abacus_config=None
                )

            filter_countries = (form.filter_countries.data or "").strip()
            config_data = {"filter_countries": filter_countries}
            os.environ["ABACUS_FILTER_COUNTRIES"] = filter_countries

            abacus_config = ExternalSystem(
                system_name="abacus",
                system_type="ea_tool",
                base_url=form.base_url.data,
                api_endpoint=f"{form.base_url.data}/api",
                auth_type="oauth2",
                credentials=json.dumps(credentials),
                config_json=json.dumps(config_data),
                enabled=form.enabled.data,
                sync_enabled=form.sync_enabled.data,
                sync_interval_minutes=form.sync_interval_minutes.data,
                updated_by_id=current_user.id,
            )
            db.session.add(abacus_config)

            flash("Abacus settings created successfully.", "success")

        db.session.commit()
        return redirect(url_for("admin.abacus_settings"))

    env_config = None
    imported_stats = {"applications": 0, "capabilities": 0}

    if request.method == "GET":
        if abacus_config or env_config:
            try:
                from app.models.application_portfolio import ApplicationComponent
                from app.models.business_capabilities import BusinessCapability

                abacus_apps = ApplicationComponent.query.filter_by(
                    abacus_source=True
                ).count()
                imported_stats["applications"] = abacus_apps

                abacus_caps = BusinessCapability.query.filter(
                    BusinessCapability.archimate_id.isnot(None),
                    BusinessCapability.archimate_id.like("ABACUS-%"),
                ).count()
                imported_stats["capabilities"] = abacus_caps
            except Exception as e:
                logger.error(f"Failed to fetch Abacus imported data counts: {e}")

        if abacus_config:
            form.base_url.data = abacus_config.base_url
            form.enabled.data = abacus_config.enabled
            form.sync_enabled.data = abacus_config.sync_enabled
            form.sync_interval_minutes.data = (
                abacus_config.sync_interval_minutes or 1440
            )

            if abacus_config.config_json:
                try:
                    config_data = json.loads(abacus_config.config_json)
                    form.filter_countries.data = config_data.get("filter_countries", "")
                except (json.JSONDecodeError, KeyError, TypeError):
                    form.filter_countries.data = ""

            if abacus_config.credentials:
                try:
                    creds = json.loads(abacus_config.credentials)
                    form.client_id.data = creds.get("client_id", "")
                except (json.JSONDecodeError, KeyError, TypeError):
                    logger.exception("Failed to JSON parsing")
                    pass
        else:
            import os

            env_base_url = os.getenv("ABACUS_BASE_URL")
            env_client_id = os.getenv("ABACUS_CLIENT_ID")
            env_client_secret = os.getenv("ABACUS_CLIENT_SECRET")
            env_enabled = os.getenv("ABACUS_ENABLED", "false").lower() == "true"
            env_filter_countries = os.getenv("ABACUS_FILTER_COUNTRIES", "")

            if all([env_base_url, env_client_id, env_client_secret]):
                form.base_url.data = env_base_url
                form.client_id.data = env_client_id
                form.enabled.data = env_enabled
                form.sync_enabled.data = False
                form.sync_interval_minutes.data = 1440
                form.filter_countries.data = env_filter_countries

                env_config = {
                    "base_url": env_base_url,
                    "client_id": env_client_id,
                    "enabled": env_enabled,
                    "filter_countries": env_filter_countries,
                }

    # Load cached filter dimensions (populated by sync or discover button)
    filter_dimensions = {}
    if abacus_config and abacus_config.config_json:
        try:
            cached = json.loads(abacus_config.config_json)
            filter_dimensions = cached.get("_cached_dimensions", {})
        except (json.JSONDecodeError, TypeError):
            logger.exception("Failed to JSON parsing")
            pass

    return render_template(
        "admin/abacus_settings.html",
        form=form,
        abacus_config=abacus_config,
        env_config=env_config,
        imported_stats=imported_stats,
        filter_dimensions=filter_dimensions,
    )


@admin_bp.route("/abacus-settings/test-connection", methods=["POST"])
@login_required
@admin_required
@audit_log("test_abacus_connection")
def test_abacus_connection():
    """Test Abacus connection."""
    import asyncio
    import os

    from app.models.models import ExternalSystem

    abacus_config = ExternalSystem.query.filter_by(
        system_name="abacus", system_type="ea_tool"
    ).first()

    if not abacus_config:
        env_base_url = os.getenv("ABACUS_BASE_URL")
        env_client_id = os.getenv("ABACUS_CLIENT_ID")
        env_client_secret = os.getenv("ABACUS_CLIENT_SECRET")

        if not all([env_base_url, env_client_id, env_client_secret]):
            flash(
                "No Abacus configuration found in database or .env file. Please configure first.",
                "error",
            )
            return redirect(url_for("admin.abacus_settings"))

        class TempConfig:
            def __init__(self):
                self.base_url = env_base_url
                self.credentials = None

        abacus_config = TempConfig()
        creds = {"client_id": env_client_id, "client_secret": env_client_secret}
    else:
        creds = (
            json.loads(abacus_config.credentials) if abacus_config.credentials else {}
        )

    try:
        from app.connectors.abacus import create_abacus_connector
        from app.services.connector_framework import ConnectorConfig

        config = ConnectorConfig(
            name="abacus",
            connector_type="ea_tool",
            config={
                "base_url": abacus_config.base_url,
                "client_id": creds.get("client_id"),
                "client_secret": creds.get("client_secret"),
            },
        )

        connector = create_abacus_connector(config)

        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        result = loop.run_until_complete(connector.test_connection())
        loop.close()

        if hasattr(abacus_config, "last_connection_test"):
            abacus_config.last_connection_test = datetime.utcnow()
            abacus_config.connection_status = "connected" if result else "error"
            abacus_config.last_error = None if result else "Connection test failed"
            db.session.commit()

        if result:
            flash("Connection test successful! Abacus API is accessible.", "success")
        else:
            flash(
                "Connection test failed. Please check your credentials and URL.",
                "error",
            )

    except Exception as e:
        logger.error(f"Abacus connection test exception: {e}", exc_info=True)

        if hasattr(abacus_config, "last_connection_test"):
            abacus_config.last_connection_test = datetime.utcnow()
            abacus_config.connection_status = "error"
            abacus_config.last_error = str(e)
            db.session.commit()

        flash("Connection test failed. Please try again.", "error")

    return redirect(url_for("admin.abacus_settings"))


@admin_bp.route("/abacus-settings/trigger-sync", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_abacus_sync_trigger")
def trigger_abacus_sync():
    """Trigger manual Abacus synchronization."""
    import os

    from app.models.models import ExternalSystem

    abacus_config = ExternalSystem.query.filter_by(system_name="abacus").first()

    if not abacus_config:
        env_base_url = os.getenv("ABACUS_BASE_URL")
        env_client_id = os.getenv("ABACUS_CLIENT_ID")
        env_client_secret = os.getenv("ABACUS_CLIENT_SECRET")

        if not all([env_base_url, env_client_id, env_client_secret]):
            flash(
                "No Abacus configuration found in database or .env file. Please configure first.",
                "error",
            )
            return redirect(url_for("admin.abacus_settings"))

        flash("Using Abacus credentials from .env file.", "info")
    elif not abacus_config.enabled:
        flash("Abacus integration is disabled. Enable it first.", "error")
        return redirect(url_for("admin.abacus_settings"))

    try:
        from app.services.job_queue_service import get_job_queue_service

        job_queue = get_job_queue_service()
        job = job_queue.create_job(
            name="Abacus Full Sync", task="abacus_sync", payload={"sync_type": "full"}
        )

        flash(
            f"Sync job created (ID: {job.id}). "
            "The sync will run in the background. Check the status below.",
            "info",
        )

    except Exception as e:
        flash("Failed to create sync job. Please try again.", "error")

    return redirect(url_for("admin.abacus_settings"))


@admin_bp.route("/abacus-settings/sync-status", methods=["GET"])
@login_required
@admin_required
def abacus_sync_status():
    """API endpoint to check current sync job status."""
    from app.models import Job

    latest_job = (
        Job.query.filter_by(task="abacus_sync").order_by(Job.created_at.desc()).first()
    )

    if not latest_job:
        return jsonify({"status": "no_jobs", "message": "No sync jobs found"})

    return jsonify(
        {
            "job_id": latest_job.id,
            "status": latest_job.status,
            "created_at": latest_job.created_at.isoformat()
            if latest_job.created_at
            else None,
            "started_at": latest_job.started_at.isoformat()
            if latest_job.started_at
            else None,
            "finished_at": latest_job.finished_at.isoformat()
            if latest_job.finished_at
            else None,
            "result": latest_job.result,
            "error": latest_job.error,
            "name": latest_job.name,
        }
    )


@admin_bp.route("/abacus-settings/cancel-job/<int:job_id>", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_abacus_job_cancel")
def cancel_abacus_job(job_id):
    """Cancel a running or pending Abacus sync job."""
    try:
        from app.services.job_queue_service import get_job_queue_service

        service = get_job_queue_service()
        success = service.cancel_job(job_id)

        if success:
            return (
                jsonify(
                    {"success": True, "message": f"Job {job_id} cancelled successfully"}
                ),
                200,
            )
        else:
            return (
                jsonify(
                    {
                        "success": False,
                        "message": f"Job {job_id} is already finished or cannot be cancelled",
                    }
                ),
                400,
            )
    except Exception as e:
        logger.error(f"Error cancelling job {job_id}: {e}", exc_info=True)
        return jsonify({"success": False, "message": "An internal error occurred"}), 500


@admin_bp.route("/abacus-settings/stats", methods=["GET"])
@login_required
@admin_required
def abacus_stats():
    """Get Abacus import statistics."""
    try:
        from app.models.application_portfolio import ApplicationComponent
        from app.models.business_capabilities import BusinessCapability

        abacus_apps = ApplicationComponent.query.filter_by(abacus_source=True).count()

        abacus_caps = BusinessCapability.query.filter(
            BusinessCapability.discovery_source == "abacus",
        ).count()

        return (
            jsonify(
                {
                    "success": True,
                    "applications": abacus_apps,
                    "capabilities": abacus_caps,
                }
            ),
            200,
        )

    except Exception as e:
        logger.error(f"Error fetching Abacus stats: {e}", exc_info=True)
        return (
            jsonify(
                {
                    "success": False,
                    "applications": 0,
                    "capabilities": 0,
                    "error": str(e),
                }
            ),
            500,
        )




# ============================================================================
# Governance Gates Management (GOV-03)
# ============================================================================


@admin_bp.route("/governance-gates")
@login_required
@admin_required
def governance_gates():
    """Governance gates configuration page."""
    from app.modules.solutions_strategic.v2.services.governance_gate_service import (
        DEFAULT_GATES,
    )

    return render_template(
        "admin/governance_gates.html",
        default_gates=DEFAULT_GATES,
    )


@admin_bp.route("/api/governance-gates", methods=["GET"])
@login_required
@admin_required
def governance_gates_list():
    """List all governance gates from DB."""
    from app.models.governance_gates import GovernanceGate

    try:
        gates = GovernanceGate.query.order_by(GovernanceGate.gate_name).all()
        return jsonify({"success": True, "gates": [g.to_dict() for g in gates]})
    except Exception as e:
        logger.error("Error listing governance gates: %s", e, exc_info=True)
        return jsonify({"success": False, "error": "Failed to load gates"}), 500


@admin_bp.route("/api/governance-gates", methods=["POST"])
@login_required
@admin_required
def governance_gates_create():
    """Create a new governance gate."""
    from app.models.governance_gates import GovernanceGate

    data = request.get_json()
    if not data or not data.get("gate_name"):
        return jsonify({"success": False, "error": "gate_name is required"}), 400

    gate_name = data["gate_name"].strip()
    if GovernanceGate.query.filter_by(gate_name=gate_name).first():
        return jsonify({"success": False, "error": f"Gate '{gate_name}' already exists"}), 409

    try:
        gate = GovernanceGate(
            gate_name=gate_name,
            description=data.get("description", ""),
            required_sections=data.get("required_sections", []),
            min_completeness=int(data.get("min_completeness", 60)),
            required_decisions_count=int(data.get("required_decisions_count", 0)),
            require_risk_mitigations=bool(data.get("require_risk_mitigations", False)),
            enabled=bool(data.get("enabled", True)),
        )
        db.session.add(gate)
        db.session.commit()
        return jsonify({"success": True, "gate": gate.to_dict()}), 201
    except Exception as e:
        db.session.rollback()
        logger.error("Error creating governance gate: %s", e, exc_info=True)
        return jsonify({"success": False, "error": "Failed to create gate"}), 500


@admin_bp.route("/api/governance-gates/<int:gate_id>", methods=["PUT"])
@login_required
@admin_required
def governance_gates_update(gate_id):
    """Update an existing governance gate."""
    from app.models.governance_gates import GovernanceGate

    gate = GovernanceGate.query.get_or_404(gate_id)
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "error": "No data provided"}), 400

    try:
        if "gate_name" in data:
            new_name = data["gate_name"].strip()
            existing = GovernanceGate.query.filter_by(gate_name=new_name).first()
            if existing and existing.id != gate_id:
                return jsonify({"success": False, "error": f"Gate '{new_name}' already exists"}), 409
            gate.gate_name = new_name
        if "description" in data:
            gate.description = data["description"]
        if "required_sections" in data:
            gate.required_sections = data["required_sections"]
        if "min_completeness" in data:
            gate.min_completeness = int(data["min_completeness"])
        if "required_decisions_count" in data:
            gate.required_decisions_count = int(data["required_decisions_count"])
        if "require_risk_mitigations" in data:
            gate.require_risk_mitigations = bool(data["require_risk_mitigations"])
        if "enabled" in data:
            gate.enabled = bool(data["enabled"])

        db.session.commit()
        return jsonify({"success": True, "gate": gate.to_dict()})
    except Exception as e:
        db.session.rollback()
        logger.error("Error updating governance gate %s: %s", gate_id, e, exc_info=True)
        return jsonify({"success": False, "error": "Failed to update gate"}), 500


@admin_bp.route("/api/governance-gates/<int:gate_id>", methods=["DELETE"])
@login_required
@admin_required
def governance_gates_delete(gate_id):
    """Soft-delete a governance gate by disabling it."""
    from app.models.governance_gates import GovernanceGate

    gate = GovernanceGate.query.get_or_404(gate_id)
    try:
        gate.enabled = False
        db.session.commit()
        return jsonify({"success": True, "message": f"Gate '{gate.gate_name}' disabled"})
    except Exception as e:
        db.session.rollback()
        logger.error("Error disabling governance gate %s: %s", gate_id, e, exc_info=True)
        return jsonify({"success": False, "error": "Failed to disable gate"}), 500


@admin_bp.route("/abacus-settings/discover-filters", methods=["POST"])
@login_required
@admin_required
def discover_abacus_filters():
    """Discover available filter dimensions from the Abacus API."""
    import asyncio

    from app.models.models import ExternalSystem

    try:
        abacus_config = ExternalSystem.query.filter_by(system_name="abacus").first()
        if not abacus_config:
            return jsonify({"success": False, "error": "No Abacus configuration found"}), 400

        credentials = json.loads(abacus_config.credentials)

        class TempConfig:
            def __init__(self):
                self.config = {
                    "base_url": abacus_config.base_url,
                    "client_id": credentials.get("client_id"),
                    "client_secret": credentials.get("client_secret"),
                }

        from app.connectors.abacus import AbacusConnector

        connector = AbacusConnector(TempConfig())

        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            options = loop.run_until_complete(connector.discover_filter_options())
        finally:
            loop.close()

        # Cache results so page shows them without re-fetching
        if options:
            config = {}
            if abacus_config.config_json:
                try:
                    config = json.loads(abacus_config.config_json)
                except (json.JSONDecodeError, TypeError):
                    config = {}
            config["_cached_dimensions"] = options
            abacus_config.config_json = json.dumps(config)
            db.session.commit()

        return jsonify({"success": True, "dimensions": options})

    except Exception as e:
        logger.error("discover_abacus_filters failed: %s", e, exc_info=True)
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp.route("/abacus-dashboard", methods=["GET"])
@login_required
@admin_required
def abacus_dashboard():
    """Display Abacus sync dashboard with health metrics and statistics."""
    from app.models.application_portfolio import ApplicationComponent
    from app.models.business_capabilities import BusinessCapability
    from app.models.models import ExternalSystem
    from app.services.abacus_sync_service import get_sync_service

    sync_service = get_sync_service()
    sync_status = sync_service.get_sync_status()

    stats = {
        "applications_synced": 0,
        "applications_created": 0,
        "applications_updated": 0,
        "applications_skipped": 0,
        "capabilities_synced": 0,
        "capabilities_created": 0,
        "capabilities_updated": 0,
        "capabilities_skipped": 0,
    }

    try:
        abacus_apps = ApplicationComponent.query.filter_by(abacus_source=True).count()
        stats["applications_synced"] = abacus_apps

        abacus_caps = BusinessCapability.query.filter(
            BusinessCapability.discovery_source == "abacus",
        ).count()
        stats["capabilities_synced"] = abacus_caps

        external_system = ExternalSystem.query.filter_by(system_name="abacus").first()
        if external_system and external_system.metadata:
            try:
                metadata = (
                    json.loads(external_system.metadata)
                    if isinstance(external_system.metadata, str)
                    else external_system.metadata
                )
                last_stats = metadata.get("last_sync_stats", {})
                if last_stats:
                    stats.update(
                        {
                            "applications_created": last_stats.get(
                                "applications_created", 0
                            ),
                            "applications_updated": last_stats.get(
                                "applications_updated", 0
                            ),
                            "applications_skipped": last_stats.get(
                                "applications_skipped", 0
                            ),
                            "capabilities_created": last_stats.get(
                                "capabilities_created", 0
                            ),
                            "capabilities_updated": last_stats.get(
                                "capabilities_updated", 0
                            ),
                            "capabilities_skipped": last_stats.get(
                                "capabilities_skipped", 0
                            ),
                        }
                    )
            except (json.JSONDecodeError, TypeError, AttributeError):
                logger.exception("Failed to operation")
                pass

    except Exception as e:
        logger.error(f"Failed to fetch sync statistics: {e}", exc_info=True)

    return render_template(
        "admin/abacus_dashboard.html", sync_status=sync_status, stats=stats
    )


# ============================================================================
# Seed Management
# ============================================================================


@admin_bp.route("/seed-management")
@login_required
@admin_required
def seed_management():
    """Seed management dashboard."""
    from app.services.seed_management_service import SeedManagementService

    service = SeedManagementService()
    categories = service.get_seed_status()

    return render_template("admin/seed_management.html", categories=categories)


@admin_bp.route("/api/seed-status")
@login_required
@admin_required
def seed_status():
    """API: Get current seed status."""
    from app.services.seed_management_service import SeedManagementService

    service = SeedManagementService()
    categories = service.get_seed_status()

    return jsonify({"success": True, "categories": categories})


@admin_bp.route("/api/seed/<key>", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_seed_run")
def seed(key):
    """API: Trigger seeding for a specific category."""
    from app.services.seed_management_service import SeedManagementService

    service = SeedManagementService()
    result = service.seed_category(key)
    status_code = 200 if result.get("success") else 400
    return jsonify(result), status_code


@admin_bp.route("/api/seed-all", methods=["POST"])
@login_required
@admin_required
@audit_log("admin_seed_all")
def seed_all():
    """API: Trigger all seeders."""
    from app.services.seed_management_service import SeedManagementService

    service = SeedManagementService()
    result = service.seed_all()
    status_code = 200 if result.get("success") else 400
    return jsonify(result), status_code


# ============================================================================
# Feature Auto-Discovery Helper
# ============================================================================


def _auto_discover_features(app):
    """Auto-discover features from Flask app blueprints and routes.

    Creates a hierarchical feature tree:
    - Level 1: Blueprints (parent features)
    - Level 2: Routes within blueprints (child features)
    - Level 3: Sub-routes (grandchild features)

    Filters out infrastructure routes to reduce noise.

    Args:
        app: Flask application instance

    Returns:
        int: Number of features created
    """
    from collections import defaultdict

    created_count = 0
    blueprint_features = {}

    existing_keys = {f.key for f in FeatureFlag.query.all()}

    SKIP_ENDPOINTS = {
        "static",
        "_debug_toolbar",
        "debugtoolbar",
        "health",
        "metrics",
        "status",
    }
    SKIP_PREFIXES = ("_", "static", "/health", "/metrics", "/_", "/api/")
    SKIP_BLUEPRINTS = {"debugtoolbar", "static"}

    routes_by_blueprint = defaultdict(list)
    skipped_count = 0
    total_count = 0

    for rule in app.url_map.iter_rules():
        total_count += 1

        if any(rule.endpoint.startswith(skip) for skip in SKIP_ENDPOINTS):
            skipped_count += 1
            continue
        if any(rule.rule.startswith(skip) for skip in SKIP_PREFIXES):
            skipped_count += 1
            continue

        parts = rule.endpoint.split(".")
        if len(parts) >= 2:
            blueprint_name = parts[0]

            if blueprint_name in SKIP_BLUEPRINTS:
                skipped_count += 1
                continue

            routes_by_blueprint[blueprint_name].append(
                {
                    "endpoint": rule.endpoint,
                    "rule": rule.rule,
                    "methods": list(rule.methods - {"HEAD", "OPTIONS"}),
                }
            )
        else:
            skipped_count += 1

    app.logger.info(
        f"Route discovery: {total_count} total routes, {skipped_count} skipped, {len(routes_by_blueprint)} blueprints found"
    )

    for blueprint_name in sorted(routes_by_blueprint.keys()):
        blueprint_key = f"blueprint_{blueprint_name}"

        if blueprint_key in existing_keys:
            parent_feature = FeatureFlag.query.filter_by(key=blueprint_key).first()
            blueprint_features[blueprint_name] = parent_feature
            continue

        display_name = blueprint_name.replace("_", " ").title()

        parent_feature = FeatureFlag(
            key=blueprint_key,
            name=f"{display_name} Module",
            description=f"All features within the {display_name} blueprint",
            feature_type=FeatureType.BLUEPRINT,
            state=FeatureState.STABLE,
            enabled=True,
            sort_order=created_count,
        )

        try:
            db.session.add(parent_feature)
            db.session.flush()
            blueprint_features[blueprint_name] = parent_feature
            existing_keys.add(blueprint_key)
            created_count += 1
        except Exception as e:
            db.session.rollback()
            app.logger.warning(
                f"Failed to create blueprint feature {blueprint_key}: {e}"
            )
            parent_feature = FeatureFlag.query.filter_by(key=blueprint_key).first()
            if parent_feature:
                blueprint_features[blueprint_name] = parent_feature
            continue

        base_paths = defaultdict(list)
        for route_info in routes_by_blueprint[blueprint_name]:
            rule = route_info["rule"]
            path_parts = [p for p in rule.split("/") if p and not p.startswith("<")]
            base_path = "/" + "/".join(path_parts[:3]) if len(path_parts) >= 2 else rule
            base_paths[base_path].append(route_info)

        for base_path, route_list in sorted(base_paths.items()):
            if len(route_list) == 1 and route_list[0]["rule"] == base_path:
                route_info = route_list[0]
                route_key = f"route_{route_info['endpoint'].replace('.', '_')}"

                if route_key in existing_keys:
                    continue

                function_name = route_info["endpoint"].split(".")[-1]
                display_name = function_name.replace("_", " ").title()

                route_feature = FeatureFlag(
                    key=route_key,
                    name=display_name,
                    description=f"{display_name} ({route_info['rule']})",
                    feature_type=FeatureType.ROUTE,
                    state=FeatureState.STABLE,
                    enabled=True,
                    routes=[route_info["rule"]],
                    parent_id=parent_feature.id,
                    sort_order=created_count,
                )

                try:
                    db.session.add(route_feature)
                    existing_keys.add(route_key)
                    created_count += 1
                except Exception as e:
                    db.session.rollback()
                    app.logger.warning(
                        f"Failed to create route feature {route_key}: {e}"
                    )
                    continue
            else:
                section_key = (
                    f"section_{blueprint_name}_{base_path.replace('/', '_').strip('_')}"
                )

                section_feature = None
                if section_key in existing_keys:
                    section_feature = FeatureFlag.query.filter_by(
                        key=section_key
                    ).first()
                else:
                    section_name = (
                        base_path.split("/")[-1]
                        .replace("-", " ")
                        .replace("_", " ")
                        .title()
                    )

                    section_feature = FeatureFlag(
                        key=section_key,
                        name=f"{section_name} Section",
                        description=f"All routes under {base_path}",
                        feature_type=FeatureType.SIDEBAR_SECTION,
                        state=FeatureState.STABLE,
                        enabled=True,
                        routes=[base_path + "*"],
                        parent_id=parent_feature.id,
                        sort_order=created_count,
                    )

                    try:
                        db.session.add(section_feature)
                        db.session.flush()
                        existing_keys.add(section_key)
                        created_count += 1
                    except Exception as e:
                        db.session.rollback()
                        app.logger.warning(
                            f"Failed to create section feature {section_key}: {e}"
                        )
                        section_feature = FeatureFlag.query.filter_by(
                            key=section_key
                        ).first()

                if not section_feature:
                    continue

                for route_info in route_list:
                    route_key = f"route_{route_info['endpoint'].replace('.', '_')}"

                    if route_key in existing_keys:
                        continue

                    function_name = route_info["endpoint"].split(".")[-1]
                    display_name = function_name.replace("_", " ").title()

                    route_feature = FeatureFlag(
                        key=route_key,
                        name=display_name,
                        description=f"{display_name} ({route_info['rule']}) - {', '.join(route_info['methods'])}",
                        feature_type=FeatureType.ROUTE,
                        state=FeatureState.STABLE,
                        enabled=True,
                        routes=[route_info["rule"]],
                        parent_id=section_feature.id,
                        sort_order=created_count,
                    )

                    try:
                        db.session.add(route_feature)
                        existing_keys.add(route_key)
                        created_count += 1
                    except Exception as e:
                        db.session.rollback()
                        app.logger.warning(
                            f"Failed to create route feature {route_key}: {e}"
                        )
                        continue

    try:
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"Failed to commit feature discovery: {e}")
        raise

    return created_count


@admin_bp.route("/api/users", methods=["GET"])
@login_required
@admin_required
def api_list_users():
    """Paginated user list API for data table."""
    page = request.args.get("page", 1, type=int)
    per_page = min(request.args.get("per_page", 25, type=int), 100)
    search = request.args.get("q") or request.args.get("search", "")
    role_filter = request.args.get("role", "")
    sort_by = request.args.get("sort", "id")
    sort_dir = request.args.get("dir", "desc")

    ALLOWED_SORT = {"id", "first_name", "last_name", "email", "created_at"}
    if sort_by not in ALLOWED_SORT:
        sort_by = "id"

    from sqlalchemy.orm import joinedload
    query = User.query.options(joinedload(User.role))
    if search:
        term = f"%{search}%"
        query = query.filter(
            User.first_name.ilike(term) | User.last_name.ilike(term) | User.email.ilike(term)
        )
    if role_filter:
        query = query.join(User.role).filter(Role.name == role_filter)

    sort_col = getattr(User, sort_by, User.id)
    if sort_dir == "asc":
        query = query.order_by(sort_col.asc())
    else:
        query = query.order_by(sort_col.desc())

    paginated = query.paginate(page=page, per_page=per_page, error_out=False)

    items = []
    for idx, u in enumerate(paginated.items):
        items.append({
            "id": u.id,
            "row_number": (page - 1) * per_page + idx + 1,
            "first_name": u.first_name or "",
            "last_name": u.last_name or "",
            "email": u.email,
            "role": u.role.name if u.role else "No Role",
            "confirmed": u.confirmed,
        })

    return jsonify({
        "users": items,
        "total": paginated.total,
        "page": page,
        "pages": paginated.pages,
        "per_page": per_page,
    })


@admin_bp.route("/api/users/bulk", methods=["DELETE"])
@login_required
@admin_required
def api_bulk_delete_users():
    """Bulk delete users by IDs.

    Handles the ~347 NO ACTION FK references to the users table:
    - Personal/junction tables (user_roles, user_preferences, notifications) → DELETE rows
    - All content tables (created_by_id, owner_id, updated_by_id, etc.) → SET NULL
    - Audit/log tables (audit_logs, ai_interaction_logs, etc.) → SET NULL (preserve for audit trail)

    Uses per-user savepoints so a failure on one user rolls back only that user.
    """
    data = request.get_json() or {}
    ids = data.get("ids", [])
    if not ids or not isinstance(ids, list):
        return jsonify({"error": "ids list required"}), 400

    from flask_login import current_user as cu
    safe_ids = [int(i) for i in ids if int(i) != cu.id]
    if not safe_ids:
        return jsonify({"deleted": 0})

    # Tables where rows BELONG to the user → delete them entirely
    _DELETE_OWNED = [
        ("user_roles", "user_id"),
        ("user_preferences", "user_id"),
        ("notifications", "user_id"),
    ]

    # Tables with multiple nullable FK columns referencing users
    # Each entry: (table_name, [col1, col2, ...])
    # All of these preserve the row content (SET NULL), critical for audit trails and org records.
    _SET_NULL_MULTI = [
        ("adm_board_portfolio_links", ["added_by_id"]),
        ("adm_board_portfolios", ["created_by_id", "portfolio_manager_id", "portfolio_owner_id"]),
        ("adm_compliance_checkpoints", ["completed_by_id", "verified_by_id"]),
        ("adm_cross_board_dependencies", ["created_by_id", "resolved_by_id"]),
        ("adm_deliverables", ["approved_by_id", "created_by_id"]),
        ("adm_phase_approvals", ["decided_by_id", "requested_by_id", "reviewer_id"]),
        ("adm_rida_logs", ["assigned_to_id", "created_by_id", "owner_id", "resolved_by_id"]),
        ("adm_stakeholder_concurrences", ["stakeholder_user_id"]),
        ("adm_transition_history", ["transitioned_by_id"]),
        ("ai_chat_audit_logs", ["user_id"]),
        ("ai_chat_crud_approvals", ["approved_by_id", "user_id"]),
        ("ai_chat_document_uploads", ["uploaded_by_id"]),
        ("ai_interaction_logs", ["user_id"]),
        ("ai_recommendations", ["user_id"]),
        ("ai_suggestions", ["reviewed_by_id"]),
        ("analysis_audit_logs", ["performed_by_id"]),
        ("analysis_recommendations", ["reviewed_by_id"]),
        ("analysis_scenarios", ["created_by_id"]),
        ("api_settings", ["updated_by_id"]),
        ("application_capability_mapping", ["created_by_id", "updated_by_id"]),
        ("application_costs", ["created_by_id"]),
        ("application_disposition_records", ["approved_by_id", "created_by_id"]),
        ("application_import_history", ["imported_by_id"]),
        ("application_replacements", ["created_by_id"]),
        ("application_roi", ["created_by_id"]),
        ("application_versioning", ["created_by_id"]),
        ("arb_adversarial_reviews", ["acknowledged_by_id", "assigned_by_id", "reviewer_id"]),
        ("arb_board_members", ["user_id"]),
        ("arb_compliance_checks", ["overridden_by_id"]),
        ("arb_conditions", ["fulfilled_by_id", "verified_by_id", "waived_by_id"]),
        ("arb_exceptions", ["approved_by_id", "denied_by_id", "requested_by_id", "reviewed_by_id", "revoked_by_id"]),
        ("arb_governance_standards", ["owner_id"]),
        ("arb_review_comments", ["resolved_by_id", "user_id"]),
        ("arb_review_items", ["decided_by_id", "reviewer_id", "submitter_id"]),
        ("arb_submission_packs", ["approved_by_id", "created_by_id"]),
        ("archimate_audit_log", ["user_id"]),
        ("archimate_element_comments", ["user_id"]),
        ("archimate_patterns", ["created_by"]),
        ("archimate_resources", ["owner_id"]),
        ("archimate_viewpoint_templates", ["created_by"]),
        ("archimate_viewpoints", ["created_by_id"]),
        ("architecture_decisions", ["created_by_id"]),
        ("architecture_documents", ["uploaded_by_id"]),
        ("architecture_models", ["user_id"]),
        ("architecture_patterns", ["approved_by_id", "created_by_id"]),
        ("architecture_policies", ["owner_id"]),
        ("architecture_review_boards", ["chair_id", "secretary_id"]),
        ("architecture_review_findings", ["approved_by_id", "created_by_id", "resolved_by_id"]),
        ("architecture_sessions", ["user_id"]),
        ("architecture_vision_documents", ["approved_by_id", "created_by_id"]),
        ("audit_logs", ["user_id"]),
        ("audit_trail", ["user_id"]),
        ("batch_import_batch", ["reviewed_by_id"]),
        ("batch_import_element", ["approved_by_id"]),
        ("batch_import_job", ["user_id"]),
        ("batch_job_errors", ["resolved_by_id"]),
        ("batch_jobs", ["created_by_id"]),
        ("business_processes", ["created_by_id"]),
        ("capability_archimate_classifications", ["validated_by_id"]),
        ("capability_assessments", ["assessor_id"]),
        ("capability_cost_allocations", ["created_by_id"]),
        ("capability_health_overrides", ["created_by_id"]),
        ("capability_investment_plans", ["approved_by_id", "created_by_id"]),
        ("capability_sets", ["user_id"]),
        ("capability_taxonomy_audit", ["changed_by_id"]),
        ("change_management_records", ["approved_by_id", "created_by_id"]),
        ("chat_message_embeddings", ["user_id"]),
        ("code_reviews", ["reviewer_id"]),
        ("code_templates", ["approved_by_id", "created_by_id"]),
        ("codegen_generation_history", ["generated_by_id"]),
        ("codegen_system_boundaries", ["created_by_id"]),
        ("codegen_template_sets", ["created_by_id"]),
        ("compliance_checks", ["reviewed_by_id"]),
        ("compliance_gaps", ["assigned_to_id", "identified_by_id"]),
        ("compliance_governance_reports", ["approved_by_id", "created_by_id"]),
        ("compliance_requirements", ["created_by_id", "verified_by_id", "waiver_approved_by_id"]),
        ("compliance_scan_reports", ["approved_by_id", "created_by_id"]),
        ("compliance_violations", ["assigned_to_id", "created_by_id"]),
        ("conceptual_data_models", ["created_by_id"]),
        ("confidence_thresholds", ["created_by_id"]),
        ("consolidation_candidates", ["reviewed_by"]),
        ("consolidation_opportunities", ["owner_id"]),
        ("cost_budgets", ["approver_id"]),
        ("cost_transactions", ["created_by_id"]),
        ("courses_of_action", ["created_by_id", "sponsor_id"]),
        ("custom_field_definitions", ["created_by"]),
        ("data_access_controls", ["granted_by_id", "user_id"]),
        ("data_catalogs", ["created_by_id"]),
        ("data_domains", ["created_by_id"]),
        ("data_entities", ["created_by_id"]),
        ("data_governance_workflows", ["approver_id", "data_steward_id", "requester_id"]),
        ("data_lineage", ["created_by_id"]),
        ("data_quality_metrics", ["created_by_id"]),
        ("data_retention_policies", ["created_by_id"]),
        ("data_transformations", ["created_by_id"]),
        ("decision_authority", ["user_id"]),
        ("deployment_pipelines", ["created_by_id"]),
        ("derivation_audit", ["reviewed_by"]),
        ("derivation_session", ["initiated_by"]),
        ("design_patterns", ["created_by_id"]),
        ("document_analyses", ["analyzed_by_id"]),
        ("document_analysis_edits", ["edited_by_id"]),
        ("drivers", ["created_by_id"]),
        ("ea_workflow_definitions", ["created_by_id"]),
        ("ea_workflow_instances", ["approved_by_id", "started_by_id", "triggered_by_user_id"]),
        ("ea_workflow_notifications", ["recipient_id"]),
        ("ea_workflow_schedules", ["created_by_id"]),
        ("ea_workflow_step_executions", ["approved_by_id"]),
        ("element_template_recommendations", ["created_by"]),
        ("element_template_usage", ["created_by_id", "instantiated_by"]),
        ("element_templates", ["created_by", "updated_by"]),
        ("enterprise_architecture_frameworks", ["created_by_id", "updated_by_id"]),
        ("enterprise_capabilities", ["created_by"]),
        ("enterprise_initiatives", ["created_by_id"]),
        ("exception_process", ["approver_id"]),
        ("external_systems", ["updated_by_id"]),
        ("gap_remediation_reports", ["approved_by_id", "created_by_id"]),
        ("generation_pipelines", ["created_by_id"]),
        ("git_repositories", ["created_by_id"]),
        ("goals", ["created_by_id"]),
        ("impact_analysis_results", ["created_by_id"]),
        ("implementation_plateaus", ["updated_by"]),
        ("import_audit_log", ["rolled_back_by_id", "user_id"]),
        ("import_sessions", ["user_id"]),
        ("industry_frameworks", ["created_by_id"]),
        ("industry_process_recommendation", ["accepted_by_id"]),
        ("integration_impact_registers", ["approved_by_id", "created_by_id"]),
        ("jira_projects", ["created_by_id"]),
        ("kanban_boards", ["created_by_id"]),
        ("kanban_card_attachments", ["uploaded_by_id"]),
        ("kanban_card_comments", ["user_id"]),
        ("kanban_cards", ["assigned_to_id", "created_by_id"]),
        ("llm_interactions", ["user_id"]),
        ("llm_reviews", ["reviewed_by"]),
        ("logical_data_models", ["created_by_id"]),
        ("migration_plan_documents", ["approved_by_id", "created_by_id"]),
        ("missing_business_collaborations", ["created_by_id"]),
        ("missing_business_interactions", ["created_by_id"]),
        ("missing_business_interfaces", ["created_by_id"]),
        ("options_analysis", ["approved_by_id", "created_by_id"]),
        ("physical_data_models", ["created_by_id"]),
        ("platform_configurations", ["created_by_id"]),
        ("platform_migration_scopes", ["approved_by_id", "created_by_id"]),
        ("policy_exemptions", ["approved_by", "requested_by"]),
        ("policy_violations", ["acknowledged_by", "exemption_approved_by"]),
        ("portfolio_initiatives", ["created_by_id"]),
        ("products", ["created_by_id"]),
        ("project_constraints", ["constraint_owner_id", "created_by_id"]),
        ("published_api_specs", ["published_by_id"]),
        ("quality_attributes", ["created_by_id"]),
        ("quality_frameworks", ["created_by_id"]),
        ("refactoring_tracking", ["created_by_id"]),
        ("regulatory_frameworks", ["created_by_id"]),
        ("relationship_suggestions", ["reviewed_by_id"]),
        ("representations", ["created_by_id"]),
        ("requirements_traceability_matrices", ["approved_by_id", "created_by_id"]),
        ("review_decisions", ["reviewer_id"]),
        ("review_queue_items", ["assigned_to_id", "escalated_to_id", "reviewed_by_id"]),
        ("rfp_templates", ["created_by_id"]),
        ("roadmap_audit", ["user_id"]),
        ("roadmap_gaps", ["created_by", "updated_by"]),
        ("roadmap_scenarios", ["created_by"]),
        ("roadmap_work_packages", ["created_by", "updated_by"]),
        ("runtime_compliance_checks", ["checked_by_id"]),
        ("savings_realizations", ["verified_by"]),
        ("scoring_configurations", ["created_by_id"]),
        ("service_level_agreements", ["created_by_id"]),
        ("sharepoint_config", ["updated_by_id"]),
        ("sla_violations", ["reported_by_id"]),
        ("soc2_audit_log", ["user_id"]),
        ("software_dependencies", ["created_by_id"]),
        ("software_modules", ["created_by_id"]),
        ("solution_adr_direct", ["linked_by_id"]),
        ("solution_adr_links", ["linked_by_id"]),
        ("solution_analysis_sessions", ["created_by_id"]),
        ("solution_app_elements", ["created_by_id"]),
        ("solution_apqc_processes", ["linked_by_id"]),
        ("solution_arb_reviews", ["decided_by_id", "submitted_by_id"]),
        ("solution_archimate_elements", ["created_by_id"]),
        ("solution_assessments_sad", ["created_by_id"]),
        ("solution_benefit_realizations", ["created_by_id", "owner_id"]),
        ("solution_business_elements", ["created_by_id"]),
        ("solution_capability_mappings", ["created_by_id"]),
        ("solution_change_requests", ["created_by_id", "decision_by_id", "submitted_by_id"]),
        ("solution_comments", ["author_id"]),
        ("solution_compliance_mappings", ["assessor_id", "created_by_id"]),
        ("solution_compositions", ["created_by_id"]),
        ("solution_contracts_model", ["created_by_id"]),
        ("solution_cost_comparisons", ["created_by_id"]),
        ("solution_cost_models", ["approved_by_id", "created_by_id"]),
        ("solution_deployment_architectures", ["created_by_id"]),
        ("solution_execution_tracking", ["last_updated_by_id"]),
        ("solution_feasibility_reviews", ["created_by_id", "reviewer_id"]),
        ("solution_governance_exceptions", ["approver_id", "created_by_id"]),
        ("solution_integration_flows", ["created_by_id"]),
        ("solution_investment_phases", ["created_by_id"]),
        ("solution_issues", ["assigned_to_id", "created_by_id", "escalated_to_id", "resolved_by_id"]),
        ("solution_lessons_learned", ["created_by_id", "owner_id"]),
        ("solution_migration_dependencies", ["created_by_id"]),
        ("solution_org_impacts", ["created_by_id"]),
        ("solution_outcome_measurements", ["measured_by_id"]),
        ("solution_outcome_tracking", ["recorded_by_id"]),
        ("solution_outcomes", ["created_by_id"]),
        ("solution_patterns", ["created_by_id"]),
        ("solution_principles_sad", ["created_by_id"]),
        ("solution_quality_attributes", ["created_by_id"]),
        ("solution_risk_snapshots", ["created_by_id"]),
        ("solution_risks", ["created_by_id"]),
        ("solution_scoring_configs", ["created_by_id"]),
        ("solution_session_versions", ["created_by_id"]),
        ("solution_slas", ["created_by_id"]),
        ("solution_spec_generations", ["generated_by_id"]),
        ("solution_stakeholders_sad", ["created_by_id"]),
        ("solution_tech_elements", ["created_by_id"]),
        ("solution_versions", ["approved_by_id", "created_by_id"]),
        ("solutions", ["created_by_id"]),
        ("stakeholder_inputs", ["stakeholder_id"]),
        ("stakeholders", ["created_by_id"]),
        ("strategic_initiatives", ["owner_id"]),
        ("strategic_recommendations", ["created_by_id", "rated_by_id"]),
        ("strategy_resources", ["created_by_id", "custodian_id"]),
        ("system_boundaries", ["created_by_id"]),
        ("system_deployments", ["created_by_id"]),
        ("system_hierarchies", ["created_by_id"]),
        ("system_interfaces", ["created_by_id"]),
        ("system_lifecycles", ["created_by_id"]),
        ("tco_calculations", ["created_by_id"]),
        ("technical_capabilities", ["created_by"]),
        ("technical_capability_vendor_mappings", ["validated_by_id"]),
        ("technical_debt", ["assigned_to_id", "created_by_id"]),
        ("technology_stacks", ["approved_by_id"]),
        ("threshold_configurations", ["created_by_id"]),
        ("traceability_links", ["created_by_id"]),
        ("unified_capability_vendor_organization_mappings", ["contract_manager_id", "executive_sponsor_id"]),
        ("unified_work_packages", ["created_by", "updated_by"]),
        ("usage_analytics", ["user_id"]),
        ("validation_violations", ["resolved_by_id"]),
        ("vendor_capability_risks", ["assessed_by_id"]),
        ("vendor_comparison_criteria", ["created_by_id"]),
        ("vendor_data_sources", ["created_by_id", "verified_by_id"]),
        ("vendor_organizations", ["created_by_id"]),
        ("vendor_process_mappings", ["validated_by_id"]),
        ("vendor_product_capabilities", ["validated_by_id"]),
        ("vendor_product_details", ["created_by_id"]),
        ("vendor_product_families", ["created_by_id"]),
        ("vendor_product_pricing", ["created_by_id"]),
        ("vendor_proof_points", ["verified_by_id"]),
        ("vendor_responses", ["evaluated_by_id"]),
        ("vendor_risk_assessments", ["created_by_id"]),
        ("vendor_selection_reports", ["approved_by_id", "created_by_id"]),
        ("vendor_stack_templates", ["created_by_id", "updated_by_id"]),
        ("viewpoint_views", ["owner_id"]),
        ("work_package_dependencies", ["created_by"]),
        ("workflow_completion_summaries", ["created_by_id"]),
        ("workflow_executions", ["triggered_by_user_id"]),
        ("workflow_nodes", ["approved_by_id"]),
        ("workflow_pipelines", ["approved_by_id", "created_by_id"]),
        ("workflow_run_watchers", ["user_id"]),
        ("workflow_templates", ["created_by_id"]),
    ]

    deleted_count = 0
    errors = []

    for uid in safe_ids:
        try:
            db.session.execute(db.text(f"SAVEPOINT del_user_{uid}"))

            # Step 1: delete rows that are owned by / belong to this user
            for table, col in _DELETE_OWNED:
                db.session.execute(
                    db.text(f"DELETE FROM {table} WHERE {col} = :uid"),
                    {"uid": uid},
                )

            # Step 2: SET NULL on all content/audit tables (one UPDATE per table)
            for table, cols in _SET_NULL_MULTI:
                set_clause = ", ".join(f"{c} = NULL" for c in cols)
                # Build a WHERE that matches any of the columns
                where_clause = " OR ".join(f"{c} = :uid" for c in cols)
                db.session.execute(
                    db.text(f"UPDATE {table} SET {set_clause} WHERE {where_clause}"),
                    {"uid": uid},
                )

            # Step 3: delete the user row itself
            db.session.execute(
                db.text("DELETE FROM users WHERE id = :uid"),
                {"uid": uid},
            )

            db.session.execute(db.text(f"RELEASE SAVEPOINT del_user_{uid}"))
            deleted_count += 1

        except Exception as exc:  # noqa: BLE001
            db.session.execute(db.text(f"ROLLBACK TO SAVEPOINT del_user_{uid}"))
            logger.error("Failed to delete user %s: %s", uid, exc)
            errors.append({"id": uid, "error": str(exc)})

    db.session.commit()

    response = {"deleted": deleted_count}
    if errors:
        response["errors"] = errors
    return jsonify(response)


# ============================================================================
# Jira Push Integration
# ============================================================================


@admin_bp.route("/jira-settings", methods=["GET", "POST"])
@login_required
@admin_required
def jira_settings():
    """Manage Jira push integration configuration."""
    from flask_wtf import FlaskForm
    from wtforms import BooleanField, PasswordField, StringField
    from wtforms.validators import DataRequired

    from app.models.models import ExternalSystem

    class JiraSettingsForm(FlaskForm):
        base_url = StringField("Jira Base URL", validators=[DataRequired()])
        username = StringField("Username / Email", validators=[DataRequired()])
        api_token = PasswordField("API Token")
        project_key = StringField("Project Key", validators=[DataRequired()])
        issue_type = StringField("Issue Type", default="Task")
        filter_countries = StringField("Country Filter", default="United Kingdom")
        enabled = BooleanField("Enable Integration", default=False)

    jira_config = ExternalSystem.query.filter_by(
        system_name="jira", system_type="alm"
    ).first()

    form = JiraSettingsForm()

    if form.validate_on_submit():
        credentials = {}
        if form.api_token.data and form.api_token.data.strip():
            credentials = {
                "username": form.username.data,
                "api_token": form.api_token.data,
            }
        elif jira_config and jira_config.credentials:
            try:
                existing_creds = json.loads(jira_config.credentials)
                credentials = {
                    "username": form.username.data,
                    "api_token": existing_creds.get("api_token", ""),
                }
            except (json.JSONDecodeError, TypeError):
                logger.exception("Failed to JSON parsing")
                pass

        filter_countries = (form.filter_countries.data or "").strip()
        config_data = {
            "project_key": form.project_key.data,
            "issue_type": form.issue_type.data or "Task",
            "filter_countries": [c.strip() for c in filter_countries.split(",") if c.strip()],
            "custom_field_map": {},
        }

        if jira_config:
            jira_config.base_url = form.base_url.data
            jira_config.api_endpoint = f"{form.base_url.data}/rest/api/3"
            jira_config.auth_type = "basic"
            if credentials:
                jira_config.credentials = json.dumps(credentials)
            jira_config.config_json = json.dumps(config_data)
            jira_config.enabled = form.enabled.data
            jira_config.updated_by_id = current_user.id
            jira_config.updated_at = datetime.utcnow()
            flash("Jira settings updated successfully.", "success")
        else:
            if not credentials or not credentials.get("api_token"):
                flash("API token is required for new configuration.", "error")
                _env_url = os.environ.get("JIRA_BASE_URL", "").strip()
                _env_user = os.environ.get("JIRA_USERNAME", "").strip()
                _env_project = os.environ.get("JIRA_PROJECT_KEY", "").strip()
                _env_configured = bool(_env_url and _env_user and _env_project)
                return render_template(
                    "admin/jira_settings.html", form=form, jira_config=None, push_status=None,
                    env_configured=_env_configured,
                )
            jira_config = ExternalSystem(
                system_name="jira",
                system_type="alm",
                base_url=form.base_url.data,
                api_endpoint=f"{form.base_url.data}/rest/api/3",
                auth_type="basic",
                credentials=json.dumps(credentials),
                config_json=json.dumps(config_data),
                enabled=form.enabled.data,
                updated_by_id=current_user.id,
            )
            db.session.add(jira_config)
            flash("Jira settings created successfully.", "success")

        db.session.commit()
        return redirect(url_for("admin.jira_settings"))

    env_jira_url = os.environ.get("JIRA_BASE_URL", "").strip()
    env_jira_user = os.environ.get("JIRA_USERNAME", "").strip()
    env_configured = bool(env_jira_url and env_jira_user and os.environ.get("JIRA_PROJECT_KEY", "").strip())

    push_status = None
    if request.method == "GET":
        if jira_config:
            try:
                creds = json.loads(jira_config.credentials) if jira_config.credentials else {}
            except (json.JSONDecodeError, TypeError):
                creds = {}
            try:
                cfg = json.loads(jira_config.config_json) if jira_config.config_json else {}
            except (json.JSONDecodeError, TypeError):
                cfg = {}

            form.base_url.data = jira_config.base_url
            form.username.data = creds.get("username", "")
            form.project_key.data = cfg.get("project_key", "")
            form.issue_type.data = cfg.get("issue_type", "Task")
            countries = cfg.get("filter_countries", [])
            form.filter_countries.data = ", ".join(countries) if isinstance(countries, list) else str(countries)
            form.enabled.data = jira_config.enabled
        else:
            form.base_url.data = os.environ.get("JIRA_BASE_URL", "")
            form.username.data = os.environ.get("JIRA_USERNAME", "")
            form.project_key.data = os.environ.get("JIRA_PROJECT_KEY", "")

        try:
            from app.services.jira_push_service import get_jira_push_service
            svc = get_jira_push_service()
            push_status = svc.get_push_status()
            if push_status and jira_config and jira_config.config_json:
                cfg_data = json.loads(jira_config.config_json) if isinstance(jira_config.config_json, str) else {}
                last_wh = cfg_data.get("last_webhook_at")
                if last_wh:
                    push_status["last_webhook_at"] = last_wh
        except Exception:
            push_status = None

    return render_template(
        "admin/jira_settings.html",
        form=form,
        jira_config=jira_config,
        push_status=push_status,
        env_configured=env_configured,
    )


@admin_bp.route("/jira-settings/test-connection", methods=["POST"])
@login_required
@admin_required
def jira_test_connection():
    """Test Jira API connectivity."""
    import asyncio

    from app.models.models import ExternalSystem

    jira_config = ExternalSystem.query.filter_by(
        system_name="jira", system_type="alm"
    ).first()

    if not jira_config:
        try:
            from app.services.jira_push_service import get_jira_push_service
            svc = get_jira_push_service()
            if not svc._initialize_connector():
                return jsonify({
                    "success": False,
                    "message": "Jira not configured. Please fill in and save your Jira credentials first, then test the connection."
                }), 400
            loop = asyncio.new_event_loop()
            try:
                connected = loop.run_until_complete(svc.connector.test_connection())
            finally:
                loop.close()
            return jsonify({
                "success": connected,
                "message": "Connection successful (via .env)" if connected else "Connection failed",
            })
        except Exception as e:
            return jsonify({"success": False, "message": str(e)}), 500

    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        if not svc._initialize_connector():
                return jsonify({
                    "success": False,
                    "message": "Failed to connect: Jira credentials are saved but invalid or incomplete. Check the URL, username, API token and project key."
                }), 400

        loop = asyncio.new_event_loop()
        try:
            connected = loop.run_until_complete(svc.connector.test_connection())
        finally:
            loop.close()

        jira_config.last_connection_test = datetime.utcnow()
        jira_config.connection_status = "connected" if connected else "error"
        db.session.commit()

        return jsonify({
            "success": connected,
            "message": "Connection successful" if connected else "Connection failed",
        })
    except Exception as e:
        jira_config.connection_status = "error"
        jira_config.last_error = str(e)
        db.session.commit()
        return jsonify({"success": False, "message": str(e)}), 500


@admin_bp.route("/jira-settings/webhook", methods=["POST"])
@login_required
# csrf.exempt: webhook receiver — Jira sends POST events without CSRF tokens, signature verified via HMAC
@csrf.exempt
def jira_webhook():
    """Receive Jira webhook events and sync status back to KanbanCard."""
    from app.models.models import ExternalSystem

    secret = os.environ.get("JIRA_WEBHOOK_SECRET", "")
    if secret:
        sig = request.headers.get("X-Hub-Signature", "")
        expected = "sha256=" + hmac.new(secret.encode(), request.data, hashlib.sha256).hexdigest()
        if not hmac.compare_digest(expected, sig):
            return jsonify({"error": "Invalid signature"}), 403

    data = request.get_json(silent=True) or {}
    issue_key = data.get("issue", {}).get("key", "")
    if not issue_key:
        return jsonify({"ok": True}), 200

    changelog = data.get("changelog", {})
    items = changelog.get("items", [])
    new_status = None
    for item in items:
        if item.get("field") == "status":
            new_status = item.get("toString")
            break

    if new_status:
        try:
            from app.config.kanban_jira_field_mapping import COLUMN_TO_JIRA_STATUS
            from app.models.adm_kanban_models import KanbanCard
            jira_to_column = {v: k for k, v in COLUMN_TO_JIRA_STATUS.items()}
            new_column = jira_to_column.get(new_status)
            if new_column:
                card = KanbanCard.query.filter_by(jira_issue_key=issue_key).first()
                if card:
                    card.status = new_column
                    card.updated_at = datetime.utcnow()
                    db.session.commit()
        except Exception as _card_exc:
            logger.warning("Jira webhook card update failed for %s: %s", issue_key, _card_exc)

    try:
        jira_config = ExternalSystem.query.filter_by(system_name="jira").first()
        if jira_config:
            cfg = json.loads(jira_config.config_json) if jira_config.config_json else {}
            cfg["last_webhook_at"] = datetime.utcnow().isoformat()
            jira_config.config_json = json.dumps(cfg)
            db.session.commit()
    except Exception as _ts_exc:
        logger.warning("Jira webhook timestamp update failed: %s", _ts_exc)

    return jsonify({"ok": True}), 200


@admin_bp.route("/jira-settings/save-env-config", methods=["POST"])
@login_required
@admin_required
def save_env_jira_config():
    """Save .env Jira credentials to database."""
    from app.models.models import ExternalSystem

    url = os.environ.get("JIRA_BASE_URL", "").strip()
    user = os.environ.get("JIRA_USERNAME", "").strip()
    token = os.environ.get("JIRA_API_TOKEN", "").strip()
    proj = os.environ.get("JIRA_PROJECT_KEY", "").strip()

    missing = [k for k, v in [
        ("JIRA_BASE_URL", url), ("JIRA_USERNAME", user),
        ("JIRA_API_TOKEN", token), ("JIRA_PROJECT_KEY", proj),
    ] if not v]
    if missing:
        flash(f"Missing required Jira environment variables: {', '.join(missing)}", "error")
        return redirect(url_for("admin.jira_settings"))

    jira_config = ExternalSystem.query.filter_by(
        system_name="jira", system_type="alm"
    ).first()
    if jira_config is None:
        jira_config = ExternalSystem(system_name="jira", system_type="alm")
        db.session.add(jira_config)

    jira_config.credentials = json.dumps({"base_url": url, "username": user, "api_token": token, "project_key": proj})
    jira_config.enabled = True
    db.session.commit()

    flash("Jira configuration saved to database.", "success")
    return redirect(url_for("admin.jira_settings"))


@admin_bp.route("/jira-settings/trigger-push", methods=["POST"])
@login_required
@admin_required
def jira_trigger_push():
    """Create a Job and start pushing applications to Jira."""
    from app.models.job import Job, JobStatus

    force = request.json.get("force", False) if request.is_json else False
    app_ids = request.json.get("app_ids") if request.is_json else None

    job = Job(
        name="jira_push",
        task="push_applications",
        payload={"app_ids": app_ids, "force": force},
        status=JobStatus.IN_PROGRESS.value,
        started_at=datetime.utcnow(),
    )
    db.session.add(job)
    db.session.commit()

    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_applications(app_ids=app_ids, force=force)
        job.status = JobStatus.COMPLETED.value
        job.result = result.as_dict()
        job.finished_at = datetime.utcnow()
        db.session.commit()
        return jsonify({"success": True, "job_id": job.id, "result": result.as_dict()})
    except Exception as e:
        job.status = JobStatus.FAILED.value
        job.error = str(e)
        job.finished_at = datetime.utcnow()
        db.session.commit()
        return jsonify({"success": False, "job_id": job.id, "error": str(e)}), 500


@admin_bp.route("/jira-settings/push-status", methods=["GET"])
@login_required
@admin_required
def jira_push_status():
    """Return JSON push status for polling."""
    from app.models.job import Job

    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        status = svc.get_push_status()
    except Exception:
        status = {}

    latest_job = Job.query.filter_by(name="jira_push").order_by(Job.created_at.desc()).first()
    return jsonify({
        "push_status": status,
        "latest_job": latest_job.as_dict() if latest_job else None,
    })


@admin_bp.route("/jira-settings/kanban-push-status", methods=["GET"])
@login_required
@admin_required
def jira_kanban_push_status():
    """Return JSON kanban push status for polling."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        kanban_status = svc.get_kanban_push_status()
        return jsonify({"success": True, "kanban_status": kanban_status})
    except Exception as exc:
        logger.error("Kanban push status error: %s", exc)
        return jsonify({"success": False, "error": str(exc)}), 500


@admin_bp.route("/jira-settings/trigger-kanban-push", methods=["POST"])
@login_required
@admin_required
def jira_trigger_kanban_push():
    """Push all unpushed KanbanCard rows to Jira."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_all_unpushed_cards()
        return jsonify({"success": True, **result})
    except Exception as exc:
        logger.error("Trigger kanban push error: %s", exc)
        return jsonify({"success": False, "error": str(exc)}), 500


@admin_bp.route("/jira-settings/push-epics", methods=["POST"])
@login_required
@admin_required
def jira_push_epics():
    """Create one Jira Epic per ADM phase."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_adm_phase_epics()
        return jsonify({"success": True, **result})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp.route("/jira-settings/push-applications", methods=["POST"])
@login_required
@admin_required
def jira_push_applications():
    """Push ApplicationComponents to Jira."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_applications()
        return jsonify({"success": True, **result.as_dict()})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp.route("/jira-settings/push-dependencies", methods=["POST"])
@login_required
@admin_required
def jira_push_dependencies():
    """Create Jira Subtasks from KanbanCard dependencies."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_card_dependencies()
        return jsonify({"success": True, **result})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp.route("/jira-settings/field-discovery", methods=["GET"])
@login_required
@admin_required
def jira_field_discovery():
    """Return available Jira fields for the configured project."""
    import asyncio

    from app.models.models import ExternalSystem

    jira_config = ExternalSystem.query.filter_by(
        system_name="jira", system_type="alm"
    ).first()

    if not jira_config:
        return jsonify({"success": False, "message": "Jira not configured"}), 400

    try:
        cfg = json.loads(jira_config.config_json) if jira_config.config_json else {}
    except (json.JSONDecodeError, TypeError):
        cfg = {}

    project_key = cfg.get("project_key", "EA")
    issue_type_id = request.args.get("issue_type_id", "10001")

    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        if not svc._initialize_connector():
            return jsonify({"success": False, "message": "Connector init failed"}), 400

        loop = asyncio.new_event_loop()
        try:
            fields = loop.run_until_complete(
                svc.connector.discover_fields(project_key, issue_type_id)
            )
        finally:
            loop.close()

        return jsonify({"success": True, "fields": fields})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


# ============================================================================
# PLT-024 / PLT-034: Portfolio Export (PDF print page + PPTX download)
# ============================================================================


def _get_portfolio_export_data() -> dict:
    """Gather portfolio data for PDF and PPTX exports.

    Returns a dict with solutions_by_status, top10, bottom10, junction_coverage,
    total_solutions, and avg_completeness.
    """
    from app.models.solution_models import Solution

    solutions = Solution.query.limit(500).all()
    total = len(solutions)
    if total == 0:
        return {
            "total_solutions": 0,
            "avg_completeness": 0,
            "solutions_by_status": {},
            "top10": [],
            "bottom10": [],
            "junction_coverage": {},
            "generated_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        }

    scored = []
    for sol in solutions:
        try:
            cs = sol.architecture_completeness_score
            score = cs["score"]
            filled = cs["filled"]
        except Exception:
            score = 0
            filled = []
        scored.append({
            "name": sol.name or f"Solution #{sol.id}",
            "governance_status": sol.governance_status or "draft",
            "score": score,
            "filled": filled,
        })

    avg_completeness = round(sum(s["score"] for s in scored) / total) if total else 0

    # Solutions by governance status
    status_counts = {}
    for s in scored:
        st = s["governance_status"]
        status_counts[st] = status_counts.get(st, 0) + 1

    sorted_desc = sorted(scored, key=lambda s: s["score"], reverse=True)
    sorted_asc = sorted(scored, key=lambda s: s["score"])

    junction_names = [
        "drivers", "goals", "constraints", "requirements",
        "applications", "archimate_elements", "vendor_products",
        "recommendations", "risks", "tco_items", "metrics",
        "plateaus", "apqc_processes", "capability_mappings",
    ]
    junction_coverage = {}
    for jname in junction_names:
        with_count = sum(1 for s in scored if jname in s["filled"])
        junction_coverage[jname] = {
            "with": with_count,
            "without": total - with_count,
        }

    return {
        "total_solutions": total,
        "avg_completeness": avg_completeness,
        "solutions_by_status": status_counts,
        "top10": [{"name": s["name"], "score": s["score"]} for s in sorted_desc[:10]],
        "bottom10": [{"name": s["name"], "score": s["score"]} for s in sorted_asc[:10]],
        "junction_coverage": junction_coverage,
        "generated_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
    }


@admin_bp.route("/export-portfolio-pdf")
@login_required
@admin_required
def export_portfolio_pdf():
    """PLT-024: Render a print-optimised HTML page for board-ready portfolio PDF.

    Uses window.print() with @media print CSS -- no new dependencies.
    """
    try:
        data = _get_portfolio_export_data()
    except Exception as exc:
        logger.warning("PLT-024: Portfolio PDF data gathering failed: %s", exc)
        data = {
            "total_solutions": 0,
            "avg_completeness": 0,
            "solutions_by_status": {},
            "top10": [],
            "bottom10": [],
            "junction_coverage": {},
            "generated_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        }

    # Render inline HTML (no new template per constraint)
    html_parts = []
    html_parts.append("""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>Architecture Portfolio Summary</title>
<style>
body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
       margin: 2rem; color: #1a1a2e; }
h1 { font-size: 1.75rem; margin-bottom: 0.25rem; }
h2 { font-size: 1.2rem; margin-top: 2rem; border-bottom: 2px solid #3b82f6;
     padding-bottom: 0.25rem; color: #1e3a5f; }
.subtitle { color: #6b7280; font-size: 0.9rem; margin-bottom: 1.5rem; }
.kpi-row { display: flex; gap: 2rem; margin-bottom: 1.5rem; }
.kpi { text-align: center; }
.kpi .value { font-size: 2rem; font-weight: 700; }
.kpi .label { font-size: 0.8rem; color: #6b7280; }
table { width: 100%; border-collapse: collapse; margin-top: 0.5rem; font-size: 0.85rem; }
th, td { border: 1px solid #d1d5db; padding: 0.4rem 0.75rem; text-align: left; }
th { background: #f1f5f9; font-weight: 600; }
tr:nth-child(even) { background: #f9fafb; }
.confidential { color: #9ca3af; font-size: 0.7rem; text-align: center;
                margin-top: 2rem; border-top: 1px solid #e5e7eb; padding-top: 0.5rem; }
.print-btn { background: #3b82f6; color: #fff; border: none; padding: 0.5rem 1.5rem;
             border-radius: 0.375rem; font-size: 0.9rem; cursor: pointer; margin-bottom: 1rem; }
.print-btn:hover { background: #2563eb; }
.back-link { margin-right: 1rem; color: #3b82f6; text-decoration: none; font-size: 0.9rem; }
@media print {
  .no-print { display: none !important; }
  body { margin: 0; font-size: 11pt; }
  h1 { font-size: 16pt; }
  h2 { font-size: 13pt; }
}
</style>
</head>
<body>
<div class="no-print" style="margin-bottom:1rem;">
  <a href="/admin" class="back-link">&larr; Back to Admin</a>
  <button class="print-btn" onclick="window.print()">Print / Save as PDF</button>
</div>
""")

    html_parts.append(f"""<h1>Architecture Portfolio Summary</h1>
<p class="subtitle">Generated {data['generated_at']} &mdash; CONFIDENTIAL</p>

<div class="kpi-row">
  <div class="kpi"><div class="value">{data['total_solutions']}</div><div class="label">Total Solutions</div></div>
  <div class="kpi"><div class="value">{data['avg_completeness']}%</div><div class="label">Avg Completeness</div></div>
  <div class="kpi"><div class="value">{len(data['solutions_by_status'])}</div><div class="label">Status Categories</div></div>
</div>
""")

    # Solutions by status table
    html_parts.append("<h2>Solutions by Governance Status</h2><table><tr><th>Status</th><th>Count</th></tr>")
    if data["solutions_by_status"]:
        for status, count in sorted(data["solutions_by_status"].items()):
            label = status.replace("_", " ").title()
            html_parts.append(f"<tr><td>{label}</td><td>{count}</td></tr>")
    else:
        html_parts.append("<tr><td colspan='2'>No data available</td></tr>")
    html_parts.append("</table>")

    # Top 10
    html_parts.append("<h2>Top 10 Highest Completeness</h2><table><tr><th>#</th><th>Solution</th><th>Score</th></tr>")
    if data["top10"]:
        for i, s in enumerate(data["top10"], 1):
            html_parts.append(f"<tr><td>{i}</td><td>{s['name']}</td><td>{s['score']}%</td></tr>")
    else:
        html_parts.append("<tr><td colspan='3'>No data available</td></tr>")
    html_parts.append("</table>")

    # Bottom 10
    html_parts.append("<h2>Top 10 Gaps (Lowest Completeness)</h2><table><tr><th>#</th><th>Solution</th><th>Score</th></tr>")
    if data["bottom10"]:
        for i, s in enumerate(data["bottom10"], 1):
            html_parts.append(f"<tr><td>{i}</td><td>{s['name']}</td><td>{s['score']}%</td></tr>")
    else:
        html_parts.append("<tr><td colspan='3'>No data available</td></tr>")
    html_parts.append("</table>")

    # Junction coverage
    html_parts.append("<h2>Junction Coverage</h2><table><tr><th>Junction Type</th><th>With Data</th><th>Without Data</th></tr>")
    if data["junction_coverage"]:
        for jname, jdata in data["junction_coverage"].items():
            label = jname.replace("_", " ").title()
            html_parts.append(f"<tr><td>{label}</td><td>{jdata['with']}</td><td>{jdata['without']}</td></tr>")
    else:
        html_parts.append("<tr><td colspan='3'>No data available</td></tr>")
    html_parts.append("</table>")

    html_parts.append('<p class="confidential">CONFIDENTIAL &mdash; For authorised recipients only</p>')
    html_parts.append("</body></html>")

    return "".join(html_parts)


@admin_bp.route("/export-portfolio-pptx")
@login_required
@admin_required
def export_portfolio_pptx():
    """PLT-034: Generate a PowerPoint deck for portfolio summary.

    5 slides: Title, Key Metrics, Top 10, Bottom 10, Junction Coverage.
    """
    import io as _io

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.error("PLT-034: python-pptx not available")
        return jsonify({"error": "PowerPoint export unavailable — python-pptx not installed"}), 500

    try:
        data = _get_portfolio_export_data()
    except Exception as exc:
        logger.warning("PLT-034: Portfolio PPTX data gathering failed: %s", exc)
        data = {
            "total_solutions": 0,
            "avg_completeness": 0,
            "solutions_by_status": {},
            "top10": [],
            "bottom10": [],
            "junction_coverage": {},
            "generated_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        }

    BRAND_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
    ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
    GRAY = RGBColor(0x6B, 0x72, 0x80)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _add_title_footer(slide, footer_text="CONFIDENTIAL"):
        """Add a small footer to a slide."""
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(8)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    def _set_cell(table, row, col, text, bold=False, font_size=11, color=None):
        """Set a table cell text and formatting."""
        cell = table.cell(row, col)
        cell.text = str(text)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold
            if color:
                paragraph.font.color.rgb = color

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Architecture Portfolio Summary"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = data["generated_at"]
    p2.font.size = Pt(16)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    _add_title_footer(slide1)

    # Slide 2: Key Metrics + Status table
    slide2 = prs.slides.add_slide(slide_layout)
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Metrics"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    kpis = [
        ("Total Solutions", str(data["total_solutions"])),
        ("Avg Completeness", f"{data['avg_completeness']}%"),
    ]
    for idx, (label, value) in enumerate(kpis):
        txBox = slide2.shapes.add_textbox(Inches(1 + idx * 4), Inches(1.3), Inches(3.5), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    status_items = sorted(data["solutions_by_status"].items()) if data["solutions_by_status"] else [("No data", 0)]
    rows = len(status_items) + 1
    tbl = slide2.shapes.add_table(rows, 2, Inches(1), Inches(3.2), Inches(6), Inches(0.4 * rows)).table
    tbl.columns[0].width = Inches(4)
    tbl.columns[1].width = Inches(2)
    _set_cell(tbl, 0, 0, "Governance Status", bold=True, font_size=12)
    _set_cell(tbl, 0, 1, "Count", bold=True, font_size=12)
    for row_idx, (status, count) in enumerate(status_items, 1):
        label = status.replace("_", " ").title() if isinstance(status, str) else str(status)
        _set_cell(tbl, row_idx, 0, label, font_size=11)
        _set_cell(tbl, row_idx, 1, str(count), font_size=11)
    _add_title_footer(slide2)

    # Slide 3: Top 10 Highest Completeness
    slide3 = prs.slides.add_slide(slide_layout)
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Top 10 — Highest Completeness"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    top_items = data["top10"] if data["top10"] else [{"name": "No data available", "score": "-"}]
    rows = len(top_items) + 1
    tbl = slide3.shapes.add_table(rows, 3, Inches(1), Inches(1.3), Inches(10), Inches(0.4 * rows)).table
    tbl.columns[0].width = Inches(1)
    tbl.columns[1].width = Inches(7)
    tbl.columns[2].width = Inches(2)
    _set_cell(tbl, 0, 0, "#", bold=True, font_size=12)
    _set_cell(tbl, 0, 1, "Solution", bold=True, font_size=12)
    _set_cell(tbl, 0, 2, "Score", bold=True, font_size=12)
    for row_idx, s in enumerate(top_items, 1):
        _set_cell(tbl, row_idx, 0, str(row_idx), font_size=11)
        _set_cell(tbl, row_idx, 1, s["name"], font_size=11)
        score_text = f"{s['score']}%" if isinstance(s["score"], (int, float)) else str(s["score"])
        _set_cell(tbl, row_idx, 2, score_text, font_size=11)
    _add_title_footer(slide3)

    # Slide 4: Top 10 Gaps (Lowest Completeness)
    slide4 = prs.slides.add_slide(slide_layout)
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Top 10 Gaps — Lowest Completeness"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    bottom_items = data["bottom10"] if data["bottom10"] else [{"name": "No data available", "score": "-"}]
    rows = len(bottom_items) + 1
    tbl = slide4.shapes.add_table(rows, 3, Inches(1), Inches(1.3), Inches(10), Inches(0.4 * rows)).table
    tbl.columns[0].width = Inches(1)
    tbl.columns[1].width = Inches(7)
    tbl.columns[2].width = Inches(2)
    _set_cell(tbl, 0, 0, "#", bold=True, font_size=12)
    _set_cell(tbl, 0, 1, "Solution", bold=True, font_size=12)
    _set_cell(tbl, 0, 2, "Score", bold=True, font_size=12)
    for row_idx, s in enumerate(bottom_items, 1):
        _set_cell(tbl, row_idx, 0, str(row_idx), font_size=11)
        _set_cell(tbl, row_idx, 1, s["name"], font_size=11)
        score_text = f"{s['score']}%" if isinstance(s["score"], (int, float)) else str(s["score"])
        _set_cell(tbl, row_idx, 2, score_text, font_size=11)
    _add_title_footer(slide4)

    # Slide 5: Junction Coverage
    slide5 = prs.slides.add_slide(slide_layout)
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Junction Coverage (14 Types)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    jc = data["junction_coverage"]
    jc_items = list(jc.items()) if jc else [("No data", {"with": 0, "without": 0})]
    rows = len(jc_items) + 1
    tbl = slide5.shapes.add_table(rows, 3, Inches(1), Inches(1.3), Inches(10), Inches(0.35 * rows)).table
    tbl.columns[0].width = Inches(5)
    tbl.columns[1].width = Inches(2.5)
    tbl.columns[2].width = Inches(2.5)
    _set_cell(tbl, 0, 0, "Junction Type", bold=True, font_size=12)
    _set_cell(tbl, 0, 1, "Solutions With", bold=True, font_size=12)
    _set_cell(tbl, 0, 2, "Solutions Without", bold=True, font_size=12)
    for row_idx, (jname, jdata) in enumerate(jc_items, 1):
        label = jname.replace("_", " ").title()
        _set_cell(tbl, row_idx, 0, label, font_size=11)
        _set_cell(tbl, row_idx, 1, str(jdata["with"]), font_size=11)
        _set_cell(tbl, row_idx, 2, str(jdata["without"]), font_size=11)
    _add_title_footer(slide5)

    # Serialize and return
    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    from flask import send_file
    filename = f"portfolio_summary_{datetime.utcnow().strftime('%Y%m%d')}.pptx"
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename,
    )


# ============================================================================
# Vendor Pricing Import (FAR-008)
# ============================================================================


@admin_bp.route("/vendor-pricing/import", methods=["GET", "POST"])
@login_required
@admin_required
def vendor_pricing_import():
    """Admin page for importing vendor pricing from CSV/Excel or contracts."""
    from werkzeug.utils import secure_filename

    staging_items = []
    upload_error = None

    if request.method == "POST":
        file = request.files.get("file")
        if not file or file.filename == "":
            flash("No file selected", "error")
            return render_template(
                "admin/vendor_pricing_import.html",
                staging_items=[],
                upload_error="No file selected",
            )

        filename = secure_filename(file.filename)
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if ext in ("csv", "xlsx", "xls"):
            try:
                from app.services.import_core.file_parser import ImportFileParser

                parser = ImportFileParser()
                rows, errors = parser.parse_file(file, filename)
                if errors:
                    upload_error = "; ".join(errors[:5])
                else:
                    required = {"vendor_name", "product_name", "tier", "annual_cost"}
                    validation = parser.validate_file_structure(
                        rows, list(rows[0].keys()) if rows else [], required
                    )
                    if not validation.get("valid"):
                        upload_error = "; ".join(validation.get("errors", [])[:5])
                    else:
                        staging_items = rows
            except Exception as exc:
                upload_error = f"Parse error: {exc}"

        elif ext in ("pdf", "docx", "txt", "md"):
            try:
                from app.modules.vendors.services.pricing_extraction_service import (
                    PricingExtractionService,
                )

                svc = PricingExtractionService()
                result = svc.extract_from_file(file)
                if result["success"]:
                    staging_items = result["items"]
                else:
                    upload_error = result.get("error", "Extraction failed")
            except Exception as exc:
                upload_error = f"Extraction error: {exc}"
        else:
            upload_error = f"Unsupported file type: .{ext}"

    return render_template(
        "admin/vendor_pricing_import.html",
        staging_items=staging_items,
        upload_error=upload_error,
    )


@admin_bp.route("/vendor-pricing/confirm", methods=["POST"])
@login_required
@admin_required
def vendor_pricing_confirm():
    """Confirm staged pricing items — write to VendorProductPricing as contract_verified."""
    from difflib import SequenceMatcher
    from datetime import datetime as _dt, timezone
    from app.models.vendor.vendor_organization import (
        VendorOrganization,
        VendorProduct,
        VendorProductPricing,
    )

    data = request.get_json()
    if not data or "items" not in data:
        return jsonify({"error": "No items to confirm"}), 400

    created = 0
    updated = 0
    errors = []

    for item in data["items"]:
        vendor_name = item.get("vendor_name", "").strip()
        product_name = item.get("product_name", "").strip()
        if not vendor_name or not product_name:
            errors.append("Missing vendor/product name in item")
            continue

        vendor = VendorOrganization.query.filter(
            db.func.lower(VendorOrganization.name) == vendor_name.lower()
        ).first()
        if not vendor:
            all_vendors = VendorOrganization.query.all()
            best_match = None
            best_score = 0.0
            for v in all_vendors:
                score = SequenceMatcher(None, vendor_name.lower(), v.name.lower()).ratio()
                if score > best_score:
                    best_score = score
                    best_match = v
            if best_match and best_score >= 0.7:
                vendor = best_match
            else:
                errors.append(f"Vendor not found: {vendor_name}")
                continue

        product = VendorProduct.query.filter(
            VendorProduct.vendor_organization_id == vendor.id,
            db.func.lower(VendorProduct.name) == product_name.lower(),
        ).first()
        if not product:
            products = VendorProduct.query.filter_by(vendor_organization_id=vendor.id).all()
            best_match = None
            best_score = 0.0
            for p in products:
                score = SequenceMatcher(None, product_name.lower(), p.name.lower()).ratio()
                if score > best_score:
                    best_score = score
                    best_match = p
            if best_match and best_score >= 0.6:
                product = best_match
            else:
                errors.append(f"Product not found: {vendor_name} / {product_name}")
                continue

        tier = item.get("tier", "Standard")
        annual_cost = item.get("annual_cost")
        existing = VendorProductPricing.query.filter_by(
            product_id=product.id, tier_name=tier
        ).first()

        if existing:
            existing.list_price_annual = annual_cost
            existing.data_source_type = "contract_verified"
            existing.typical_discount_percent = item.get("discount_percent")
            existing.contract_term_months = item.get("contract_term_months", 12)
            existing.setup_fee = item.get("setup_fee")
            existing.implementation_fee = item.get("implementation_fee")
            existing.source = "Contract Import"
            existing.last_verified_at = _dt.now(timezone.utc)
            existing.organization_id = 1
            if item.get("expiry_date"):
                try:
                    from datetime import date as _date
                    existing.expiry_date = _date.fromisoformat(item["expiry_date"])
                except (ValueError, TypeError):
                    logger.exception("Failed to operation")
                    pass
            updated += 1
        else:
            pricing = VendorProductPricing(
                product_id=product.id,
                pricing_model=item.get("unit_type", "per_user"),
                tier_name=tier,
                list_price_annual=annual_cost,
                min_users=1,
                max_users=item.get("unit_quantity"),
                currency=item.get("currency", "USD"),
                source="Contract Import",
                billing_frequency="annual",
                contract_term_months=item.get("contract_term_months", 12),
                setup_fee=item.get("setup_fee"),
                implementation_fee=item.get("implementation_fee"),
                typical_discount_percent=item.get("discount_percent"),
                data_source_type="contract_verified",
                last_verified_at=_dt.now(timezone.utc),
                organization_id=1,
                created_by_id=current_user.id,
            )
            if item.get("expiry_date"):
                try:
                    from datetime import date as _date
                    pricing.expiry_date = _date.fromisoformat(item["expiry_date"])
                except (ValueError, TypeError):
                    logger.exception("Failed to operation")
                    pass
            db.session.add(pricing)
            created += 1

    db.session.commit()
    return jsonify({"success": True, "created": created, "updated": updated, "errors": errors}), 200


# ============================================================================
# Pricing Analytics (FAR-008)
# ============================================================================


@admin_bp.route("/pricing-analytics")
@login_required
@admin_required
def pricing_analytics():
    """Vendor pricing analytics dashboard."""
    try:
        from app.modules.vendors.services.confidence_engine import ConfidenceEngine

        engine = ConfidenceEngine()
        analytics = engine.get_analytics()
    except Exception as exc:
        logger.warning("pricing_analytics: ConfidenceEngine unavailable: %s", exc)
        analytics = {
            "coverage_pct": 0,
            "covered_capabilities": 0,
            "total_capabilities": 0,
            "confidence_distribution": {},
            "stale_count": 0,
            "conflict_count": 0,
        }
    return render_template("admin/pricing_analytics.html", analytics=analytics)


# ---------------------------------------------------------------------------
# Power Platform CoE Integration
# ---------------------------------------------------------------------------

_PP_PROVIDER = "power_platform_coe"
_PP_LABEL = "default"


def _pp_settings_row():
    """Return the APISettings row for Power Platform CoE, or None."""
    return APISettings.query.filter_by(
        provider=_PP_PROVIDER, key_label=_PP_LABEL
    ).first()


@admin_bp.route("/integrations/power-platform", methods=["GET"])
@login_required
def power_platform_integration():
    """GET /admin/integrations/power-platform — CoE configuration and discovery UI."""
    row = _pp_settings_row()
    config = {}
    if row:
        config = {
            "tenant_id": row.jira_url or "",
            "client_id": row.jira_email or "",
            "env_url": row.custom_endpoint_url or "",
            "configured": bool(row.jira_url),
        }
    return render_template("admin/integrations_power_platform.html", config=config)


@admin_bp.route("/integrations/power-platform/save", methods=["POST"])
@login_required
def power_platform_save_credentials():
    """POST /admin/integrations/power-platform/save — persist credentials to api_settings."""
    data = request.get_json() or request.form
    row = _pp_settings_row()
    if not row:
        row = APISettings(provider=_PP_PROVIDER, key_label=_PP_LABEL)
        db.session.add(row)
    row.jira_url = data.get("tenant_id", "")
    row.jira_email = data.get("client_id", "")
    # Only update secret if a non-empty value was provided (avoids blanking stored secret)
    raw_secret = data.get("client_secret", "")
    if raw_secret:
        row.api_key = raw_secret
    row.custom_endpoint_url = data.get("env_url", "")
    db.session.commit()
    return jsonify({"success": True})


@admin_bp.route("/integrations/power-platform/test", methods=["POST"])
@login_required
def power_platform_test_connection():
    """POST /admin/integrations/power-platform/test — test credentials."""
    from app.modules.solutions_strategic.v2.services.power_platform_coe_service import (
        PowerPlatformCoeService,
    )

    data = request.get_json() or {}
    tenant_id = data.get("tenant_id", "")
    client_id = data.get("client_id", "")
    client_secret = data.get("client_secret", "")
    env_url = data.get("env_url", "")

    # Fall back to stored credentials for any missing field
    if not (tenant_id and client_id and client_secret):
        row = _pp_settings_row()
        if row:
            tenant_id = tenant_id or row.jira_url or ""
            client_id = client_id or row.jira_email or ""
            client_secret = client_secret or row.api_key or ""
            env_url = env_url or row.custom_endpoint_url or ""

    if not (tenant_id and client_id and client_secret):
        return jsonify({"status": "error", "message": "Credentials not configured"}), 400

    result = PowerPlatformCoeService.test_connection(
        tenant_id, client_id, client_secret, env_url
    )
    return jsonify(result)


@admin_bp.route("/integrations/power-platform/discover", methods=["POST"])
@login_required
def power_platform_discover():
    """POST /admin/integrations/power-platform/discover — trigger discovery, return app list."""
    from app.models.application_portfolio import ApplicationComponent
    from app.modules.solutions_strategic.v2.services.power_platform_coe_service import (
        PowerPlatformCoeService,
    )

    row = _pp_settings_row()
    if not row or not row.jira_url:
        return jsonify({"error": "Power Platform credentials not configured"}), 400

    apps = PowerPlatformCoeService.discover_apps(
        row.jira_url or "", row.jira_email or "", row.api_key or ""
    )

    # Annotate with ARCHIE link status
    linked_ids = {
        r.source_identifier
        for r in ApplicationComponent.query.filter(
            ApplicationComponent.source_identifier.isnot(None)
        )
        .with_entities(ApplicationComponent.source_identifier)
        .all()
    }
    for app in apps:
        app["status"] = "linked" if app["id"] in linked_ids else "unregistered"

    ungoverned = sum(1 for a in apps if not a.get("owner_email"))
    linked = sum(1 for a in apps if a["status"] == "linked")

    return jsonify({
        "apps": apps,
        "total": len(apps),
        "ungoverned": ungoverned,
        "linked": linked,
    })


@admin_bp.route("/integrations/power-platform/import", methods=["POST"])
@login_required
def power_platform_import():
    """POST /admin/integrations/power-platform/import — import selected app_ids."""
    from app.modules.solutions_strategic.v2.services.power_platform_coe_service import (
        PowerPlatformCoeService,
    )

    data = request.get_json() or {}
    app_ids = data.get("app_ids", [])
    if not app_ids:
        return jsonify({"error": "No app_ids provided"}), 400

    row = _pp_settings_row()
    if not row or not row.jira_url:
        return jsonify({"error": "Credentials not configured"}), 400

    discovered = PowerPlatformCoeService.discover_apps(
        row.jira_url or "", row.jira_email or "", row.api_key or ""
    )
    result = PowerPlatformCoeService.import_apps(app_ids, discovered, current_user.id)
    status_code = 200 if (result["imported"] > 0 or result["already_exists"] > 0) else 400
    return jsonify(result), status_code


# ============================================================================
# ServiceNow CMDB Integration — ENT-074
# ============================================================================

@admin_bp.route("/integrations/servicenow", methods=["GET", "POST"])
@login_required
@admin_required
def servicenow_integration():
    """Manage ServiceNow CMDB integration configuration."""
    from flask_wtf import FlaskForm
    from wtforms import BooleanField, IntegerField, PasswordField, StringField
    from wtforms.validators import DataRequired, NumberRange, Optional as OptionalValidator, URL
    
    from app.models.models import ExternalSystem
    
    class ServiceNowSettingsForm(FlaskForm):
        instance_url = StringField(
            "ServiceNow Instance URL",
            validators=[DataRequired(), URL()],
            description="e.g., https://yourorg.service-now.com"
        )
        username = StringField(
            "API Username",
            validators=[DataRequired()],
            description="ServiceNow user with CMDB read access"
        )
        password = PasswordField(
            "API Password / Token",
            description="Leave blank to keep existing password"
        )
        batch_size = IntegerField(
            "Batch Size",
            validators=[OptionalValidator(), NumberRange(min=1, max=1000)],
            default=100,
            description="Records per API call (1-1000)"
        )
        query_filter = StringField(
            "Query Filter",
            description="Optional ServiceNow query (e.g., operational_status=1)"
        )
        enabled = BooleanField("Enable Integration", default=False)
    
    # Load existing ServiceNow config
    snow_config = ExternalSystem.query.filter_by(
        system_name="servicenow",
        system_type="cmdb"
    ).first()
    
    form = ServiceNowSettingsForm()
    
    if form.validate_on_submit():
        # Build credentials dict
        credentials = {}
        if form.password.data and form.password.data.strip():
            credentials = {
                "username": form.username.data,
                "password": form.password.data,
            }
        elif snow_config and snow_config.credentials:
            # Keep existing password if not changed
            try:
                existing_creds = json.loads(snow_config.credentials)
                credentials = {
                    "username": form.username.data,
                    "password": existing_creds.get("password", ""),
                }
            except (json.JSONDecodeError, TypeError):
                logger.exception("Failed to parse existing ServiceNow credentials")
                pass
        
        # Build config JSON
        config_data = {
            "batch_size": form.batch_size.data or 100,
            "query_filter": form.query_filter.data or "",
            "cmdb_tables": ["cmdb_ci_appl", "cmdb_ci_server", "cmdb_ci_database"],
        }
        
        if snow_config:
            # Update existing config
            snow_config.base_url = form.instance_url.data.rstrip('/')
            snow_config.api_endpoint = f"{form.instance_url.data.rstrip('/')}/api/now/table"
            snow_config.auth_type = "basic"
            if credentials:
                snow_config.credentials = json.dumps(credentials)
            snow_config.config_json = json.dumps(config_data)
            snow_config.enabled = form.enabled.data
            snow_config.updated_by_id = current_user.id
            snow_config.updated_at = datetime.utcnow()
            flash("ServiceNow settings updated successfully.", "success")
        else:
            # Create new config
            if not credentials or not credentials.get("password"):
                flash("Password is required for new configuration.", "error")
                return render_template(
                    "admin/servicenow_integration.html",
                    form=form,
                    snow_config=None,
                    last_sync=None
                )
            
            snow_config = ExternalSystem(
                system_name="servicenow",
                system_type="cmdb",
                base_url=form.instance_url.data.rstrip('/'),
                api_endpoint=f"{form.instance_url.data.rstrip('/')}/api/now/table",
                auth_type="basic",
                credentials=json.dumps(credentials),
                config_json=json.dumps(config_data),
                enabled=form.enabled.data,
                updated_by_id=current_user.id,
            )
            db.session.add(snow_config)
            flash("ServiceNow settings created successfully.", "success")
        
        db.session.commit()
        return redirect(url_for("admin.servicenow_integration"))
    
    # Pre-populate form with existing config
    if snow_config and request.method == "GET":
        form.instance_url.data = snow_config.base_url
        form.enabled.data = snow_config.enabled
        
        # Load credentials
        if snow_config.credentials:
            try:
                creds = json.loads(snow_config.credentials)
                form.username.data = creds.get("username", "")
            except (json.JSONDecodeError, TypeError):
                pass
        
        # Load config
        if snow_config.config_json:
            try:
                config = json.loads(snow_config.config_json)
                form.batch_size.data = config.get("batch_size", 100)
                form.query_filter.data = config.get("query_filter", "")
            except (json.JSONDecodeError, TypeError):
                pass
    
    # Get last sync info
    last_sync = None
    if snow_config:
        last_sync = {
            "last_sync_at": snow_config.updated_at,
            "status": "Enabled" if snow_config.enabled else "Disabled"
        }
    
    return render_template(
        "admin/servicenow_integration.html",
        form=form,
        snow_config=snow_config,
        last_sync=last_sync
    )


@admin_bp.route("/integrations/servicenow/test-connection", methods=["POST"])
@login_required
@admin_required
def servicenow_test_connection():
    """Test ServiceNow CMDB connection."""
    from app.modules.vendors.connectors.servicenow_connector import ServiceNowConnector
    
    snow_config = ExternalSystem.query.filter_by(
        system_name="servicenow",
        system_type="cmdb"
    ).first()
    
    if not snow_config:
        return jsonify({
            "success": False,
            "message": "ServiceNow not configured"
        }), 400
    
    # Load credentials
    try:
        creds = json.loads(snow_config.credentials or '{}')
        username = creds.get('username')
        password = creds.get('password')
    except (json.JSONDecodeError, TypeError):
        return jsonify({
            "success": False,
            "message": "Invalid credentials configuration"
        }), 400
    
    # Test connection
    connector = ServiceNowConnector(
        instance_url=snow_config.base_url,
        username=username,
        password=password
    )
    
    if connector.health_check():
        return jsonify({
            "success": True,
            "message": "Connection successful! ServiceNow CMDB is accessible."
        })
    else:
        return jsonify({
            "success": False,
            "message": "Connection failed. Check instance URL and credentials."
        }), 400


@admin_bp.route("/integrations/servicenow/trigger-sync", methods=["POST"])
@login_required
@admin_required
@audit_log
def servicenow_trigger_sync():
    """Trigger immediate ServiceNow CMDB sync."""
    from app.modules.vendors.connectors.servicenow_connector import ServiceNowConnector
    
    snow_config = ExternalSystem.query.filter_by(
        system_name="servicenow",
        system_type="cmdb"
    ).first()
    
    if not snow_config or not snow_config.enabled:
        return jsonify({
            "success": False,
            "message": "ServiceNow integration is not enabled"
        }), 400
    
    # Load credentials
    try:
        creds = json.loads(snow_config.credentials or '{}')
        username = creds.get('username')
        password = creds.get('password')
    except (json.JSONDecodeError, TypeError):
        return jsonify({
            "success": False,
            "message": "Invalid credentials configuration"
        }), 400
    
    # Load config
    config_data = {}
    if snow_config.config_json:
        try:
            config_data = json.loads(snow_config.config_json)
        except (json.JSONDecodeError, TypeError):
            pass
    
    # Create connector and sync
    connector = ServiceNowConnector(
        instance_url=snow_config.base_url,
        username=username,
        password=password,
        batch_size=config_data.get('batch_size', 100)
    )
    
    result = connector.sync_applications(
        query_filter=config_data.get('query_filter')
    )
    
    # Update last sync timestamp
    snow_config.updated_at = datetime.utcnow()
    db.session.commit()
    
    if result.healthy:
        return jsonify({
            "success": True,
            "message": f"Sync completed: {result.applications_created} created, "
                      f"{result.applications_updated} updated",
            "result": result.to_dict()
        })
    else:
        return jsonify({
            "success": False,
            "message": "Sync completed with errors",
            "result": result.to_dict()
        }), 500


@admin_bp.route("/integrations/servicenow/sync-status", methods=["GET"])
@login_required
@admin_required
def servicenow_sync_status():
    """Get ServiceNow sync status and statistics."""
    from app.models.application_portfolio import ApplicationComponent
    
    snow_config = ExternalSystem.query.filter_by(
        system_name="servicenow",
        system_type="cmdb"
    ).first()
    
    if not snow_config:
        return jsonify({
            "configured": False,
            "enabled": False
        })
    
    # Count applications synced from ServiceNow (have external_id)
    synced_apps = ApplicationComponent.query.filter(
        ApplicationComponent.external_id.isnot(None)
    ).count()
    
    return jsonify({
        "configured": True,
        "enabled": snow_config.enabled,
        "last_sync": snow_config.updated_at.isoformat() if snow_config.updated_at else None,
        "synced_applications": synced_apps,
        "instance_url": snow_config.base_url
    })
