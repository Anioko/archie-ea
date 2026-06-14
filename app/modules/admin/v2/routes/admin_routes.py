"""
Admin Routes v2 — guardrail-enabled.

Uses the new architecture:
- @timed_route for automatic metrics collection on all endpoints
- Observability (request_id in response headers)
- Consistent error handling via exception mappers

URL prefix preserved: /admin (applied via register() in v2/__init__.py)
Blueprint name: admin (same as v1 for url_for compatibility — cross-module
refs exist in app/application_mgmt/impact_routes.py calling
url_for("admin.api_settings"))

All routes preserved exactly from v1 admin_routes.py.
"""

import logging
import os

from werkzeug.utils import secure_filename

from flask import (
    Blueprint,
    flash,
    jsonify,
    redirect,
    render_template,
    request,
    url_for,
)
from flask_login import current_user, login_required
from sqlalchemy.orm import aliased

try:
    from flask_rq import get_queue

    HAS_RQ = True
except ImportError:
    HAS_RQ = False
    get_queue = None

import hmac
import hashlib
import json
from datetime import datetime, timedelta

from app import csrf
from app.extensions import db
from app.core.compat import mark_blueprint_guardrailed
from app.core.decorators import timed_route
from ...forms.admin_forms import (
    APISettingsForm,
    ChangeAccountTypeForm,
    ChangeUserEmailForm,
    FeatureFlagForm,
    InviteUserForm,
    NewUserForm,
)
from app.modules.account.forms.account_forms import CreatePasswordForm
from app.decorators import admin_required, audit_log
from app.middleware.tenant_decorators import org_admin_required, platform_admin_required
from app.models import APISettings, EditableHTML, Role, User
from app.models.organization import Organization
from app.models.ai_service import AIPromptTemplate
from app.models.feature_flags import FeatureFlag, FeatureState, FeatureType
from app.modules.admin.v2.services.llm_service_v2 import test_api_key
from app.utils.sidebar_parser import SidebarSubmenu, parse_sidebar_template
from app.modules.admin.v2.services.admin_user_service_v2 import AdminUserService

# Blueprint name MUST be "admin" (not "admin_v2") because cross-module code
# uses url_for("admin.api_settings") etc.  The 3-tier fallback in
# _bootstrap/blueprints.py guarantees only one tier is active at a time,
# so there is no name collision.
admin_bp_v2 = Blueprint("admin", __name__)
mark_blueprint_guardrailed(admin_bp_v2)

logger = logging.getLogger(__name__)

_svc = AdminUserService


@admin_bp_v2.route("/send-digest", methods=["POST"])
@login_required
@admin_required
@timed_route
def send_digest():
    """PLT-009/PLT-031: Manual trigger for weekly digest emails."""
    from flask import current_app
    from app._bootstrap._digest_emails import (
        send_data_maturity_digest,
        send_executive_summary,
    )

    digest_type = request.form.get("type", "both")
    results = {}

    if digest_type in ("maturity", "both"):
        data = send_data_maturity_digest(current_app._get_current_object())
        results["maturity"] = {
            "total_solutions": data["total"],
            "avg_score": data["avg_score"],
            "zero_connections": len(data["zero_connections"]),
        }

    if digest_type in ("executive", "both"):
        data = send_executive_summary(current_app._get_current_object())
        results["executive"] = {
            "total_solutions": data["total_solutions"],
            "new_this_week": data["new_solutions_count"],
            "arb_decisions": data["arb_decisions_count"],
        }

    return jsonify({"success": True, "results": results})


_TEST_ARTIFACT_PATTERNS = [
    "J1-AutoTest-%", "New Solution%", "J1 Bootstrap%",
    "J1 Test%", "J1 Regression%", "J1-Debug%", "J1-Test-%",
    "%E2E Test%", "%Journey Test%", "AI Test%", "Minimal Test%",
    "QA Test%", "Driver Test%", "Blueprint Test%", "%JDD Test%",
    "%Gap2 Persistence%", "Post-Deploy%", "%Smoke Test%",
    "% Test Solution%", "% Test Solution", "%Audit Test%",
    "%Forensic Audit%", "%Verification Test%", "%Uniformity Verification%",
    "% Test Programme%", "%PESTLE News Analyser%", "MDM Test%",
    "Create an architecture for%",
]


def _get_solutions_portfolio_stats() -> dict:
    """ENT-013: Compute solutions portfolio stats for admin dashboard."""
    try:
        from app.models.solution_models import Solution
        from app.models.solution_lifecycle_models import SolutionRisk
        from sqlalchemy import not_, or_
        _real_filter = [
            Solution.name.isnot(None),
            not_(or_(*[Solution.name.like(p) for p in _TEST_ARTIFACT_PATTERNS])),
        ]
        total = Solution.query.filter(*_real_filter).count()
        phases = {}
        for sol in Solution.query.filter(*_real_filter).with_entities(Solution.adm_phase).all():
            p = sol.adm_phase or "A"
            phases[p] = phases.get(p, 0) + 1
        pending_arb = Solution.query.filter(*_real_filter).filter(
            Solution.governance_status.in_(["pending_review", "in_review"])
        ).count()
        top_risks = (
            SolutionRisk.query
            .filter(SolutionRisk.impact.in_(["critical", "high"]))
            .filter_by(status="open")
            .order_by(SolutionRisk.created_at.desc())
            .limit(3)
            .all()
        )
        return {
            "total": total,
            "phases": phases,
            "pending_arb": pending_arb,
            "top_risks": [{"description": r.risk_description[:80], "impact": r.impact} for r in top_risks],
        }
    except Exception as e:
        logger.warning(f"Could not compute solution portfolio stats: {e}")
        return {"total": 0, "phases": {}, "pending_arb": 0, "top_risks": []}


def _get_data_maturity_stats() -> dict:
    """PLT-006: Compute portfolio-wide data maturity metrics for admin dashboard."""
    try:
        from app.models.solution_models import Solution
        from sqlalchemy import not_, or_

        solutions = Solution.query.filter(
            Solution.name.isnot(None),
            not_(or_(*[Solution.name.like(p) for p in _TEST_ARTIFACT_PATTERNS])),
        ).limit(50).all()
        total = len(solutions)
        if total == 0:
            return {
                "total": 0,
                "avg_score": 0,
                "buckets": {"minimal": 0, "partial": 0, "good": 0},
                "top5": [],
                "bottom5": [],
                "junction_fill_rates": {},
            }

        scores = []
        for sol in solutions:
            try:
                cs = sol.architecture_completeness_score
                scores.append({
                    "id": sol.id,
                    "name": sol.name or f"Solution #{sol.id}",
                    "score": cs["score"],
                    "filled": cs["filled"],
                    "filled_count": cs["filled_count"],
                    "total_junctions": cs["total"],
                })
            except Exception:
                scores.append({
                    "id": sol.id,
                    "name": sol.name or f"Solution #{sol.id}",
                    "score": 0,
                    "filled": [],
                    "filled_count": 0,
                    "total_junctions": 14,
                })

        avg_score = round(sum(s["score"] for s in scores) / total) if total else 0

        minimal = sum(1 for s in scores if s["score"] < 25)
        partial = sum(1 for s in scores if 25 <= s["score"] <= 75)
        good = sum(1 for s in scores if s["score"] > 75)

        sorted_asc = sorted(scores, key=lambda s: s["score"])
        sorted_desc = sorted(scores, key=lambda s: s["score"], reverse=True)
        top5 = sorted_desc[:5]
        bottom5 = sorted_asc[:5]

        # Junction fill rates: how many solutions have each junction populated
        junction_names = [
            "drivers", "goals", "constraints", "requirements",
            "applications", "archimate_elements", "vendor_products",
            "recommendations", "risks", "tco_items", "metrics",
            "plateaus", "apqc_processes", "capability_mappings",
        ]
        junction_fill = {}
        for jname in junction_names:
            count = sum(1 for s in scores if jname in s["filled"])
            junction_fill[jname] = {
                "count": count,
                "total": total,
                "pct": round(count / total * 100) if total else 0,
            }

        return {
            "total": total,
            "avg_score": avg_score,
            "buckets": {"minimal": minimal, "partial": partial, "good": good},
            "top5": top5,
            "bottom5": bottom5,
            "junction_fill_rates": junction_fill,
        }
    except Exception as e:
        logger.warning(f"Could not compute data maturity stats: {e}")
        return {
            "total": 0,
            "avg_score": 0,
            "buckets": {"minimal": 0, "partial": 0, "good": 0},
            "top5": [],
            "bottom5": [],
            "junction_fill_rates": {},
        }


# ============================================================================
# Dashboard & Index
# ============================================================================


@admin_bp_v2.route("/")
@timed_route
@login_required
@admin_required
def index():
    """Admin command center — system health, config status, user activity."""
    ctx = {}

    # System health
    health = {"db": "ok", "redis": "ok"}
    try:
        db.session.execute(db.text("SELECT 1"))
    except Exception:
        health["db"] = "error"
    try:
        import redis as _redis, os as _os
        _r = _redis.from_url(_os.environ.get("REDIS_URL", "redis://localhost:6379/0"))
        _r.ping()
    except Exception:
        health["redis"] = "error"
    ctx["health"] = health

    # API key status
    try:
        api_keys = APISettings.query.order_by(APISettings.provider).all()
    except Exception:
        api_keys = []
    ctx["api_keys"] = api_keys

    # Feature flags summary
    try:
        flags_enabled = FeatureFlag.query.filter_by(enabled=True).count()
        flags_disabled = FeatureFlag.query.filter_by(enabled=False).count()
        flags_total = flags_enabled + flags_disabled
    except Exception:
        flags_enabled = flags_disabled = flags_total = 0
    ctx["flags"] = {"enabled": flags_enabled, "disabled": flags_disabled, "total": flags_total}

    # Recent users
    try:
        recent_users = User.query.order_by(User.id.desc()).limit(8).all()
        total_users = User.query.count()
        unconfirmed = User.query.filter_by(confirmed=False).count()
    except Exception:
        recent_users = []
        total_users = unconfirmed = 0
    ctx["recent_users"] = recent_users
    ctx["total_users"] = total_users
    ctx["unconfirmed_users"] = unconfirmed

    return render_template("admin/index.html", **ctx)


@admin_bp_v2.route("/dashboard")
@timed_route
@login_required
@admin_required
def dashboard():
    """Admin dashboard with stats and overview."""
    page = request.args.get("page", 1, type=int)
    per_page = request.args.get("per_page", 10, type=int)
    search_query = request.args.get("search", "")

    pagination = _svc.get_paginated_users(page, per_page, search_query)
    users = pagination.items
    roles = _svc.get_all_roles()

    # PLT-006: Data maturity stats
    data_maturity = _get_data_maturity_stats()

    # Overview stats — admin/index.html (shared template) renders flag/user/key
    # cards, so supply them here too (otherwise the template hits undefined vars).
    try:
        api_keys = APISettings.query.order_by(APISettings.provider).all()
    except Exception:
        api_keys = []
    try:
        flags_enabled = FeatureFlag.query.filter_by(enabled=True).count()
        flags_disabled = FeatureFlag.query.filter_by(enabled=False).count()
    except Exception:
        flags_enabled = flags_disabled = 0
    try:
        recent_users = User.query.order_by(User.id.desc()).limit(8).all()
        total_users = User.query.count()
        unconfirmed_users = User.query.filter_by(confirmed=False).count()
    except Exception:
        recent_users = []
        total_users = unconfirmed_users = 0

    return render_template(
        "admin/index.html",
        users=users,
        roles=roles,
        pagination=pagination,
        search_query=search_query,
        per_page=per_page,
        data_maturity=data_maturity,
        api_keys=api_keys,
        flags={
            "enabled": flags_enabled,
            "disabled": flags_disabled,
            "total": flags_enabled + flags_disabled,
        },
        recent_users=recent_users,
        total_users=total_users,
        unconfirmed_users=unconfirmed_users,
    )


# ============================================================================
# User Management (thin routes -> AdminUserService)
# ============================================================================


@admin_bp_v2.route("/new-user", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("create_user")
def new_user():
    """Create a new user."""
    form = NewUserForm()
    if form.validate_on_submit():
        user = _svc.create_user(
            first_name=form.first_name.data,
            last_name=form.last_name.data,
            email=form.email.data,
            password=form.password.data,
            role=form.role.data,
        )
        flash("User {} successfully created".format(user.full_name()), "form-success")
    return render_template("admin/new_user.html", form=form)


@admin_bp_v2.route("/invite-user", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("invite_user")
def invite_user():
    """Invites a new user to create an account and set their own password."""
    form = InviteUserForm()
    if form.validate_on_submit():
        user = _svc.invite_user(
            first_name=form.first_name.data,
            last_name=form.last_name.data,
            email=form.email.data,
            role=form.role.data,
        )
        flash("User {} successfully invited".format(user.full_name()), "form-success")
    return render_template("admin/new_user.html", form=form)


@admin_bp_v2.route("/manage-users")
def manage_users_redirect():
    """Redirect legacy /admin/manage-users to canonical /admin/users."""
    return redirect(url_for("admin.registered_users"), code=301)


@admin_bp_v2.route("/users")
@timed_route
@login_required
@admin_required
def registered_users():
    """View all registered users."""
    users = _svc.get_all_users()
    roles = _svc.get_all_roles()
    return render_template("admin/registered_users.html", users=users, roles=roles)


@admin_bp_v2.route("/user/<int:user_id>")
@admin_bp_v2.route("/user/<int:user_id>/info")
@timed_route
@login_required
@admin_required
def user_info(user_id):
    """View a user's profile."""
    user = _svc.get_user_or_404(user_id)
    return render_template("admin/manage_user.html", user=user)


@admin_bp_v2.route("/user/<int:user_id>/change-email", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("change_user_email")
def change_user_email(user_id):
    """Change a user's email."""
    user = _svc.get_user_or_404(user_id)
    form = ChangeUserEmailForm()
    if form.validate_on_submit():
        _svc.change_user_email(user, form.email.data)
        flash(
            "Email for user {} successfully changed to {}.".format(
                user.full_name(), user.email
            ),
            "form-success",
        )
    return render_template("admin/manage_user.html", user=user, form=form)


@admin_bp_v2.route("/user/<int:user_id>/change-account-type", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("change_account_type")
def change_account_type(user_id):
    """Change a user's account type."""
    if current_user.id == user_id:
        flash(
            "You cannot change the type of your own account. Please ask "
            "another administrator to do this.",
            "error",
        )
        return redirect(url_for("admin.user_info", user_id=user_id))

    user = _svc.get_user_or_404(user_id)
    form = ChangeAccountTypeForm()
    if form.validate_on_submit():
        _svc.change_user_role(user, form.role.data)
        role_name = user.role.name if user.role else "No Role"
        flash(
            "Role for user {} successfully changed to {}.".format(
                user.full_name(), role_name
            ),
            "form-success",
        )
    return render_template("admin/manage_user.html", user=user, form=form)


@admin_bp_v2.route("/user/<int:user_id>/set-password", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("set_user_password")
def set_user_password(user_id):
    """Set or reset a user's password."""
    user = _svc.get_user_or_404(user_id)
    form = CreatePasswordForm()
    if form.validate_on_submit():
        _svc.set_user_password(user, form.password.data)
        flash(
            "Password for user {} successfully updated. The account is now login-ready.".format(
                user.full_name()
            ),
            "form-success",
        )
        return redirect(url_for("admin.user_info", user_id=user_id))
    return render_template("admin/manage_user.html", user=user, form=form)


@admin_bp_v2.route("/user/<int:user_id>/delete")
@timed_route
@login_required
@admin_required
def delete_user_request(user_id):
    """Request deletion of a user's account."""
    user = _svc.get_user_or_404(user_id)
    return render_template("admin/manage_user.html", user=user)


@admin_bp_v2.route("/user/<int:user_id>/_delete")
@timed_route
@login_required
@admin_required
def delete_user(user_id):
    """Delete a user's account."""
    if current_user.id == user_id:
        flash(
            "You cannot delete your own account. Please ask another "
            "administrator to do this.",
            "error",
        )
    else:
        user = _svc.get_user_or_404(user_id)
        success, message = _svc.delete_user(user)
        flash(message, "success")
    return redirect(url_for("admin.registered_users"))


# ============================================================================
# Editor Contents
# ============================================================================


@admin_bp_v2.route("/_update_editor_contents", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("update_editor_contents")
def update_editor_contents():
    """Update the contents of an editor."""
    edit_data = request.form.get("edit_data")
    editor_name = request.form.get("editor_name")

    editor_contents = EditableHTML.query.filter_by(editor_name=editor_name).first()
    if editor_contents is None:
        editor_contents = EditableHTML(editor_name=editor_name)
    editor_contents.value = edit_data

    db.session.add(editor_contents)
    db.session.commit()
    return "OK", 200


# ============================================================================
# API Settings
# ============================================================================


@admin_bp_v2.route("/api-settings", methods=["GET", "POST"])
@timed_route
@org_admin_required
@admin_required
@audit_log("update_api_settings")
def api_settings():
    """Manage API settings for LLM providers."""
    all_settings = APISettings.query.order_by(APISettings.provider, APISettings.key_label).all()
    edit_id = request.args.get("edit", type=int)
    edit_settings = APISettings.query.get(edit_id) if edit_id else None
    form = APISettingsForm()

    if form.validate_on_submit():
        provider = form.provider.data
        key_label = (form.key_label.data or "").strip()
        # For edit, look up by ID; for create, look up by (provider, key_label) pair
        if edit_settings:
            existing = edit_settings
        else:
            existing = APISettings.query.filter_by(provider=provider, key_label=key_label).first()

        if existing:
            if form.api_key.data and form.api_key.data.strip():
                existing.api_key = form.api_key.data
            existing.key_label = key_label
            existing.enabled = form.enabled.data
            existing.default_model = form.default_model.data or None
            existing.max_tokens = form.max_tokens.data or 4000
            existing.temperature = form.temperature.data or 0.7
            existing.jira_url = form.jira_url.data or None
            existing.jira_email = form.jira_email.data or None
            existing.hf_model_id = form.hf_model_id.data or None
            existing.hf_endpoint_url = form.hf_endpoint_url.data or None
            existing.custom_endpoint_url = form.custom_endpoint_url.data or None
            existing.custom_auth_method = form.custom_auth_method.data or "bearer"
            existing.custom_headers = form.custom_headers.data or None
            existing.updated_at = datetime.utcnow()
            existing.updated_by_id = current_user.id
            settings = existing
        else:
            if not form.api_key.data or not form.api_key.data.strip():
                flash("API key is required when creating new settings.", "error")
                return render_template(
                    "admin/api_settings.html",
                    form=form,
                    all_settings=all_settings,
                    edit_id=edit_id,
                    edit_settings=edit_settings,
                )

            settings = APISettings(
                provider=provider,
                key_label=key_label,
                api_key=form.api_key.data,
                enabled=form.enabled.data,
                default_model=form.default_model.data or None,
                max_tokens=form.max_tokens.data or 4000,
                temperature=form.temperature.data or 0.7,
                jira_url=form.jira_url.data or None,
                jira_email=form.jira_email.data or None,
                hf_model_id=form.hf_model_id.data or None,
                hf_endpoint_url=form.hf_endpoint_url.data or None,
                custom_endpoint_url=form.custom_endpoint_url.data or None,
                custom_auth_method=form.custom_auth_method.data or "bearer",
                custom_headers=form.custom_headers.data or None,
                updated_by_id=current_user.id,
            )
            db.session.add(settings)

        api_key_to_test = (
            form.api_key.data
            if form.api_key.data and form.api_key.data.strip()
            else settings.api_key
        )
        if form.test.data:
            if not api_key_to_test:
                flash("No API key available to test.", "error")
                return redirect(url_for("admin.api_settings"))
            test_result = test_api_key(provider, api_key_to_test)
            settings.test_status = "success" if test_result["success"] else "failed"
            settings.test_message = test_result.get("message", "")
            settings.last_tested_at = datetime.utcnow()

            if test_result["success"]:
                flash(
                    f"Connection test successful: {test_result['message']}", "success"
                )
            else:
                flash(
                    f"Connection test failed: {test_result.get('message', 'Unknown error')}",
                    "error",
                )
        else:
            label_suffix = f" ({key_label})" if key_label else ""
            flash(f"API settings for {provider}{label_suffix} saved successfully.", "success")

        db.session.commit()
        return redirect(url_for("admin.api_settings"))

    if edit_settings:
        form.provider.data = edit_settings.provider
        form.key_label.data = edit_settings.key_label or ""
        form.settings_id.data = str(edit_settings.id)
        form.api_key.data = edit_settings.api_key
        form.enabled.data = edit_settings.enabled
        form.default_model.data = edit_settings.default_model or ""
        form.max_tokens.data = edit_settings.max_tokens or 4000
        form.temperature.data = edit_settings.temperature or 0.7
        form.jira_url.data = edit_settings.jira_url or ""
        form.jira_email.data = edit_settings.jira_email or ""
        form.hf_model_id.data = edit_settings.hf_model_id or ""
        form.hf_endpoint_url.data = edit_settings.hf_endpoint_url or ""
        form.custom_endpoint_url.data = edit_settings.custom_endpoint_url or ""
        form.custom_auth_method.data = edit_settings.custom_auth_method or "bearer"
        form.custom_headers.data = edit_settings.custom_headers or ""

    return render_template(
        "admin/api_settings.html",
        form=form,
        all_settings=all_settings,
        edit_id=edit_id,
        edit_settings=edit_settings,
    )


@admin_bp_v2.route("/api-settings/<int:settings_id>/delete", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("delete_api_settings")
def delete_api_settings(settings_id):
    """Delete API settings."""
    settings = APISettings.query.get_or_404(settings_id)
    provider = settings.provider
    db.session.delete(settings)
    db.session.commit()
    flash(f"API settings for {provider} deleted successfully.", "success")
    return redirect(url_for("admin.api_settings"))


@admin_bp_v2.route("/api-settings/<int:settings_id>/test", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("test_api_settings")
def test_api_settings(settings_id):
    """Test API settings."""
    settings = APISettings.query.get_or_404(settings_id)

    try:
        result = test_api_key(settings.api_key, settings.api_provider)
        flash(f"API test successful: {result}", "success")
    except Exception as e:
        flash("API test failed. Please try again.", "error")

    return redirect(url_for("admin.api_settings"))


@admin_bp_v2.route("/api-settings/env-keys", methods=["GET"])
@timed_route
@login_required
@admin_required
def preview_env_keys():
    """Preview API keys found in environment variables but not yet in the database."""
    import os

    env_key_map = {
        "OPENAI_API_KEY": "openai",
        "ANTHROPIC_API_KEY": "anthropic",
        "GOOGLE_API_KEY": "gemini",
        "GEMINI_API_KEY": "gemini",
        "DEEPSEEK_API_KEY": "deepseek",
        "HUGGINGFACE_API_KEY": "huggingface",
        "HF_TOKEN": "huggingface",
        "AZURE_OPENAI_API_KEY": "azure",
        "OPENROUTER_API_KEY": "openrouter",
    }

    default_models = {
        "openai": "gpt-4o",
        "anthropic": "claude-3-5-sonnet-20241022",
        "gemini": "gemini-2.0-flash-exp",
        "deepseek": "deepseek-chat",
        "huggingface": "meta-llama/Llama-3.1-8B-Instruct",
        "azure": "gpt-4",
        "openrouter": "google/gemini-2.5-flash-preview:free",
    }

    existing = {s.provider for s in APISettings.query.all()}

    found_keys = []
    for env_var, provider in env_key_map.items():
        value = os.environ.get(env_var, "")
        if value and value.strip():
            masked = value[:8] + "..." + value[-4:] if len(value) > 12 else "****"
            found_keys.append(
                {
                    "env_var": env_var,
                    "provider": provider,
                    "masked_key": masked,
                    "in_database": provider in existing,
                    "default_model": default_models.get(provider, ""),
                }
            )

    return jsonify(
        {
            "success": True,
            "keys": found_keys,
            "count": len(found_keys),
            "new_count": sum(1 for k in found_keys if not k["in_database"]),
        }
    )


@admin_bp_v2.route("/api-settings/update-model", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("update_provider_model")
def update_provider_model():
    """Quick-update the default model for a provider without the full form."""
    data = request.get_json()
    if not data or "provider" not in data or "model" not in data:
        return jsonify(
            {"success": False, "error": "provider and model are required"}
        ), 400

    provider = data["provider"]
    model = data["model"].strip()

    settings = APISettings.query.filter_by(provider=provider).first()
    if not settings:
        return jsonify(
            {"success": False, "error": f"No settings found for {provider}"}
        ), 404

    settings.default_model = model
    settings.updated_at = datetime.utcnow()
    settings.updated_by_id = current_user.id
    db.session.commit()

    return jsonify({"success": True, "provider": provider, "model": model})


@admin_bp_v2.route("/api-settings/load-env", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("load_env_keys")
def load_env_keys():
    """Import selected API keys from environment variables into the database."""
    import os

    data = request.get_json()
    if not data or "keys" not in data:
        return jsonify({"success": False, "error": "No keys specified"}), 400

    env_key_map = {
        "OPENAI_API_KEY": "openai",
        "ANTHROPIC_API_KEY": "anthropic",
        "GOOGLE_API_KEY": "gemini",
        "GEMINI_API_KEY": "gemini",
        "DEEPSEEK_API_KEY": "deepseek",
        "HUGGINGFACE_API_KEY": "huggingface",
        "HF_TOKEN": "huggingface",
        "AZURE_OPENAI_API_KEY": "azure",
        "OPENROUTER_API_KEY": "openrouter",
    }

    default_models = {
        "openai": "gpt-4o",
        "anthropic": "claude-3-5-sonnet-20241022",
        "gemini": "gemini-2.0-flash-exp",
        "deepseek": "deepseek-chat",
        "huggingface": "meta-llama/Llama-3.1-8B-Instruct",
        "azure": "gpt-4",
        "openrouter": "google/gemini-2.5-flash-preview:free",
    }

    requested_keys = data["keys"]
    update_existing = data.get("update_existing", False)

    imported = []
    skipped = []
    errors = []

    existing_map = {s.provider: s for s in APISettings.query.all()}

    for env_var in requested_keys:
        provider = env_key_map.get(env_var)
        if not provider:
            errors.append(f"Unknown env var: {env_var}")
            continue

        api_key = os.environ.get(env_var, "").strip()
        if not api_key:
            errors.append(f"{env_var} is empty")
            continue

        if provider in existing_map and not update_existing:
            skipped.append(f"{provider} already in database (skipped)")
            continue

        if provider in existing_map:
            existing_map[provider].api_key = api_key
            existing_map[provider].enabled = True
            existing_map[provider].updated_at = datetime.utcnow()
            existing_map[provider].updated_by_id = current_user.id
            imported.append(f"{provider} (updated)")
        else:
            settings = APISettings(
                provider=provider,
                api_key=api_key,
                enabled=True,
                default_model=default_models.get(provider, ""),
                max_tokens=4000,
                temperature=0.7,
                updated_by_id=current_user.id,
            )
            db.session.add(settings)
            imported.append(f"{provider} (created)")

    db.session.commit()

    return jsonify(
        {
            "success": True,
            "imported": imported,
            "skipped": skipped,
            "errors": errors,
            "message": f"Imported {len(imported)} key(s)",
        }
    )


# ============================================================================
# Consolidation Status
# ============================================================================


@admin_bp_v2.route("/consolidation")
@timed_route
@login_required
@admin_required
def consolidation_status():
    """Blueprint Consolidation Status Dashboard"""
    from app.api_registry import (
        API_REGISTRY,
        CONSOLIDATION_TARGETS,
        BlueprintStatus,
        get_consolidation_status,
    )

    status = get_consolidation_status()

    blueprints = []
    for name, bp_info in API_REGISTRY.items():
        bp_data = {
            "name": name,
            "status": bp_info.status.value,
            "url_prefix": bp_info.url_prefix,
            "endpoints_count": bp_info.endpoints_count,
            "consolidate_into": bp_info.consolidate_into,
            "notes": bp_info.notes or "",
            "is_active": bp_info.status == BlueprintStatus.ACTIVE,
            "is_deprecated": bp_info.status == BlueprintStatus.DEPRECATED,
            "is_consolidated": bp_info.status == BlueprintStatus.CONSOLIDATED,
        }
        blueprints.append(bp_data)

    status_order = {"deprecated": 0, "active": 1, "consolidated": 2, "legacy": 3}
    blueprints.sort(key=lambda x: (status_order.get(x["status"], 4), x["name"]))

    targets = []
    for target_name, target_info in CONSOLIDATION_TARGETS.items():
        keep_bp = target_info["keep"]
        consolidate_list = target_info["consolidate"]
        priority = target_info["priority"]

        completed_count = 0
        total_count = len(consolidate_list)

        for bp_name in consolidate_list:
            if bp_name in API_REGISTRY:
                bp_info = API_REGISTRY[bp_name]
                if bp_info.status in [
                    BlueprintStatus.CONSOLIDATED,
                    BlueprintStatus.DEPRECATED,
                ]:
                    completed_count += 1

        targets.append(
            {
                "name": target_name,
                "keep_blueprint": keep_bp,
                "consolidate_blueprints": consolidate_list,
                "priority": priority,
                "completed_count": completed_count,
                "total_count": total_count,
                "progress_percentage": (completed_count / total_count * 100)
                if total_count > 0
                else 100,
                "is_complete": completed_count == total_count,
                "estimated_effort": target_info.get("estimated_effort", "unknown"),
            }
        )

    priority_order = {"HIGH": 0, "MEDIUM": 1, "LOW": 2}
    targets.sort(
        key=lambda x: (priority_order.get(x["priority"], 3), -x["progress_percentage"])
    )

    return render_template(
        "admin/consolidation_status.html",
        status=status,
        blueprints=blueprints,
        targets=targets,
        BlueprintStatus=BlueprintStatus,
    )


# ============================================================================
# Feature Flags Management
# ============================================================================


@admin_bp_v2.route("/feature-flags")
@timed_route
@platform_admin_required
@admin_required
def feature_flags():
    """Feature flags management page with pagination."""
    page = request.args.get("page", 1, type=int)
    per_page = request.args.get("per_page", 50, type=int)
    search = request.args.get("search", "")
    filter_type = request.args.get("type", "")
    filter_state = request.args.get("state", "")
    filter_status = request.args.get("status", "")
    filter_hierarchy = request.args.get("hierarchy", "")

    query = FeatureFlag.query

    if search:
        search_term = f"%{search}%"
        query = query.filter(
            db.or_(
                FeatureFlag.name.ilike(search_term),
                FeatureFlag.key.ilike(search_term),
                FeatureFlag.description.ilike(search_term),
            )
        )

    if filter_type:
        query = query.filter(FeatureFlag.feature_type == filter_type)

    if filter_state:
        query = query.filter(FeatureFlag.state == filter_state)

    if filter_status:
        if filter_status == "enabled":
            query = query.filter(FeatureFlag.enabled == True)
        elif filter_status == "disabled":
            query = query.filter(FeatureFlag.enabled == False)

    if filter_hierarchy:
        if filter_hierarchy == "parent":
            query = query.filter(FeatureFlag.parent_id == None)
        elif filter_hierarchy == "child":
            parent_table = aliased(FeatureFlag)
            query = query.join(
                parent_table, FeatureFlag.parent_id == parent_table.id
            ).filter(FeatureFlag.parent_id != None, parent_table.parent_id == None)
        elif filter_hierarchy == "grandchild":
            parent_table = aliased(FeatureFlag)
            query = query.join(
                parent_table, FeatureFlag.parent_id == parent_table.id
            ).filter(FeatureFlag.parent_id != None, parent_table.parent_id != None)

    query = query.order_by(FeatureFlag.sort_order, FeatureFlag.name)
    pagination = query.paginate(page=page, per_page=per_page, error_out=False)

    if pagination.total == 0 and not search:
        flash(
            "No features found. Run 'flask discover-features' to auto-discover routes, "
            "or create features manually.",
            "info",
        )

    # Build sections for grouped view using parent_id hierarchy
    # Top-level = sidebar_section with no parent
    top_sections = (
        FeatureFlag.query.filter(
            FeatureFlag.feature_type == FeatureType.SIDEBAR_SECTION,
            FeatureFlag.parent_id.is_(None),
        )
        .order_by(FeatureFlag.sort_order, FeatureFlag.name)
        .all()
    )
    sections = []
    for sec in top_sections:
        children = (
            FeatureFlag.query.filter(FeatureFlag.parent_id == sec.id)
            .order_by(FeatureFlag.sort_order, FeatureFlag.name)
            .all()
        )
        sections.append({"section": sec, "children": children})

    # Default to list view when filters are active, otherwise grouped
    view_mode = "grouped"
    if search or filter_type or filter_state or filter_status or filter_hierarchy:
        view_mode = "list"

    return render_template(
        "admin/feature_flags.html",
        features=pagination.items,
        pagination=pagination,
        sections=sections,
        view_mode=view_mode,
        search=search,
        filter_type=filter_type,
        filter_state=filter_state,
        filter_status=filter_status,
        filter_hierarchy=filter_hierarchy,
        FeatureType=FeatureType,
        FeatureState=FeatureState,
    )


@admin_bp_v2.route("/feature-flags/new", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("create_feature_flag")
def feature_flag_new():
    """Create new feature flag."""
    form = FeatureFlagForm()

    if form.validate_on_submit():
        try:
            routes_data = None
            if form.routes.data and form.routes.data.strip():
                try:
                    routes_data = json.loads(form.routes.data)
                except json.JSONDecodeError:
                    flash("Invalid JSON format for routes", "error")
                    return render_template(
                        "admin/feature_flag_form.html", form=form, action="New"
                    )

            feature = FeatureFlag(
                key=form.key.data,
                name=form.name.data,
                description=form.description.data,
                feature_type=FeatureType(form.feature_type.data),
                state=FeatureState(form.state.data),
                enabled=form.enabled.data,
                sidebar_label=form.sidebar_label.data,
                sidebar_icon=form.sidebar_icon.data,
                routes=routes_data,
                parent_id=form.parent_id.data or None,
                sort_order=form.sort_order.data or 0,
                last_modified_by=current_user.id,
            )

            db.session.add(feature)
            db.session.commit()

            flash(f"Feature flag '{feature.name}' created successfully", "success")
            return redirect(url_for("admin.feature_flags"))

        except Exception as e:
            db.session.rollback()
            flash("Error creating feature flag. Please try again.", "error")

    return render_template("admin/feature_flag_form.html", form=form, action="New")


@admin_bp_v2.route("/feature-flags/<int:id>/edit", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("edit_feature_flag")
def feature_flag_edit(id):
    """Edit feature flag."""
    feature = FeatureFlag.query.get_or_404(id)
    form = FeatureFlagForm(obj=feature)

    if form.validate_on_submit():
        try:
            routes_data = None
            if form.routes.data and form.routes.data.strip():
                try:
                    routes_data = json.loads(form.routes.data)
                except json.JSONDecodeError:
                    flash("Invalid JSON format for routes", "error")
                    return render_template(
                        "admin/feature_flag_form.html",
                        form=form,
                        feature=feature,
                        action="Edit",
                    )

            feature.key = form.key.data
            feature.name = form.name.data
            feature.description = form.description.data
            feature.feature_type = FeatureType(form.feature_type.data)
            feature.state = FeatureState(form.state.data)
            feature.enabled = form.enabled.data
            feature.sidebar_label = form.sidebar_label.data
            feature.sidebar_icon = form.sidebar_icon.data
            feature.routes = routes_data
            feature.parent_id = form.parent_id.data or None
            feature.sort_order = form.sort_order.data or 0
            feature.last_modified_by = current_user.id

            db.session.commit()

            flash(f"Feature flag '{feature.name}' updated successfully", "success")
            return redirect(url_for("admin.feature_flags"))

        except Exception as e:
            db.session.rollback()
            flash("Error updating feature flag. Please try again.", "error")

    if feature.routes and not form.routes.data:
        form.routes.data = json.dumps(feature.routes, indent=2)
    elif not feature.routes:
        form.routes.data = ""

    return render_template(
        "admin/feature_flag_form.html", form=form, feature=feature, action="Edit"
    )


@admin_bp_v2.route("/feature-flags/<int:id>/toggle", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("toggle_feature_flag")
def feature_flag_toggle(id):
    """Quick toggle feature enabled/disabled."""
    feature = FeatureFlag.query.get_or_404(id)

    try:
        feature.enabled = not feature.enabled
        feature.last_modified_by = current_user.id
        db.session.commit()

        status = "enabled" if feature.enabled else "disabled"
        return jsonify(
            {
                "success": True,
                "enabled": feature.enabled,
                "message": f"Feature {status}",
            }
        )
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "error": "An internal error occurred"}), 500


@admin_bp_v2.route("/feature-flags/<int:id>/delete", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("delete_feature_flag")
def feature_flag_delete(id):
    """Delete feature flag."""
    feature = FeatureFlag.query.get_or_404(id)

    try:
        db.session.delete(feature)
        db.session.commit()
        flash(f"Feature flag '{feature.name}' deleted successfully", "success")
    except Exception as e:
        db.session.rollback()
        flash("Error deleting feature flag. Please try again.", "error")

    return redirect(url_for("admin.feature_flags"))


@admin_bp_v2.route("/feature-flags/discover-sidebar")
@timed_route
@login_required
@admin_required
def feature_flags_discover_sidebar():
    """Discover sidebar menu items for feature flagging."""
    try:
        parser = parse_sidebar_template()
        existing_flags = {ff.key: ff for ff in FeatureFlag.query.all()}

        all_items = []
        for section_title, submenu_title, link in parser.get_all_links():
            feature_key = link.feature_key
            if submenu_title:
                feature_key = f"{SidebarSubmenu(submenu_title, None, section_title, []).feature_key}_{link.feature_key}"

            all_items.append(
                {
                    "section": section_title,
                    "submenu": submenu_title,
                    "text": link.text,
                    "href": link.href,
                    "icon": link.icon,
                    "endpoint": link.endpoint,
                    "feature_key": feature_key,
                    "already_flagged": feature_key in existing_flags,
                    "flag_id": existing_flags[feature_key].id
                    if feature_key in existing_flags
                    else None,
                }
            )

        return render_template(
            "admin/discover_sidebar.html",
            items=all_items,
            sections=parser.sections,
            parser_data=parser.to_dict(),
        )

    except Exception as e:
        flash("Error parsing sidebar. Please try again.", "error")
        return redirect(url_for("admin.feature_flags"))


@admin_bp_v2.route("/feature-flags/discover-sidebar/create", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("create_feature_flags_from_sidebar")
def feature_flags_create_from_sidebar():
    """Create feature flags from selected sidebar items."""
    try:
        selected_keys = request.form.getlist("selected_items")

        if not selected_keys:
            flash("No items selected", "warning")
            return redirect(url_for("admin.feature_flags_discover_sidebar"))

        parser = parse_sidebar_template()

        items_map = {}
        for section_title, submenu_title, link in parser.get_all_links():
            feature_key = link.feature_key
            if submenu_title:
                submenu_obj = SidebarSubmenu(submenu_title, None, section_title, [])
                feature_key = f"{submenu_obj.feature_key}_{link.feature_key}"

            items_map[feature_key] = {
                "section": section_title,
                "submenu": submenu_title,
                "link": link,
            }

        created_count = 0
        skipped_count = 0

        for key in selected_keys:
            existing = FeatureFlag.query.filter_by(key=key).first()
            if existing:
                skipped_count += 1
                continue

            item_data = items_map.get(key)
            if not item_data:
                continue

            link = item_data["link"]
            submenu = item_data["submenu"]
            section = item_data["section"]

            feature = FeatureFlag(
                key=key,
                name=link.text,
                description=f"Controls visibility of '{link.text}' in {section}"
                + (f" > {submenu}" if submenu else ""),
                feature_type=FeatureType.SIDEBAR_LINK,
                state=FeatureState.BETA,
                enabled=True,
                sidebar_label=link.text,
                sidebar_icon=link.icon,
                routes=[link.endpoint] if link.endpoint else [],
                last_modified_by=current_user.id,
            )

            db.session.add(feature)
            created_count += 1

        db.session.commit()

        message = f"Created {created_count} feature flag(s)"
        if skipped_count > 0:
            message += f", skipped {skipped_count} existing"

        flash(message, "success")
        return redirect(url_for("admin.feature_flags"))

    except Exception as e:
        db.session.rollback()
        flash("Error creating feature flags. Please try again.", "error")
        return redirect(url_for("admin.feature_flags_discover_sidebar"))


# ============================================================================
# Abacus Integration Management
# ============================================================================


@admin_bp_v2.route("/abacus-settings", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("update_abacus_settings")
def abacus_settings():
    """Manage Abacus connector configuration."""
    from flask_wtf import FlaskForm
    from wtforms import BooleanField, IntegerField, PasswordField, StringField
    from wtforms.validators import URL, DataRequired

    from app.models.models import ExternalSystem

    class AbacusSettingsForm(FlaskForm):
        base_url = StringField("Base URL", validators=[DataRequired(), URL()])
        client_id = StringField("Client ID", validators=[DataRequired()])
        client_secret = PasswordField("Client Secret")
        enabled = BooleanField("Enable Integration", default=False)
        sync_enabled = BooleanField("Enable Auto-Sync", default=False)
        sync_interval_minutes = IntegerField("Sync Interval (minutes)", default=1440)
        filter_countries = StringField("Country Filter", default="")

    abacus_config = ExternalSystem.query.filter_by(
        system_name="abacus", system_type="ea_tool"
    ).first()

    form = AbacusSettingsForm()

    if form.validate_on_submit():
        import os

        credentials = {}
        if form.client_secret.data and form.client_secret.data.strip():
            credentials = {
                "client_id": form.client_id.data,
                "client_secret": form.client_secret.data,
            }
        elif abacus_config and abacus_config.credentials:
            try:
                existing_creds = json.loads(abacus_config.credentials)
                credentials = {
                    "client_id": form.client_id.data,
                    "client_secret": existing_creds.get("client_secret", ""),
                }
            except (json.JSONDecodeError, KeyError, TypeError):
                logger.exception("Failed to JSON parsing")
                pass

        if abacus_config:
            abacus_config.base_url = form.base_url.data
            abacus_config.api_endpoint = f"{form.base_url.data}/api"
            abacus_config.auth_type = "oauth2"
            if credentials:
                abacus_config.credentials = json.dumps(credentials)

            filter_countries = (form.filter_countries.data or "").strip()
            # Preserve existing dynamic filters from config_json
            existing_config = {}
            if abacus_config.config_json:
                try:
                    existing_config = json.loads(abacus_config.config_json)
                except (json.JSONDecodeError, TypeError):
                    logger.exception("Failed to JSON parsing")
                    pass
            config_data = {
                "filter_countries": filter_countries,
                "filters": existing_config.get("filters", []),
                "discovered_dimensions": existing_config.get("discovered_dimensions", {}),
            }
            # Also accept dynamic filters from form if posted as JSON
            dynamic_filters_json = request.form.get("dynamic_filters", "")
            if dynamic_filters_json:
                try:
                    config_data["filters"] = json.loads(dynamic_filters_json)
                except (json.JSONDecodeError, TypeError):
                    logger.exception("Failed to parse dynamic_filters JSON")
                    pass
            # Save selected entity types
            entity_types_json = request.form.get("sync_entity_types", "")
            if entity_types_json:
                try:
                    config_data["sync_entity_types"] = json.loads(entity_types_json)
                except (json.JSONDecodeError, TypeError):
                    pass
            else:
                config_data["sync_entity_types"] = existing_config.get("sync_entity_types", [])
            # Save sync_relationships flag
            config_data["sync_relationships"] = request.form.get("sync_relationships") == "true"
            abacus_config.config_json = json.dumps(config_data)
            os.environ["ABACUS_FILTER_COUNTRIES"] = filter_countries

            abacus_config.enabled = form.enabled.data
            abacus_config.sync_enabled = form.sync_enabled.data
            abacus_config.sync_interval_minutes = form.sync_interval_minutes.data
            abacus_config.updated_by_id = current_user.id
            abacus_config.updated_at = datetime.utcnow()

            flash("Abacus settings updated successfully.", "success")
        else:
            if not credentials or not credentials.get("client_secret"):
                flash("Client secret is required for new configuration.", "error")
                return render_template(
                    "admin/abacus_settings.html", form=form, abacus_config=None
                )

            filter_countries = (form.filter_countries.data or "").strip()
            config_data = {"filter_countries": filter_countries}
            os.environ["ABACUS_FILTER_COUNTRIES"] = filter_countries

            abacus_config = ExternalSystem(
                system_name="abacus",
                system_type="ea_tool",
                base_url=form.base_url.data,
                api_endpoint=f"{form.base_url.data}/api",
                auth_type="oauth2",
                credentials=json.dumps(credentials),
                config_json=json.dumps(config_data),
                enabled=form.enabled.data,
                sync_enabled=form.sync_enabled.data,
                sync_interval_minutes=form.sync_interval_minutes.data,
                updated_by_id=current_user.id,
            )
            db.session.add(abacus_config)

            flash("Abacus settings created successfully.", "success")

        db.session.commit()
        return redirect(url_for("admin.abacus_settings"))

    env_config = None
    imported_stats = {"applications": 0, "capabilities": 0}

    if request.method == "GET":
        if abacus_config or env_config:
            try:
                from app.models.application_portfolio import ApplicationComponent
                from app.models.business_capabilities import BusinessCapability

                abacus_apps = ApplicationComponent.query.filter_by(
                    abacus_source=True
                ).count()
                imported_stats["applications"] = abacus_apps

                abacus_caps = BusinessCapability.query.filter(
                    BusinessCapability.discovery_source == "abacus",
                ).count()
                imported_stats["capabilities"] = abacus_caps
            except Exception as e:
                logger.error(f"Failed to fetch Abacus imported data counts: {e}")

        if abacus_config:
            form.base_url.data = abacus_config.base_url
            form.enabled.data = abacus_config.enabled
            form.sync_enabled.data = abacus_config.sync_enabled
            form.sync_interval_minutes.data = (
                abacus_config.sync_interval_minutes or 1440
            )

            if abacus_config.config_json:
                try:
                    config_data = json.loads(abacus_config.config_json)
                    form.filter_countries.data = config_data.get("filter_countries", "")
                except (json.JSONDecodeError, KeyError, TypeError):
                    form.filter_countries.data = ""

            if abacus_config.credentials:
                try:
                    creds = json.loads(abacus_config.credentials)
                    form.client_id.data = creds.get("client_id", "")
                except (json.JSONDecodeError, KeyError, TypeError):
                    logger.exception("Failed to parse credentials JSON")
                    pass

        # Load filter dimensions from config_json if previously discovered
        filter_dimensions = {}
        if abacus_config and abacus_config.config_json:
            try:
                cfg = json.loads(abacus_config.config_json)
                filter_dimensions = cfg.get("discovered_dimensions", {})
            except (json.JSONDecodeError, TypeError):
                pass
        else:
            import os

            env_base_url = os.getenv("ABACUS_BASE_URL")
            env_client_id = os.getenv("ABACUS_CLIENT_ID")
            env_client_secret = os.getenv("ABACUS_CLIENT_SECRET")
            env_enabled = os.getenv("ABACUS_ENABLED", "false").lower() == "true"
            env_filter_countries = os.getenv("ABACUS_FILTER_COUNTRIES", "")

            if all([env_base_url, env_client_id, env_client_secret]):
                form.base_url.data = env_base_url
                form.client_id.data = env_client_id
                form.enabled.data = env_enabled
                form.sync_enabled.data = False
                form.sync_interval_minutes.data = 1440
                form.filter_countries.data = env_filter_countries

                env_config = {
                    "base_url": env_base_url,
                    "client_id": env_client_id,
                    "enabled": env_enabled,
                    "filter_countries": env_filter_countries,
                }

    return render_template(
        "admin/abacus_settings.html",
        form=form,
        abacus_config=abacus_config,
        env_config=env_config,
        imported_stats=imported_stats,
        filter_dimensions=filter_dimensions,
    )


@admin_bp_v2.route("/abacus-settings/test-connection", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("test_abacus_connection")
def test_abacus_connection():
    """Test Abacus connection."""
    import asyncio
    import os

    from app.models.models import ExternalSystem

    abacus_config = ExternalSystem.query.filter_by(
        system_name="abacus", system_type="ea_tool"
    ).first()

    if not abacus_config:
        env_base_url = os.getenv("ABACUS_BASE_URL")
        env_client_id = os.getenv("ABACUS_CLIENT_ID")
        env_client_secret = os.getenv("ABACUS_CLIENT_SECRET")

        if not all([env_base_url, env_client_id, env_client_secret]):
            flash(
                "No Abacus configuration found in database or .env file. Please configure first.",
                "error",
            )
            return redirect(url_for("admin.abacus_settings"))

        class TempConfig:
            def __init__(self):
                self.base_url = env_base_url
                self.credentials = None

        abacus_config = TempConfig()
        creds = {"client_id": env_client_id, "client_secret": env_client_secret}
    else:
        creds = (
            json.loads(abacus_config.credentials) if abacus_config.credentials else {}
        )

    try:
        from app.connectors.abacus import create_abacus_connector
        from app.modules.admin.v2.services.connector_framework_v2 import ConnectorConfig

        config = ConnectorConfig(
            name="abacus",
            connector_type="ea_tool",
            config={
                "base_url": abacus_config.base_url,
                "client_id": creds.get("client_id"),
                "client_secret": creds.get("client_secret"),
            },
        )

        connector = create_abacus_connector(config)

        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        result = loop.run_until_complete(connector.test_connection())
        loop.close()

        if hasattr(abacus_config, "last_connection_test"):
            abacus_config.last_connection_test = datetime.utcnow()
            abacus_config.connection_status = "connected" if result else "error"
            abacus_config.last_error = None if result else "Connection test failed"
            db.session.commit()

        if result:
            flash("Connection test successful! Abacus API is accessible.", "success")
        else:
            flash(
                "Connection test failed. Please check your credentials and URL.",
                "error",
            )

    except Exception as e:
        logger.error(f"Abacus connection test exception: {e}", exc_info=True)

        if hasattr(abacus_config, "last_connection_test"):
            abacus_config.last_connection_test = datetime.utcnow()
            abacus_config.connection_status = "error"
            abacus_config.last_error = str(e)
            db.session.commit()

        flash("Connection test failed. Please try again.", "error")

    return redirect(url_for("admin.abacus_settings"))


@admin_bp_v2.route("/abacus-settings/trigger-sync", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("trigger_abacus_sync")
def trigger_abacus_sync():
    """Trigger manual Abacus synchronization."""
    import os

    from app.models.models import ExternalSystem

    abacus_config = ExternalSystem.query.filter_by(system_name="abacus").first()

    if not abacus_config:
        env_base_url = os.getenv("ABACUS_BASE_URL")
        env_client_id = os.getenv("ABACUS_CLIENT_ID")
        env_client_secret = os.getenv("ABACUS_CLIENT_SECRET")

        if not all([env_base_url, env_client_id, env_client_secret]):
            flash(
                "No Abacus configuration found in database or .env file. Please configure first.",
                "error",
            )
            return redirect(url_for("admin.abacus_settings"))

        flash("Using Abacus credentials from .env file.", "info")
    elif not abacus_config.enabled:
        flash("Abacus integration is disabled. Enable it first.", "error")
        return redirect(url_for("admin.abacus_settings"))

    try:
        from app.modules.admin.v2.services.job_queue_service_v2 import (
            get_job_queue_service,
        )

        job_queue = get_job_queue_service()
        job = job_queue.create_job(
            name="Abacus Full Sync", task="abacus_sync", payload={"sync_type": "full"}
        )

        flash(
            f"Sync job created (ID: {job.id}). "
            "The sync will run in the background. Check the status below.",
            "info",
        )

    except Exception as e:
        flash("Failed to create sync job. Please try again.", "error")

    return redirect(url_for("admin.abacus_settings"))


@admin_bp_v2.route("/abacus-settings/sync-status", methods=["GET"])
@timed_route
@login_required
@admin_required
def abacus_sync_status():
    """API endpoint to check current sync job status."""
    from app.models import Job

    latest_job = (
        Job.query.filter_by(task="abacus_sync").order_by(Job.created_at.desc()).first()
    )

    if not latest_job:
        return jsonify({"status": "no_jobs", "message": "No sync jobs found"})

    return jsonify(
        {
            "job_id": latest_job.id,
            "status": latest_job.status,
            "created_at": latest_job.created_at.isoformat()
            if latest_job.created_at
            else None,
            "started_at": latest_job.started_at.isoformat()
            if latest_job.started_at
            else None,
            "finished_at": latest_job.finished_at.isoformat()
            if latest_job.finished_at
            else None,
            "result": latest_job.result,
            "error": latest_job.error,
            "name": latest_job.name,
        }
    )


@admin_bp_v2.route("/abacus-settings/cancel-job/<int:job_id>", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("cancel_abacus_job")
def cancel_abacus_job(job_id):
    """Cancel a running or pending Abacus sync job."""
    try:
        from app.modules.admin.v2.services.job_queue_service_v2 import (
            get_job_queue_service,
        )

        service = get_job_queue_service()
        success = service.cancel_job(job_id)

        if success:
            return (
                jsonify(
                    {"success": True, "message": f"Job {job_id} cancelled successfully"}
                ),
                200,
            )
        else:
            return (
                jsonify(
                    {
                        "success": False,
                        "message": f"Job {job_id} is already finished or cannot be cancelled",
                    }
                ),
                400,
            )
    except Exception as e:
        logger.error(f"Error cancelling job {job_id}: {e}", exc_info=True)
        return jsonify({"success": False, "message": "An internal error occurred"}), 500


@admin_bp_v2.route("/abacus-settings/clear-stale-jobs", methods=["POST"])
@login_required
@admin_required
def clear_stale_abacus_jobs():
    """Force-clear sync jobs stuck in_progress for more than 1 hour."""
    from app.models import Job  # noqa: local import to match pattern

    cutoff = datetime.utcnow() - timedelta(hours=1)
    stale_jobs = Job.query.filter(
        Job.task == "abacus_sync",
        Job.status.in_(["in_progress", "pending"]),
        Job.created_at < cutoff,
    ).all()

    cleared = 0
    for job in stale_jobs:
        job.status = "failed"
        job.finished_at = datetime.utcnow()
        job.error = "Force-cleared: job was stale (>1 hour with no completion)"
        cleared += 1

    if cleared:
        db.session.commit()

    return jsonify({"success": True, "cleared": cleared})


@admin_bp_v2.route("/abacus-settings/discover-types", methods=["POST"])
@login_required
@admin_required
def discover_abacus_types():
    """Discover available ComponentType names from the Abacus API."""
    import asyncio

    from app.models.models import ExternalSystem

    try:
        abacus_config = ExternalSystem.query.filter_by(system_name="abacus").first()
        if not abacus_config:
            return jsonify({"success": False, "error": "No Abacus configuration found"}), 400

        credentials = json.loads(abacus_config.credentials)

        class TempConfig:
            def __init__(self):
                self.config = {
                    "base_url": abacus_config.base_url,
                    "client_id": credentials.get("client_id"),
                    "client_secret": credentials.get("client_secret"),
                }

        from app.connectors.abacus import AbacusConnector

        connector = AbacusConnector(TempConfig())

        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            types = loop.run_until_complete(connector.discover_component_types())
        finally:
            loop.close()

        return jsonify({"success": True, "types": types})

    except Exception as e:
        logger.error("discover_abacus_types failed: %s", e, exc_info=True)
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp_v2.route("/abacus-settings/stats", methods=["GET"])
@timed_route
@login_required
@admin_required
def abacus_stats():
    """Get Abacus import statistics."""
    try:
        from app.models.application_portfolio import ApplicationComponent
        from app.models.business_capabilities import BusinessCapability

        abacus_apps = ApplicationComponent.query.filter_by(abacus_source=True).count()

        abacus_caps = BusinessCapability.query.filter(
            BusinessCapability.discovery_source == "abacus",
        ).count()

        return (
            jsonify(
                {
                    "success": True,
                    "applications": abacus_apps,
                    "capabilities": abacus_caps,
                }
            ),
            200,
        )

    except Exception as e:
        logger.error(f"Error fetching Abacus stats: {e}", exc_info=True)
        return (
            jsonify(
                {
                    "success": False,
                    "applications": 0,
                    "capabilities": 0,
                    "error": str(e),
                }
            ),
            500,
        )


@admin_bp_v2.route("/abacus-settings/discover-filters", methods=["POST"])
@timed_route
@login_required
@admin_required
def discover_abacus_filters():
    """Discover available filter dimensions from the Abacus API.

    Fetches a sample of 50 apps and extracts all unique ConnectionTypeName
    values with their SinkComponentName values. Returns the options so the
    admin can configure dynamic sync filters.
    """
    import asyncio

    from app.models.models import ExternalSystem

    try:
        abacus_config = ExternalSystem.query.filter_by(system_name="abacus").first()
        if not abacus_config:
            return jsonify({"success": False, "error": "No Abacus configuration found"}), 400

        credentials = json.loads(abacus_config.credentials)

        class TempConfig:
            def __init__(self):
                self.config = {
                    "base_url": abacus_config.base_url,
                    "client_id": credentials.get("client_id"),
                    "client_secret": credentials.get("client_secret"),
                }

        from app.connectors.abacus import AbacusConnector

        connector = AbacusConnector(TempConfig())

        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            options = loop.run_until_complete(connector.discover_filter_options())
        finally:
            loop.close()

        # Persist discovered dimensions in config_json so they render on page load
        if options and abacus_config:
            try:
                cfg = json.loads(abacus_config.config_json or "{}")
                cfg["discovered_dimensions"] = options
                abacus_config.config_json = json.dumps(cfg)
                db.session.commit()
            except Exception:
                logger.exception("Failed to persist discovered dimensions")

        return jsonify({"success": True, "dimensions": options})

    except Exception as e:
        logger.error("discover_abacus_filters failed: %s", e, exc_info=True)
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp_v2.route("/abacus-dashboard", methods=["GET"])
@timed_route
@login_required
@admin_required
def abacus_dashboard():
    """Display Abacus sync dashboard with health metrics and statistics."""
    from app.models.application_portfolio import ApplicationComponent
    from app.models.business_capabilities import BusinessCapability
    from app.models.models import ExternalSystem
    from app.modules.admin.v2.services.abacus_sync_service_v2 import get_sync_service

    sync_service = get_sync_service()
    sync_status = sync_service.get_sync_status()

    stats = {
        "applications_synced": 0,
        "applications_created": 0,
        "applications_updated": 0,
        "applications_skipped": 0,
        "capabilities_synced": 0,
        "capabilities_created": 0,
        "capabilities_updated": 0,
        "capabilities_skipped": 0,
    }

    try:
        abacus_apps = ApplicationComponent.query.filter_by(abacus_source=True).count()
        stats["applications_synced"] = abacus_apps

        abacus_caps = BusinessCapability.query.filter(
            BusinessCapability.discovery_source == "abacus",
        ).count()
        stats["capabilities_synced"] = abacus_caps

        external_system = ExternalSystem.query.filter_by(system_name="abacus").first()
        if external_system and external_system.metadata:
            try:
                metadata = (
                    json.loads(external_system.metadata)
                    if isinstance(external_system.metadata, str)
                    else external_system.metadata
                )
                last_stats = metadata.get("last_sync_stats", {})
                if last_stats:
                    stats.update(
                        {
                            "applications_created": last_stats.get(
                                "applications_created", 0
                            ),
                            "applications_updated": last_stats.get(
                                "applications_updated", 0
                            ),
                            "applications_skipped": last_stats.get(
                                "applications_skipped", 0
                            ),
                            "capabilities_created": last_stats.get(
                                "capabilities_created", 0
                            ),
                            "capabilities_updated": last_stats.get(
                                "capabilities_updated", 0
                            ),
                            "capabilities_skipped": last_stats.get(
                                "capabilities_skipped", 0
                            ),
                        }
                    )
            except (json.JSONDecodeError, TypeError, AttributeError):
                logger.exception("Failed to operation")
                pass

    except Exception as e:
        logger.error(f"Failed to fetch sync statistics: {e}", exc_info=True)

    return render_template(
        "admin/abacus_dashboard.html", sync_status=sync_status, stats=stats
    )


# ============================================================================
# Abacus Relationship Mapping Configuration (BPP-011)
# ============================================================================


@admin_bp_v2.route("/abacus-settings/save-relationship-mappings", methods=["POST"])
@login_required
@admin_required
def save_relationship_mappings():
    """Save custom OutConnection → ArchiMate relationship mappings."""
    from app.config.abacus_field_mapping import save_outconnection_mappings
    from app.models.models import ExternalSystem

    try:
        data = request.get_json(silent=True) or {}
        mappings = data.get("mappings", {})

        abacus = ExternalSystem.query.filter_by(system_name="abacus").first()
        if not abacus:
            return jsonify({"error": "Abacus integration not configured"}), 404

        save_outconnection_mappings(abacus, mappings)
        db.session.commit()

        return jsonify({"success": True, "saved_count": len(mappings)})

    except Exception as e:
        db.session.rollback()
        logger.error(f"Failed to save relationship mappings: {e}")
        return jsonify({"error": str(e)}), 500


@admin_bp_v2.route("/abacus-settings/relationship-mappings", methods=["GET"])
@login_required
@admin_required
def get_relationship_mappings():
    """Get current OutConnection → ArchiMate relationship mappings."""
    from app.config.abacus_field_mapping import (
        DEFAULT_OUTCONNECTION_MAPPINGS,
        get_outconnection_mappings,
    )
    from app.models.models import ExternalSystem

    try:
        abacus = ExternalSystem.query.filter_by(system_name="abacus").first()
        mappings = get_outconnection_mappings(abacus)

        return jsonify({
            "mappings": mappings,
            "is_default": mappings == DEFAULT_OUTCONNECTION_MAPPINGS,
            "total": len(mappings),
        })

    except Exception as e:
        logger.error(f"Failed to get relationship mappings: {e}")
        return jsonify({"error": str(e)}), 500


# ============================================================================
# SSO Group-to-Role Mapping (PLT-033)
# ============================================================================

_VALID_ROLES = [
    "solution_architect",
    "enterprise_architect",
    "arb_member",
    "portfolio_manager",
    "platform_admin",
]


@admin_bp_v2.route("/sso-settings", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("update_sso_settings")
def sso_settings():
    """Manage SSO group-to-role mappings stored in the database."""
    from app.models.miscellaneous import SSOGroupRoleMapping

    if request.method == "POST":
        action = request.form.get("action")

        if action == "add":
            group_name = (request.form.get("sso_group_name") or "").strip()
            role_name = request.form.get("role_name", "").strip()
            description = (request.form.get("description") or "").strip()

            if not group_name:
                flash("Group name is required.", "error")
                return redirect(url_for("admin.sso_settings"))

            if role_name not in _VALID_ROLES:
                flash(f"Invalid role: {role_name!r}.", "error")
                return redirect(url_for("admin.sso_settings"))

            existing = SSOGroupRoleMapping.query.filter_by(
                sso_group_name=group_name
            ).first()
            if existing:
                flash(
                    f"A mapping for group '{group_name}' already exists. Edit it instead.",
                    "error",
                )
                return redirect(url_for("admin.sso_settings"))

            mapping = SSOGroupRoleMapping(
                sso_group_name=group_name,
                role_name=role_name,
                description=description or None,
                is_active=True,
            )
            db.session.add(mapping)
            db.session.commit()
            logger.info(
                "SSO mapping added: group=%r role=%r by user=%s",
                group_name,
                role_name,
                current_user.email,
            )
            flash(f"Mapping '{group_name}' → '{role_name}' added.", "success")

        elif action == "delete":
            mapping_id = request.form.get("mapping_id", type=int)
            mapping = SSOGroupRoleMapping.query.get(mapping_id)
            if not mapping:
                flash("Mapping not found.", "error")
                return redirect(url_for("admin.sso_settings"))
            group_name = mapping.sso_group_name
            db.session.delete(mapping)
            db.session.commit()
            logger.info(
                "SSO mapping deleted: group=%r by user=%s",
                group_name,
                current_user.email,
            )
            flash(f"Mapping for group '{group_name}' deleted.", "success")

        elif action == "toggle":
            mapping_id = request.form.get("mapping_id", type=int)
            mapping = SSOGroupRoleMapping.query.get(mapping_id)
            if not mapping:
                flash("Mapping not found.", "error")
                return redirect(url_for("admin.sso_settings"))
            mapping.is_active = not mapping.is_active
            db.session.commit()
            state = "enabled" if mapping.is_active else "disabled"
            flash(f"Mapping for group '{mapping.sso_group_name}' {state}.", "success")

        elif action == "sync_all":
            # Re-apply current DB mappings to all SSO users.
            from app.models.user import User
            from app.auth.sso import sso_service

            sso_users = User.query.filter(User.sso_provider.isnot(None)).all()
            updated = 0
            for user in sso_users:
                # Re-evaluate via sso_service which now reads DB mappings.
                # We pass the user's stored external groups claim if available;
                # since we don't persist groups, re-evaluate by username match only.
                # A full re-evaluation requires an active SSO session.  Here we
                # apply the highest-priority DB mapping that matches any group
                # whose name equals the user's sso_provider+email pattern.
                # Best-effort: if user has no stored groups we skip.
                pass

            # Refresh sso_service in-memory cache from DB.
            db_mappings = SSOGroupRoleMapping.query.filter_by(is_active=True).all()
            sso_service._group_role_map = {
                m.sso_group_name: m.role_name for m in db_mappings
            }
            flash(
                f"SSO role map refreshed from database ({len(db_mappings)} active mappings). "
                "Roles will be applied on next SSO login.",
                "success",
            )
            logger.info(
                "SSO group-role map refreshed from DB by user=%s (%d mappings)",
                current_user.email,
                len(db_mappings),
            )

        else:
            flash("Unknown action.", "error")

        return redirect(url_for("admin.sso_settings"))

    # GET
    mappings = SSOGroupRoleMapping.query.order_by(
        SSOGroupRoleMapping.sso_group_name
    ).all()
    return render_template(
        "admin/sso_settings.html",
        mappings=mappings,
        valid_roles=_VALID_ROLES,
    )


# ============================================================================
# Jira Push Integration
# ============================================================================


@admin_bp_v2.route("/jira-settings", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
def jira_settings():
    """Manage Jira push integration configuration."""
    from flask_wtf import FlaskForm
    from wtforms import BooleanField, PasswordField, StringField
    from wtforms.validators import DataRequired

    from app.models.models import ExternalSystem

    class JiraSettingsForm(FlaskForm):
        base_url = StringField("Jira Base URL", validators=[DataRequired()])
        username = StringField("Username / Email", validators=[DataRequired()])
        api_token = PasswordField("API Token")
        project_key = StringField("Project Key", validators=[DataRequired()])
        issue_type = StringField("Issue Type", default="Task")
        filter_countries = StringField("Country Filter", default="United Kingdom")
        enabled = BooleanField("Enable Integration", default=False)

    jira_config = ExternalSystem.query.filter_by(
        system_name="jira", system_type="alm"
    ).first()

    form = JiraSettingsForm()

    if form.validate_on_submit():
        credentials = {}
        if form.api_token.data and form.api_token.data.strip():
            credentials = {
                "username": form.username.data,
                "api_token": form.api_token.data,
            }
        elif jira_config and jira_config.credentials:
            try:
                existing_creds = json.loads(jira_config.credentials)
                credentials = {
                    "username": form.username.data,
                    "api_token": existing_creds.get("api_token", ""),
                }
            except (json.JSONDecodeError, TypeError):
                logger.exception("Failed to JSON parsing")
                pass

        filter_countries = (form.filter_countries.data or "").strip()
        config_data = {
            "project_key": form.project_key.data,
            "issue_type": form.issue_type.data or "Task",
            "filter_countries": [c.strip() for c in filter_countries.split(",") if c.strip()],
            "custom_field_map": {},
        }

        if jira_config:
            jira_config.base_url = form.base_url.data
            jira_config.api_endpoint = f"{form.base_url.data}/rest/api/3"
            jira_config.auth_type = "basic"
            if credentials:
                jira_config.credentials = json.dumps(credentials)
            jira_config.config_json = json.dumps(config_data)
            jira_config.enabled = form.enabled.data
            jira_config.updated_by_id = current_user.id
            jira_config.updated_at = datetime.utcnow()
            flash("Jira settings updated successfully.", "success")
        else:
            if not credentials or not credentials.get("api_token"):
                flash("API token is required for new configuration.", "error")
                _env_url = os.environ.get("JIRA_BASE_URL", "").strip()
                _env_user = os.environ.get("JIRA_USERNAME", "").strip()
                _env_project = os.environ.get("JIRA_PROJECT_KEY", "").strip()
                _env_configured = bool(_env_url and _env_user and _env_project)
                return render_template(
                    "admin/jira_settings.html", form=form, jira_config=None, push_status=None,
                    env_configured=_env_configured,
                )
            jira_config = ExternalSystem(
                system_name="jira",
                system_type="alm",
                base_url=form.base_url.data,
                api_endpoint=f"{form.base_url.data}/rest/api/3",
                auth_type="basic",
                credentials=json.dumps(credentials),
                config_json=json.dumps(config_data),
                enabled=form.enabled.data,
                updated_by_id=current_user.id,
            )
            db.session.add(jira_config)
            flash("Jira settings created successfully.", "success")

        db.session.commit()
        return redirect(url_for("admin.jira_settings"))

    # Detect .env configuration regardless of DB record
    # Any non-empty JIRA_BASE_URL indicates env-based config; JIRA_API_TOKEN may load at boot
    env_jira_url = os.environ.get("JIRA_BASE_URL", "").strip()
    env_jira_user = os.environ.get("JIRA_USERNAME", "").strip()
    env_jira_token = os.environ.get("JIRA_API_TOKEN", "").strip()
    env_jira_project = os.environ.get("JIRA_PROJECT_KEY", "").strip()
    # Fallback: check if the connector can initialize (handles .env loading internally)
    env_configured = bool(env_jira_url and env_jira_user and env_jira_project)

    push_status = None
    if request.method == "GET":
        if jira_config:
            try:
                creds = json.loads(jira_config.credentials) if jira_config.credentials else {}
            except (json.JSONDecodeError, TypeError):
                creds = {}
            try:
                cfg = json.loads(jira_config.config_json) if jira_config.config_json else {}
            except (json.JSONDecodeError, TypeError):
                cfg = {}

            form.base_url.data = jira_config.base_url
            form.username.data = creds.get("username", "")
            form.project_key.data = cfg.get("project_key", "")
            form.issue_type.data = cfg.get("issue_type", "Task")
            countries = cfg.get("filter_countries", [])
            form.filter_countries.data = ", ".join(countries) if isinstance(countries, list) else str(countries)
            form.enabled.data = jira_config.enabled
        else:
            # Pre-populate from .env when no DB record exists yet
            form.base_url.data = os.environ.get("JIRA_BASE_URL", "")
            form.username.data = os.environ.get("JIRA_USERNAME", "")
            form.project_key.data = os.environ.get("JIRA_PROJECT_KEY", "")

        try:
            from app.services.jira_push_service import get_jira_push_service
            svc = get_jira_push_service()
            push_status = svc.get_push_status()
            if push_status and jira_config and jira_config.config_json:
                cfg_data = json.loads(jira_config.config_json) if isinstance(jira_config.config_json, str) else {}
                last_wh = cfg_data.get("last_webhook_at")
                if last_wh:
                    push_status["last_webhook_at"] = last_wh
        except Exception:
            push_status = None

    return render_template(
        "admin/jira_settings.html",
        form=form,
        jira_config=jira_config,
        push_status=push_status,
        env_configured=env_configured,
    )


@admin_bp_v2.route("/jira-settings/test-connection", methods=["POST"])
@timed_route
@login_required
@admin_required
def jira_test_connection():
    """Test Jira API connectivity."""
    import asyncio

    from app.models.models import ExternalSystem

    jira_config = ExternalSystem.query.filter_by(
        system_name="jira", system_type="alm"
    ).first()

    if not jira_config:
        # Try .env fallback — connector handles it via _initialize_connector()
        try:
            from app.services.jira_push_service import get_jira_push_service
            svc = get_jira_push_service()
            if not svc._initialize_connector():
                return jsonify({"success": False, "message": "Jira not configured. Please fill in and save your Jira credentials first, then test the connection."}), 400
            loop = asyncio.new_event_loop()
            try:
                connected = loop.run_until_complete(svc.connector.test_connection())
            finally:
                loop.close()
            return jsonify({
                "success": connected,
                "message": "Connection successful (via .env)" if connected else "Connection failed — check .env credentials",
            })
        except Exception as e:
            return jsonify({"success": False, "message": str(e)}), 500

    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        if not svc._initialize_connector():
            return jsonify({
                "success": False,
                "message": "Failed to connect: Jira credentials are saved but invalid or incomplete. Check the URL, username, API token and project key."
            }), 400

        loop = asyncio.new_event_loop()
        try:
            connected = loop.run_until_complete(svc.connector.test_connection())
        finally:
            loop.close()

        jira_config.last_connection_test = datetime.utcnow()
        jira_config.connection_status = "connected" if connected else "error"
        db.session.commit()

        return jsonify({
            "success": connected,
            "message": "Connection successful" if connected else "Connection failed",
        })
    except Exception as e:
        jira_config.connection_status = "error"
        jira_config.last_error = str(e)
        db.session.commit()
        return jsonify({"success": False, "message": str(e)}), 500


@admin_bp_v2.route("/jira-settings/webhook", methods=["POST"])
@login_required
# csrf.exempt: webhook receiver — external systems cannot include CSRF tokens
@csrf.exempt
def jira_webhook():
    """Receive Jira webhook events and sync status back to KanbanCard."""
    from app.models.models import ExternalSystem

    # HMAC verification
    secret = os.environ.get("JIRA_WEBHOOK_SECRET", "")
    if secret:
        sig = request.headers.get("X-Hub-Signature", "")
        expected = "sha256=" + hmac.new(secret.encode(), request.data, hashlib.sha256).hexdigest()
        if not hmac.compare_digest(expected, sig):
            return jsonify({"error": "Invalid signature"}), 403

    data = request.get_json(silent=True) or {}
    issue_key = data.get("issue", {}).get("key", "")
    if not issue_key:
        return jsonify({"ok": True}), 200

    # Extract status change from changelog
    changelog = data.get("changelog", {})
    items = changelog.get("items", [])
    new_status = None
    for item in items:
        if item.get("field") == "status":
            new_status = item.get("toString")
            break

    # Reverse-map Jira status to platform column
    if new_status:
        try:
            from app.config.kanban_jira_field_mapping import COLUMN_TO_JIRA_STATUS
            from app.models.adm_kanban_models import KanbanCard
            jira_to_column = {v: k for k, v in COLUMN_TO_JIRA_STATUS.items()}
            new_column = jira_to_column.get(new_status)
            if new_column:
                card = KanbanCard.query.filter_by(jira_issue_key=issue_key).first()
                if card:
                    card.status = new_column
                    card.updated_at = datetime.utcnow()
                    db.session.commit()
        except Exception as _card_exc:
            logger.warning("Jira webhook card update failed for %s: %s", issue_key, _card_exc)

    # Store last webhook timestamp
    try:
        jira_config = ExternalSystem.query.filter_by(system_name="jira").first()
        if jira_config:
            cfg = json.loads(jira_config.config_json) if jira_config.config_json else {}
            cfg["last_webhook_at"] = datetime.utcnow().isoformat()
            jira_config.config_json = json.dumps(cfg)
            db.session.commit()
    except Exception as _ts_exc:
        logger.warning("Jira webhook timestamp update failed: %s", _ts_exc)

    return jsonify({"ok": True}), 200


@admin_bp_v2.route("/jira-settings/save-env-config", methods=["POST"])
@login_required
@admin_required
def save_env_jira_config():
    """Save .env Jira credentials to database."""
    import os

    from app.models.models import ExternalSystem

    url = os.environ.get("JIRA_BASE_URL", "").strip()
    user = os.environ.get("JIRA_USERNAME", "").strip()
    token = os.environ.get("JIRA_API_TOKEN", "").strip()
    proj = os.environ.get("JIRA_PROJECT_KEY", "").strip()

    missing = [k for k, v in [
        ("JIRA_BASE_URL", url), ("JIRA_USERNAME", user),
        ("JIRA_API_TOKEN", token), ("JIRA_PROJECT_KEY", proj),
    ] if not v]
    if missing:
        flash(f"Missing required Jira environment variables: {', '.join(missing)}", "error")
        return redirect(url_for("admin.jira_settings"))

    jira_config = ExternalSystem.query.filter_by(
        system_name="jira", system_type="alm"
    ).first()
    if jira_config is None:
        jira_config = ExternalSystem(system_name="jira", system_type="alm")
        db.session.add(jira_config)

    jira_config.credentials = json.dumps({
        "base_url": url,
        "username": user,
        "api_token": token,
        "project_key": proj,
    })
    jira_config.enabled = True
    db.session.commit()

    flash("Jira configuration saved to database.", "success")
    return redirect(url_for("admin.jira_settings"))


@admin_bp_v2.route("/jira-settings/trigger-push", methods=["POST"])
@timed_route
@login_required
@admin_required
def jira_trigger_push():
    """Create a Job and start pushing applications to Jira."""
    from app.models.job import Job, JobStatus

    force = request.json.get("force", False) if request.is_json else False
    app_ids = request.json.get("app_ids") if request.is_json else None

    job = Job(
        name="jira_push",
        task="push_applications",
        payload={"app_ids": app_ids, "force": force},
        status=JobStatus.IN_PROGRESS.value,
        started_at=datetime.utcnow(),
    )
    db.session.add(job)
    db.session.commit()

    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_applications(app_ids=app_ids, force=force)

        job.status = JobStatus.COMPLETED.value
        job.result = result.as_dict()
        job.finished_at = datetime.utcnow()
        db.session.commit()

        return jsonify({"success": True, "job_id": job.id, "result": result.as_dict()})
    except Exception as e:
        job.status = JobStatus.FAILED.value
        job.error = str(e)
        job.finished_at = datetime.utcnow()
        db.session.commit()
        return jsonify({"success": False, "job_id": job.id, "error": str(e)}), 500


@admin_bp_v2.route("/jira-settings/push-status", methods=["GET"])
@timed_route
@login_required
@admin_required
def jira_push_status():
    """Return JSON push status for polling."""
    from app.models.job import Job

    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        status = svc.get_push_status()
    except Exception:
        status = {}

    latest_job = (
        Job.query.filter_by(name="jira_push")
        .order_by(Job.created_at.desc())
        .first()
    )

    return jsonify({
        "push_status": status,
        "latest_job": latest_job.as_dict() if latest_job else None,
    })


@admin_bp_v2.route("/jira-settings/kanban-push-status", methods=["GET"])
@timed_route
@login_required
@admin_required
def jira_kanban_push_status():
    """Return JSON kanban push status for polling."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        kanban_status = svc.get_kanban_push_status()
        return jsonify({"success": True, "kanban_status": kanban_status})
    except Exception as exc:
        logger.error("Kanban push status error: %s", exc)
        return jsonify({"success": False, "error": str(exc)}), 500


@admin_bp_v2.route("/jira-settings/trigger-kanban-push", methods=["POST"])
@timed_route
@login_required
@admin_required
def jira_trigger_kanban_push():
    """Push all unpushed KanbanCard rows to Jira."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_all_unpushed_cards()
        return jsonify({"success": True, **result})
    except Exception as exc:
        logger.error("Trigger kanban push error: %s", exc)
        return jsonify({"success": False, "error": str(exc)}), 500


@admin_bp_v2.route("/jira-settings/push-epics", methods=["POST"])
@timed_route
@login_required
@admin_required
def jira_push_epics():
    """Create one Jira Epic per ADM phase as an ArchiMate Plateau."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_adm_phase_epics()
        return jsonify({"success": True, **result})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp_v2.route("/jira-settings/push-applications", methods=["POST"])
@timed_route
@login_required
@admin_required
def jira_push_applications():
    """Push ApplicationComponents (ArchiMate Application Layer, Phase C/D) to Jira."""
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_applications()
        return jsonify({"success": True, **result.as_dict()})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp_v2.route("/jira-settings/push-dependencies", methods=["POST"])
@timed_route
@login_required
@admin_required
def jira_push_dependencies():
    """Create Jira Subtasks from KanbanCard.depends_on (ArchiMate TriggeringRelationship).

    Also creates Blocks issue links from KanbanCard.blocks (ArchiMate AssociationRelationship).
    """
    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        result = svc.push_card_dependencies()
        return jsonify({"success": True, **result})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@admin_bp_v2.route("/jira-settings/field-discovery", methods=["GET"])
@timed_route
@login_required
@admin_required
def jira_field_discovery():
    """Trigger discover_fields and return available Jira fields."""
    import asyncio

    from app.models.models import ExternalSystem

    jira_config = ExternalSystem.query.filter_by(
        system_name="jira", system_type="alm"
    ).first()

    if not jira_config:
        return jsonify({"success": False, "message": "Jira not configured"}), 400

    try:
        cfg = json.loads(jira_config.config_json) if jira_config.config_json else {}
    except (json.JSONDecodeError, TypeError):
        cfg = {}

    project_key = cfg.get("project_key", "EA")
    issue_type_id = request.args.get("issue_type_id", "10001")

    try:
        from app.services.jira_push_service import get_jira_push_service
        svc = get_jira_push_service()
        if not svc._initialize_connector():
            return jsonify({"success": False, "message": "Connector init failed"}), 400

        loop = asyncio.new_event_loop()
        try:
            fields = loop.run_until_complete(
                svc.connector.discover_fields(project_key, issue_type_id)
            )
        finally:
            loop.close()

        return jsonify({"success": True, "fields": fields})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


# ============================================================================
# Seed Management
# ============================================================================


@admin_bp_v2.route("/seed-management")
@timed_route
@login_required
@admin_required
def seed_management():
    """Seed management dashboard."""
    from app.modules.admin.v2.services.seed_management_service_v2 import (
        SeedManagementService,
    )

    service = SeedManagementService()
    categories = service.get_seed_status()

    return render_template("admin/seed_management.html", categories=categories)


@admin_bp_v2.route("/api/seed-status")
@timed_route
@login_required
@admin_required
def seed_status():
    """API: Get current seed status."""
    from app.modules.admin.v2.services.seed_management_service_v2 import (
        SeedManagementService,
    )

    service = SeedManagementService()
    categories = service.get_seed_status()

    return jsonify({"success": True, "categories": categories})


@admin_bp_v2.route("/api/seed/<key>", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("seed_data")
def seed(key):
    """API: Trigger seeding for a specific category."""
    from app.modules.admin.v2.services.seed_management_service_v2 import (
        SeedManagementService,
    )

    service = SeedManagementService()
    result = service.seed_category(key)
    status_code = 200 if result.get("success") else 400
    return jsonify(result), status_code


@admin_bp_v2.route("/api/seed-all", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("seed_all_data")
def seed_all():
    """API: Trigger all seeders."""
    from app.modules.admin.v2.services.seed_management_service_v2 import (
        SeedManagementService,
    )

    service = SeedManagementService()
    result = service.seed_all()
    status_code = 200 if result.get("success") else 400
    return jsonify(result), status_code


# ============================================================================
# Feature Auto-Discovery Helper
# ============================================================================


def _auto_discover_features(app):
    """Auto-discover features from Flask app blueprints and routes.

    Creates a hierarchical feature tree:
    - Level 1: Blueprints (parent features)
    - Level 2: Routes within blueprints (child features)
    - Level 3: Sub-routes (grandchild features)

    Filters out infrastructure routes to reduce noise.

    Args:
        app: Flask application instance

    Returns:
        int: Number of features created
    """
    from collections import defaultdict

    created_count = 0
    blueprint_features = {}

    existing_keys = {f.key for f in FeatureFlag.query.all()}

    SKIP_ENDPOINTS = {
        "static",
        "_debug_toolbar",
        "debugtoolbar",
        "health",
        "metrics",
        "status",
    }
    SKIP_PREFIXES = ("_", "static", "/health", "/metrics", "/_", "/api/")
    SKIP_BLUEPRINTS = {"debugtoolbar", "static"}

    routes_by_blueprint = defaultdict(list)
    skipped_count = 0
    total_count = 0

    for rule in app.url_map.iter_rules():
        total_count += 1

        if any(rule.endpoint.startswith(skip) for skip in SKIP_ENDPOINTS):
            skipped_count += 1
            continue
        if any(rule.rule.startswith(skip) for skip in SKIP_PREFIXES):
            skipped_count += 1
            continue

        parts = rule.endpoint.split(".")
        if len(parts) >= 2:
            blueprint_name = parts[0]

            if blueprint_name in SKIP_BLUEPRINTS:
                skipped_count += 1
                continue

            routes_by_blueprint[blueprint_name].append(
                {
                    "endpoint": rule.endpoint,
                    "rule": rule.rule,
                    "methods": list(rule.methods - {"HEAD", "OPTIONS"}),
                }
            )
        else:
            skipped_count += 1

    app.logger.info(
        f"Route discovery: {total_count} total routes, {skipped_count} skipped, {len(routes_by_blueprint)} blueprints found"
    )

    for blueprint_name in sorted(routes_by_blueprint.keys()):
        blueprint_key = f"blueprint_{blueprint_name}"

        if blueprint_key in existing_keys:
            parent_feature = FeatureFlag.query.filter_by(key=blueprint_key).first()
            blueprint_features[blueprint_name] = parent_feature
            continue

        display_name = blueprint_name.replace("_", " ").title()

        parent_feature = FeatureFlag(
            key=blueprint_key,
            name=f"{display_name} Module",
            description=f"All features within the {display_name} blueprint",
            feature_type=FeatureType.BLUEPRINT,
            state=FeatureState.STABLE,
            enabled=True,
            sort_order=created_count,
        )

        try:
            db.session.add(parent_feature)
            db.session.flush()
            blueprint_features[blueprint_name] = parent_feature
            existing_keys.add(blueprint_key)
            created_count += 1
        except Exception as e:
            db.session.rollback()
            app.logger.warning(
                f"Failed to create blueprint feature {blueprint_key}: {e}"
            )
            parent_feature = FeatureFlag.query.filter_by(key=blueprint_key).first()
            if parent_feature:
                blueprint_features[blueprint_name] = parent_feature
            continue

        base_paths = defaultdict(list)
        for route_info in routes_by_blueprint[blueprint_name]:
            rule = route_info["rule"]
            path_parts = [p for p in rule.split("/") if p and not p.startswith("<")]
            base_path = "/" + "/".join(path_parts[:3]) if len(path_parts) >= 2 else rule
            base_paths[base_path].append(route_info)

        for base_path, route_list in sorted(base_paths.items()):
            if len(route_list) == 1 and route_list[0]["rule"] == base_path:
                route_info = route_list[0]
                route_key = f"route_{route_info['endpoint'].replace('.', '_')}"

                if route_key in existing_keys:
                    continue

                function_name = route_info["endpoint"].split(".")[-1]
                display_name = function_name.replace("_", " ").title()

                route_feature = FeatureFlag(
                    key=route_key,
                    name=display_name,
                    description=f"{display_name} ({route_info['rule']})",
                    feature_type=FeatureType.ROUTE,
                    state=FeatureState.STABLE,
                    enabled=True,
                    routes=[route_info["rule"]],
                    parent_id=parent_feature.id,
                    sort_order=created_count,
                )

                try:
                    db.session.add(route_feature)
                    existing_keys.add(route_key)
                    created_count += 1
                except Exception as e:
                    db.session.rollback()
                    app.logger.warning(
                        f"Failed to create route feature {route_key}: {e}"
                    )
                    continue
            else:
                section_key = (
                    f"section_{blueprint_name}_{base_path.replace('/', '_').strip('_')}"
                )

                section_feature = None
                if section_key in existing_keys:
                    section_feature = FeatureFlag.query.filter_by(
                        key=section_key
                    ).first()
                else:
                    section_name = (
                        base_path.split("/")[-1]
                        .replace("-", " ")
                        .replace("_", " ")
                        .title()
                    )

                    section_feature = FeatureFlag(
                        key=section_key,
                        name=f"{section_name} Section",
                        description=f"All routes under {base_path}",
                        feature_type=FeatureType.SIDEBAR_SECTION,
                        state=FeatureState.STABLE,
                        enabled=True,
                        routes=[base_path + "*"],
                        parent_id=parent_feature.id,
                        sort_order=created_count,
                    )

                    try:
                        db.session.add(section_feature)
                        db.session.flush()
                        existing_keys.add(section_key)
                        created_count += 1
                    except Exception as e:
                        db.session.rollback()
                        app.logger.warning(
                            f"Failed to create section feature {section_key}: {e}"
                        )
                        section_feature = FeatureFlag.query.filter_by(
                            key=section_key
                        ).first()

                if not section_feature:
                    continue

                for route_info in route_list:
                    route_key = f"route_{route_info['endpoint'].replace('.', '_')}"

                    if route_key in existing_keys:
                        continue

                    function_name = route_info["endpoint"].split(".")[-1]
                    display_name = function_name.replace("_", " ").title()

                    route_feature = FeatureFlag(
                        key=route_key,
                        name=display_name,
                        description=f"{display_name} ({route_info['rule']}) - {', '.join(route_info['methods'])}",
                        feature_type=FeatureType.ROUTE,
                        state=FeatureState.STABLE,
                        enabled=True,
                        routes=[route_info["rule"]],
                        parent_id=section_feature.id,
                        sort_order=created_count,
                    )

                    try:
                        db.session.add(route_feature)
                        existing_keys.add(route_key)
                        created_count += 1
                    except Exception as e:
                        db.session.rollback()
                        app.logger.warning(
                            f"Failed to create route feature {route_key}: {e}"
                        )
                        continue

    try:
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"Failed to commit feature discovery: {e}")
        raise

    return created_count


# ============================================================================
# Users Table API — canonical data table endpoints
# ============================================================================


@admin_bp_v2.route("/api/users", methods=["GET"])
@timed_route
@login_required
@admin_required
def api_list_users():
    """Paginated user list API for canonical data table."""
    from sqlalchemy.orm import joinedload

    page = request.args.get("page", 1, type=int)
    per_page = min(request.args.get("per_page", 25, type=int), 100)
    search = request.args.get("q") or request.args.get("search", "")
    role_filter = request.args.get("role", "")
    sort_by = request.args.get("sort", "id")
    sort_dir = request.args.get("dir", "desc")

    ALLOWED_SORT = {"id", "first_name", "last_name", "email", "created_at"}
    if sort_by not in ALLOWED_SORT:
        sort_by = "id"

    query = User.query.options(joinedload(User.role))
    if search:
        term = f"%{search}%"
        query = query.filter(
            User.first_name.ilike(term) | User.last_name.ilike(term) | User.email.ilike(term)
        )
    if role_filter:
        query = query.join(User.role).filter(Role.name == role_filter)

    sort_col = getattr(User, sort_by, User.id)
    if sort_dir == "asc":
        query = query.order_by(sort_col.asc())
    else:
        query = query.order_by(sort_col.desc())

    paginated = query.paginate(page=page, per_page=per_page, error_out=False)

    items = []
    for idx, u in enumerate(paginated.items):
        items.append({
            "id": u.id,
            "row_number": (page - 1) * per_page + idx + 1,
            "first_name": u.first_name or "",
            "last_name": u.last_name or "",
            "email": u.email,
            "role": u.role.name if u.role else "No Role",
            "confirmed": u.confirmed,
        })

    return jsonify({
        "users": items,
        "total": paginated.total,
        "page": page,
        "pages": paginated.pages,
        "per_page": per_page,
    })


@admin_bp_v2.route("/api/users/bulk", methods=["DELETE"])
@timed_route
@login_required
@admin_required
def api_bulk_delete_users():
    """Bulk delete users by IDs (cannot delete yourself)."""
    data = request.get_json() or {}
    ids = data.get("ids", [])
    if not ids or not isinstance(ids, list):
        return jsonify({"error": "ids list required"}), 400
    safe_ids = [i for i in ids if i != current_user.id]
    deleted = User.query.filter(User.id.in_(safe_ids)).delete(synchronize_session=False)
    db.session.commit()
    return jsonify({"deleted": deleted})


# ── Roles API ──────────────────────────────────────────────────────────────────

@admin_bp_v2.route("/api/roles", methods=["GET"])
@timed_route
@login_required
@admin_required
def api_list_roles():
    """List all roles with user counts and permission flags."""
    roles = Role.query.order_by(Role.name).all()
    items = []
    for role in roles:
        users = User.query.filter_by(role_id=role.id).all()
        items.append({
            "id": role.id,
            "name": role.name,
            "type": "system" if role.name in ("Administrator", "User") else "custom",
            "status": "active",
            "description": f"{'Full system access and administration.' if role.name == 'Administrator' else 'Standard user access.'}" if role.name in ("Administrator", "User") else "",
            "permissions": {
                "administer": bool(role.permissions and (role.permissions & 0xFF) == 0xFF),
                "general": bool(role.permissions and role.permissions & 0x01),
            },
            "userCount": len(users),
            "users": [u.email for u in users[:10]],
            "created": None,
            "lastModified": None,
        })
    return jsonify({"success": True, "roles": items})


@admin_bp_v2.route("/api/roles/<int:role_id>", methods=["GET"])
@timed_route
@login_required
@admin_required
def api_get_role(role_id):
    """Get a single role by ID."""
    role = Role.query.get_or_404(role_id)
    users = User.query.filter_by(role_id=role.id).all()
    return jsonify({
        "success": True,
        "role": {
            "id": role.id,
            "name": role.name,
            "permissions": {
                "administer": bool(role.permissions and (role.permissions & 0xFF) == 0xFF),
                "general": bool(role.permissions and role.permissions & 0x01),
            },
            "userCount": len(users),
            "users": [{"id": u.id, "email": u.email, "name": f"{u.first_name or ''} {u.last_name or ''}".strip()} for u in users],
        }
    })


@admin_bp_v2.route("/api/roles", methods=["POST"])
@timed_route
@login_required
@admin_required
def api_create_role():
    """Create a new custom role."""
    data = request.get_json() or {}
    name = (data.get("name") or "").strip()
    if not name:
        return jsonify({"success": False, "error": "Role name is required"}), 400
    if Role.query.filter_by(name=name).first():
        return jsonify({"success": False, "error": f"Role '{name}' already exists"}), 409
    from app.models.user import Permission as PermBits
    perms = PermBits.ADMINISTER if data.get("is_admin") else PermBits.GENERAL
    role = Role(name=name, permissions=perms, index="main", default=False)
    db.session.add(role)
    db.session.commit()
    return jsonify({"success": True, "role": {"id": role.id, "name": role.name}}), 201


@admin_bp_v2.route("/api/roles/<int:role_id>", methods=["PUT"])
@timed_route
@login_required
@admin_required
def api_update_role(role_id):
    """Update a role name or permissions."""
    role = Role.query.get_or_404(role_id)
    if role.name in ("Administrator", "User"):
        return jsonify({"success": False, "error": "System roles cannot be modified"}), 403
    data = request.get_json() or {}
    if "name" in data and data["name"].strip():
        role.name = data["name"].strip()
    from app.models.user import Permission as PermBits
    if "is_admin" in data:
        role.permissions = PermBits.ADMINISTER if data["is_admin"] else PermBits.GENERAL
    db.session.commit()
    return jsonify({"success": True})


@admin_bp_v2.route("/api/roles/<int:role_id>", methods=["DELETE"])
@timed_route
@login_required
@admin_required
def api_delete_role(role_id):
    """Delete a custom role. Reassigns users to the default User role."""
    role = Role.query.get_or_404(role_id)
    if role.name in ("Administrator", "User"):
        return jsonify({"success": False, "error": "System roles cannot be deleted"}), 403
    default_role = Role.query.filter_by(default=True).first()
    User.query.filter_by(role_id=role.id).update({"role_id": default_role.id if default_role else None})
    db.session.delete(role)
    db.session.commit()
    return jsonify({"success": True})


# ── Permissions / Enterprise Roles API ────────────────────────────────────────

@admin_bp_v2.route("/api/enterprise-roles", methods=["GET"])
@timed_route
@login_required
@admin_required
def api_list_enterprise_roles():
    """List enterprise role assignments across all users."""
    from app.models.user import VALID_ROLES
    role_counts = {}
    for r in VALID_ROLES:
        role_counts[r] = User.query.filter_by(enterprise_role=r).count()

    items = []
    for r in VALID_ROLES:
        label = r.replace("_", " ").title()
        items.append({
            "id": r,
            "name": label,
            "type": "enterprise",
            "resource": "platform",
            "action": "access",
            "scope": "organization",
            "description": f"Enterprise platform role: {label}",
            "level": 1,
            "userCount": role_counts.get(r, 0),
        })
    return jsonify({"success": True, "permissions": items})


@admin_bp_v2.route("/api/enterprise-roles/users", methods=["GET"])
@timed_route
@login_required
@admin_required
def api_enterprise_role_users():
    """List all users with their enterprise role assignments."""
    users = User.query.order_by(User.last_name, User.first_name).all()
    items = []
    for u in users:
        items.append({
            "id": u.id,
            "email": u.email,
            "name": f"{u.first_name or ''} {u.last_name or ''}".strip() or u.email,
            "enterprise_role": u.enterprise_role,
            "role": u.role.name if u.role else "User",
        })
    return jsonify({"success": True, "users": items})


@admin_bp_v2.route("/api/enterprise-roles/assign", methods=["POST"])
@timed_route
@login_required
@admin_required
def api_assign_enterprise_role():
    """Assign an enterprise role to a user."""
    from app.models.user import VALID_ROLES
    data = request.get_json() or {}
    user_id = data.get("user_id")
    role = data.get("role")
    if not user_id or not role:
        return jsonify({"success": False, "error": "user_id and role are required"}), 400
    if role not in VALID_ROLES:
        return jsonify({"success": False, "error": f"Invalid role. Must be one of: {', '.join(VALID_ROLES)}"}), 400
    user = User.query.get_or_404(user_id)
    user.enterprise_role = role
    db.session.commit()
    return jsonify({"success": True, "user_id": user_id, "role": role})


# =============================================================================
# AUDIT LOG VIEWER (PLT-032)
# =============================================================================


@admin_bp_v2.route("/audit-log")
@login_required
@admin_required
def audit_log_viewer():
    """PLT-032: Admin audit log query UI for SOX/HIPAA compliance."""
    import logging
    from datetime import datetime as dt

    logger = logging.getLogger(__name__)

    page = request.args.get("page", 1, type=int)
    per_page = min(request.args.get("per_page", 50, type=int), 200)
    date_from = request.args.get("date_from", "")
    date_to = request.args.get("date_to", "")
    user_email = request.args.get("user_email", "").strip()
    action_filter = request.args.get("action", "").strip()
    entity_type_filter = request.args.get("entity_type", "").strip()
    search_q = request.args.get("q", "").strip()
    export_csv = request.args.get("export") == "csv"

    entries = []
    total = 0
    action_types = []
    entity_types = []

    try:
        from app.models.audit_log import AuditLog

        query = AuditLog.query.filter_by(is_deleted=False)

        if date_from:
            try:
                query = query.filter(AuditLog.timestamp >= dt.fromisoformat(date_from))
            except ValueError:
                logger.debug("PLT-032: invalid date_from: %s", date_from)

        if date_to:
            try:
                to_dt = dt.fromisoformat(date_to).replace(hour=23, minute=59, second=59)
                query = query.filter(AuditLog.timestamp <= to_dt)
            except ValueError:
                logger.debug("PLT-032: invalid date_to: %s", date_to)

        if user_email:
            query = query.filter(AuditLog.user_email.ilike(f"%{user_email}%"))

        if action_filter:
            query = query.filter(AuditLog.action == action_filter)

        if entity_type_filter:
            query = query.filter(AuditLog.entity_type == entity_type_filter)

        if search_q:
            query = query.filter(
                db.or_(
                    AuditLog.description.ilike(f"%{search_q}%"),
                    AuditLog.entity_name.ilike(f"%{search_q}%"),
                )
            )

        query = query.order_by(AuditLog.timestamp.desc())

        # Get distinct action types and entity types for filter dropdowns
        action_types = [
            r[0] for r in db.session.query(AuditLog.action).distinct().order_by(AuditLog.action).all() if r[0]
        ]
        entity_types = [
            r[0] for r in db.session.query(AuditLog.entity_type).distinct().order_by(AuditLog.entity_type).all() if r[0]
        ]

        if export_csv:
            import csv
            import io

            rows = query.limit(10000).all()
            output = io.StringIO()
            writer = csv.writer(output)
            writer.writerow(["timestamp", "user_email", "action", "entity_type", "entity_id", "entity_name", "description", "status"])
            for row in rows:
                writer.writerow([
                    row.timestamp.isoformat() if row.timestamp else "",
                    row.user_email or "",
                    row.action,
                    row.entity_type,
                    row.entity_id or "",
                    row.entity_name or "",
                    row.description or "",
                    row.status or "",
                ])

            from flask import Response
            return Response(
                output.getvalue(),
                mimetype="text/csv",
                headers={"Content-Disposition": "attachment; filename=audit_log_export.csv"},
            )

        total = query.count()
        entries = query.offset((page - 1) * per_page).limit(per_page).all()

    except Exception as exc:
        logger.warning("PLT-032: AuditLog query failed: %s", exc)

    total_pages = (total + per_page - 1) // per_page if per_page else 1

    if request.accept_mimetypes.best == "application/json":
        return jsonify({
            "entries": [e.to_dict() for e in entries],
            "total": total,
            "page": page,
            "per_page": per_page,
            "total_pages": total_pages,
        })

    return render_template(
        "admin/audit_log.html",
        entries=entries,
        total=total,
        page=page,
        per_page=per_page,
        total_pages=total_pages,
        date_from=date_from,
        date_to=date_to,
        user_email=user_email,
        action_filter=action_filter,
        entity_type_filter=entity_type_filter,
        search_q=search_q,
        action_types=action_types,
        entity_types=entity_types,
    )


# =============================================================================
# APPLICATION REPORT BUILDER (PLT-025)
# =============================================================================


@admin_bp_v2.route("/reports")
@login_required
@admin_required
def report_builder():
    """PLT-025: Self-service parameterized report builder for applications."""
    import logging

    logger = logging.getLogger(__name__)

    page = request.args.get("page", 1, type=int)
    per_page = min(request.args.get("per_page", 50, type=int), 200)
    lifecycle = request.args.get("lifecycle", "").strip()
    vendor_filter = request.args.get("vendor", "").strip()
    search_q = request.args.get("q", "").strip()
    date_from = request.args.get("date_from", "")
    date_to = request.args.get("date_to", "")
    export_csv = request.args.get("export") == "csv"

    apps = []
    total = 0
    lifecycle_options = []
    vendor_options = []

    try:
        from app.models.application_portfolio import ApplicationComponent

        query = ApplicationComponent.query

        if lifecycle:
            query = query.filter(ApplicationComponent.lifecycle_status == lifecycle)

        if vendor_filter:
            query = query.filter(ApplicationComponent.name.ilike(f"%{vendor_filter}%"))

        if search_q:
            query = query.filter(
                db.or_(
                    ApplicationComponent.name.ilike(f"%{search_q}%"),
                    ApplicationComponent.description.ilike(f"%{search_q}%"),
                )
            )

        if date_from:
            from datetime import datetime as dt
            try:
                query = query.filter(ApplicationComponent.created_at >= dt.fromisoformat(date_from))
            except ValueError:
                logger.debug("PLT-025: invalid date_from: %s", date_from)

        if date_to:
            from datetime import datetime as dt
            try:
                to_dt = dt.fromisoformat(date_to).replace(hour=23, minute=59, second=59)
                query = query.filter(ApplicationComponent.created_at <= to_dt)
            except ValueError:
                logger.debug("PLT-025: invalid date_to: %s", date_to)

        query = query.order_by(ApplicationComponent.name)

        # Get filter options
        lifecycle_options = [
            r[0] for r in db.session.query(ApplicationComponent.lifecycle_status).distinct().order_by(
                ApplicationComponent.lifecycle_status
            ).all() if r[0]
        ]
        vendor_options = []  # Populated from VendorOrganization if needed

        if export_csv:
            import csv
            import io

            rows = query.limit(10000).all()
            output = io.StringIO()
            writer = csv.writer(output)
            writer.writerow(["id", "name", "description", "lifecycle_status", "technology_stack", "created_at"])
            for row in rows:
                writer.writerow([
                    row.id,
                    row.name or "",
                    (row.description or "")[:200],
                    row.lifecycle_status or "",
                    row.technology_stack or "",
                    row.created_at.isoformat() if row.created_at else "",
                ])

            from flask import Response
            return Response(
                output.getvalue(),
                mimetype="text/csv",
                headers={"Content-Disposition": "attachment; filename=application_report.csv"},
            )

        total = query.count()
        apps = query.offset((page - 1) * per_page).limit(per_page).all()

    except Exception as exc:
        logger.warning("PLT-025: Application query failed: %s", exc)

    total_pages = (total + per_page - 1) // per_page if per_page else 1

    if request.accept_mimetypes.best == "application/json":
        return jsonify({
            "applications": [a.to_dict() if hasattr(a, 'to_dict') else {"id": a.id, "name": a.name} for a in apps],
            "total": total,
            "page": page,
            "per_page": per_page,
            "total_pages": total_pages,
        })

    return render_template(
        "admin/report_builder.html",
        apps=apps,
        total=total,
        page=page,
        per_page=per_page,
        total_pages=total_pages,
        lifecycle=lifecycle,
        vendor_filter=vendor_filter,
        search_q=search_q,
        date_from=date_from,
        date_to=date_to,
        lifecycle_options=lifecycle_options,
    )


# =============================================================================
# WEBHOOK SETTINGS (PLT-015)
# =============================================================================


@admin_bp_v2.route("/webhook-settings", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
@audit_log("update_webhook_settings")
def webhook_settings():
    """PLT-015: Manage Slack/Teams webhook subscriptions and notification settings."""
    from app.models.webhook import WebhookSubscription

    VALID_EVENTS = [
        "solution.created",
        "solution.updated",
        "solution.approved",
        "solution.rejected",
        "solution.archived",
        "application.created",
        "application.updated",
        "application.retired",
        "arb.submitted",
        "arb.decision",
        "risk.raised",
        "risk.resolved",
    ]

    if request.method == "POST":
        action = request.form.get("action", "create")

        if action == "delete":
            sub_id = request.form.get("subscription_id", "").strip()
            if sub_id:
                sub = WebhookSubscription.query.get(sub_id)
                if sub:
                    sub.is_active = False
                    sub.updated_at = datetime.utcnow()
                    db.session.commit()
                    flash("Webhook subscription deleted.", "success")
                else:
                    flash("Subscription not found.", "error")
            return redirect(url_for("admin.webhook_settings"))

        url = request.form.get("url", "").strip()
        description = request.form.get("description", "").strip()
        webhook_type = request.form.get("webhook_type", "generic").strip()
        secret = request.form.get("secret", "").strip() or None
        selected_events = request.form.getlist("events")

        if not url:
            flash("Webhook URL is required.", "error")
            subscriptions = WebhookSubscription.query.filter_by(is_active=True).order_by(
                WebhookSubscription.created_at.desc()
            ).all()
            return render_template(
                "admin/webhook_settings.html",
                subscriptions=subscriptions,
                valid_events=VALID_EVENTS,
            )

        if webhook_type not in ("generic", "teams", "slack"):
            webhook_type = "generic"

        if not selected_events:
            selected_events = ["*"]

        import uuid as _uuid

        sub = WebhookSubscription(
            id=str(_uuid.uuid4()),
            user_id=str(current_user.id),
            url=url,
            description=description,
            webhook_type=webhook_type,
            events=selected_events,
            secret=secret,
            filters={},
            headers={},
            is_active=True,
            created_at=datetime.utcnow(),
            updated_at=datetime.utcnow(),
        )
        db.session.add(sub)
        db.session.commit()
        flash("Webhook subscription created successfully.", "success")
        return redirect(url_for("admin.webhook_settings"))

    subscriptions = WebhookSubscription.query.filter_by(is_active=True).order_by(
        WebhookSubscription.created_at.desc()
    ).all()
    return render_template(
        "admin/webhook_settings.html",
        subscriptions=subscriptions,
        valid_events=VALID_EVENTS,
    )


@admin_bp_v2.route("/webhook-settings/test/<string:subscription_id>", methods=["POST"])
@timed_route
@login_required
@admin_required
def webhook_settings_test(subscription_id: str):
    """PLT-015: Send a test payload to a webhook subscription."""
    from app.models.webhook import WebhookSubscription
    from app.services.webhook_service import WebhookService

    sub = WebhookSubscription.query.filter_by(id=subscription_id, is_active=True).first()
    if not sub:
        flash("Subscription not found.", "error")
        return redirect(url_for("admin.webhook_settings"))

    try:
        svc = WebhookService()
        result = svc.test_subscription(subscription_id, str(current_user.id))
        if result and result.get("success"):
            flash("Test payload delivered successfully.", "success")
        else:
            flash("Test delivery failed — check the webhook URL and try again.", "error")
    except Exception as exc:
        logger.error("PLT-015: webhook test failed for %s: %s", subscription_id, exc)
        flash(f"Test delivery error: {exc}", "error")

    return redirect(url_for("admin.webhook_settings"))


# =============================================================================
# CONNECTION GAPS REPORT (PLT-010)
# =============================================================================


@admin_bp_v2.route("/connection-gaps")
@timed_route
@login_required
@admin_required
def connection_gaps():
    """PLT-010: Portfolio-wide connection gaps report.

    Shows all solutions sorted by architecture completeness score (ascending),
    with per-junction breakdown, filters, and CSV export.
    """
    import csv
    import io

    from flask import Response
    from app.models.solution_models import Solution

    page = request.args.get("page", 1, type=int)
    per_page = min(request.args.get("per_page", 50, type=int), 200)
    domain_filter = request.args.get("domain", "").strip()
    type_filter = request.args.get("type", "").strip()
    score_min = request.args.get("score_min", None, type=int)
    score_max = request.args.get("score_max", None, type=int)
    export_csv = request.args.get("export") == "csv"

    # Junction category labels for display
    junction_labels = {
        "drivers": "Drivers",
        "goals": "Goals",
        "constraints": "Constraints",
        "requirements": "Requirements",
        "applications": "Applications",
        "archimate_elements": "ArchiMate Elements",
        "vendor_products": "Vendor Products",
        "recommendations": "Recommendations",
        "risks": "Risks",
        "tco_items": "TCO Items",
        "metrics": "Metrics",
        "plateaus": "Plateaus",
        "apqc_processes": "APQC Processes",
        "capability_mappings": "Capability Mappings",
    }

    solutions_data = []
    total = 0
    domain_options = []
    type_options = []
    aggregate_stats = {
        "below_25": 0,
        "below_50": 0,
        "no_phase_a": 0,
        "total": 0,
        "avg_score": 0,
        "junction_gap_counts": {},
    }

    try:
        query = Solution.query

        if domain_filter:
            query = query.filter(Solution.business_domain == domain_filter)
        if type_filter:
            query = query.filter(Solution.solution_type == type_filter)

        # Get filter options from all solutions
        domain_options = [
            r[0] for r in db.session.query(Solution.business_domain)
            .distinct()
            .order_by(Solution.business_domain)
            .all()
            if r[0]
        ]
        type_options = [
            r[0] for r in db.session.query(Solution.solution_type)
            .distinct()
            .order_by(Solution.solution_type)
            .all()
            if r[0]
        ]

        # Fetch all matching solutions and compute scores
        all_solutions = query.all()
        scored = []
        gap_counts = {k: 0 for k in junction_labels}

        for sol in all_solutions:
            try:
                comp = sol.architecture_completeness_score
            except Exception:
                comp = {"score": 0, "filled": [], "missing": list(junction_labels.keys()),
                        "filled_count": 0, "total": 14}

            score = comp.get("score", 0)

            # Apply score range filter
            if score_min is not None and score < score_min:
                continue
            if score_max is not None and score > score_max:
                continue

            missing_set = set(comp.get("missing", []))
            filled_set = set(comp.get("filled", []))

            # Track aggregate gap counts
            for junction_key in junction_labels:
                if junction_key in missing_set:
                    gap_counts[junction_key] = gap_counts.get(junction_key, 0) + 1

            scored.append({
                "id": sol.id,
                "name": sol.name or "(Unnamed)",
                "business_domain": sol.business_domain or "",
                "solution_type": sol.solution_type or "",
                "score": score,
                "filled": filled_set,
                "missing": missing_set,
                "filled_count": comp.get("filled_count", 0),
                "total_junctions": comp.get("total", 14),
            })

        # Sort by score ascending (worst first)
        scored.sort(key=lambda s: s["score"])

        total = len(scored)

        # Aggregate stats
        aggregate_stats["total"] = total
        aggregate_stats["below_25"] = sum(1 for s in scored if s["score"] < 25)
        aggregate_stats["below_50"] = sum(1 for s in scored if s["score"] < 50)
        phase_a_junctions = {"drivers", "goals", "constraints"}
        aggregate_stats["no_phase_a"] = sum(
            1 for s in scored
            if not any(j in s["filled"] for j in phase_a_junctions)
        )
        aggregate_stats["avg_score"] = (
            round(sum(s["score"] for s in scored) / total) if total > 0 else 0
        )
        aggregate_stats["junction_gap_counts"] = gap_counts

        # CSV export (before pagination)
        if export_csv:
            output = io.StringIO()
            writer = csv.writer(output)
            header = ["Solution", "Domain", "Type", "Score (%)"]
            header.extend(junction_labels.values())
            writer.writerow(header)
            for s in scored[:10000]:
                row = [s["name"], s["business_domain"], s["solution_type"], s["score"]]
                for jk in junction_labels:
                    row.append("Yes" if jk in s["filled"] else "No")
                writer.writerow(row)
            return Response(
                output.getvalue(),
                mimetype="text/csv",
                headers={"Content-Disposition": "attachment; filename=connection_gaps_report.csv"},
            )

        # Paginate
        start = (page - 1) * per_page
        solutions_data = scored[start:start + per_page]

    except Exception as exc:
        logger.warning("PLT-010: Connection gaps query failed: %s", exc)

    total_pages = (total + per_page - 1) // per_page if per_page else 1

    return render_template(
        "admin/connection_gaps.html",
        solutions=solutions_data,
        total=total,
        page=page,
        per_page=per_page,
        total_pages=total_pages,
        domain_filter=domain_filter,
        type_filter=type_filter,
        score_min=score_min if score_min is not None else "",
        score_max=score_max if score_max is not None else "",
        domain_options=domain_options,
        type_options=type_options,
        junction_labels=junction_labels,
        aggregate_stats=aggregate_stats,
    )


# ============================================================================
# PLT-024 / PLT-034: Portfolio Export (PDF print page + PPTX download)
# ============================================================================


def _get_portfolio_export_data() -> dict:
    """Gather portfolio data for PDF and PPTX exports.

    Returns a dict with solutions_by_status, top10, bottom10, junction_coverage,
    total_solutions, and avg_completeness.
    """
    from app.models.solution_models import Solution

    solutions = Solution.query.limit(500).all()
    total = len(solutions)
    if total == 0:
        return {
            "total_solutions": 0,
            "avg_completeness": 0,
            "solutions_by_status": {},
            "top10": [],
            "bottom10": [],
            "junction_coverage": {},
            "generated_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        }

    scored = []
    for sol in solutions:
        try:
            cs = sol.architecture_completeness_score
            score = cs["score"]
            filled = cs["filled"]
        except Exception:
            score = 0
            filled = []
        scored.append({
            "name": sol.name or f"Solution #{sol.id}",
            "governance_status": sol.governance_status or "draft",
            "score": score,
            "filled": filled,
        })

    avg_completeness = round(sum(s["score"] for s in scored) / total) if total else 0

    # Solutions by governance status
    status_counts = {}
    for s in scored:
        st = s["governance_status"]
        status_counts[st] = status_counts.get(st, 0) + 1

    sorted_desc = sorted(scored, key=lambda s: s["score"], reverse=True)
    sorted_asc = sorted(scored, key=lambda s: s["score"])

    junction_names = [
        "drivers", "goals", "constraints", "requirements",
        "applications", "archimate_elements", "vendor_products",
        "recommendations", "risks", "tco_items", "metrics",
        "plateaus", "apqc_processes", "capability_mappings",
    ]
    junction_coverage = {}
    for jname in junction_names:
        with_count = sum(1 for s in scored if jname in s["filled"])
        junction_coverage[jname] = {
            "with": with_count,
            "without": total - with_count,
        }

    return {
        "total_solutions": total,
        "avg_completeness": avg_completeness,
        "solutions_by_status": status_counts,
        "top10": [{"name": s["name"], "score": s["score"]} for s in sorted_desc[:10]],
        "bottom10": [{"name": s["name"], "score": s["score"]} for s in sorted_asc[:10]],
        "junction_coverage": junction_coverage,
        "generated_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
    }


@admin_bp_v2.route("/export-portfolio-pdf")
@timed_route
@login_required
@admin_required
def export_portfolio_pdf():
    """PLT-024: Render a print-optimised HTML page for board-ready portfolio PDF.

    Uses window.print() with @media print CSS -- no new dependencies.
    """
    try:
        data = _get_portfolio_export_data()
    except Exception as exc:
        logger.warning("PLT-024: Portfolio PDF data gathering failed: %s", exc)
        data = {
            "total_solutions": 0,
            "avg_completeness": 0,
            "solutions_by_status": {},
            "top10": [],
            "bottom10": [],
            "junction_coverage": {},
            "generated_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        }

    # Render inline HTML (no new template per constraint)
    html_parts = []
    html_parts.append("""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>Architecture Portfolio Summary</title>
<style>
body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
       margin: 2rem; color: #1a1a2e; }
h1 { font-size: 1.75rem; margin-bottom: 0.25rem; }
h2 { font-size: 1.2rem; margin-top: 2rem; border-bottom: 2px solid #3b82f6;
     padding-bottom: 0.25rem; color: #1e3a5f; }
.subtitle { color: #6b7280; font-size: 0.9rem; margin-bottom: 1.5rem; }
.kpi-row { display: flex; gap: 2rem; margin-bottom: 1.5rem; }
.kpi { text-align: center; }
.kpi .value { font-size: 2rem; font-weight: 700; }
.kpi .label { font-size: 0.8rem; color: #6b7280; }
table { width: 100%; border-collapse: collapse; margin-top: 0.5rem; font-size: 0.85rem; }
th, td { border: 1px solid #d1d5db; padding: 0.4rem 0.75rem; text-align: left; }
th { background: #f1f5f9; font-weight: 600; }
tr:nth-child(even) { background: #f9fafb; }
.confidential { color: #9ca3af; font-size: 0.7rem; text-align: center;
                margin-top: 2rem; border-top: 1px solid #e5e7eb; padding-top: 0.5rem; }
.print-btn { background: #3b82f6; color: #fff; border: none; padding: 0.5rem 1.5rem;
             border-radius: 0.375rem; font-size: 0.9rem; cursor: pointer; margin-bottom: 1rem; }
.print-btn:hover { background: #2563eb; }
.back-link { margin-right: 1rem; color: #3b82f6; text-decoration: none; font-size: 0.9rem; }
@media print {
  .no-print { display: none !important; }
  body { margin: 0; font-size: 11pt; }
  h1 { font-size: 16pt; }
  h2 { font-size: 13pt; }
}
</style>
</head>
<body>
<div class="no-print" style="margin-bottom:1rem;">
  <a href="/admin" class="back-link">&larr; Back to Admin</a>
  <button class="print-btn" onclick="window.print()">Print / Save as PDF</button>
</div>
""")

    html_parts.append(f"""<h1>Architecture Portfolio Summary</h1>
<p class="subtitle">Generated {data['generated_at']} &mdash; CONFIDENTIAL</p>

<div class="kpi-row">
  <div class="kpi"><div class="value">{data['total_solutions']}</div><div class="label">Total Solutions</div></div>
  <div class="kpi"><div class="value">{data['avg_completeness']}%</div><div class="label">Avg Completeness</div></div>
  <div class="kpi"><div class="value">{len(data['solutions_by_status'])}</div><div class="label">Status Categories</div></div>
</div>
""")

    # Solutions by status table
    html_parts.append("<h2>Solutions by Governance Status</h2><table><tr><th>Status</th><th>Count</th></tr>")
    if data["solutions_by_status"]:
        for status, count in sorted(data["solutions_by_status"].items()):
            label = status.replace("_", " ").title()
            html_parts.append(f"<tr><td>{label}</td><td>{count}</td></tr>")
    else:
        html_parts.append("<tr><td colspan='2'>No data available</td></tr>")
    html_parts.append("</table>")

    # Top 10
    html_parts.append("<h2>Top 10 Highest Completeness</h2><table><tr><th>#</th><th>Solution</th><th>Score</th></tr>")
    if data["top10"]:
        for i, s in enumerate(data["top10"], 1):
            html_parts.append(f"<tr><td>{i}</td><td>{s['name']}</td><td>{s['score']}%</td></tr>")
    else:
        html_parts.append("<tr><td colspan='3'>No data available</td></tr>")
    html_parts.append("</table>")

    # Bottom 10
    html_parts.append("<h2>Top 10 Gaps (Lowest Completeness)</h2><table><tr><th>#</th><th>Solution</th><th>Score</th></tr>")
    if data["bottom10"]:
        for i, s in enumerate(data["bottom10"], 1):
            html_parts.append(f"<tr><td>{i}</td><td>{s['name']}</td><td>{s['score']}%</td></tr>")
    else:
        html_parts.append("<tr><td colspan='3'>No data available</td></tr>")
    html_parts.append("</table>")

    # Junction coverage
    html_parts.append("<h2>Junction Coverage</h2><table><tr><th>Junction Type</th><th>With Data</th><th>Without Data</th></tr>")
    if data["junction_coverage"]:
        for jname, jdata in data["junction_coverage"].items():
            label = jname.replace("_", " ").title()
            html_parts.append(f"<tr><td>{label}</td><td>{jdata['with']}</td><td>{jdata['without']}</td></tr>")
    else:
        html_parts.append("<tr><td colspan='3'>No data available</td></tr>")
    html_parts.append("</table>")

    html_parts.append('<p class="confidential">CONFIDENTIAL &mdash; For authorised recipients only</p>')
    html_parts.append("</body></html>")

    return "".join(html_parts)


@admin_bp_v2.route("/export-portfolio-pptx")
@timed_route
@login_required
@admin_required
def export_portfolio_pptx():
    """PLT-034: Generate a PowerPoint deck for portfolio summary.

    5 slides: Title, Key Metrics, Top 10, Bottom 10, Junction Coverage.
    """
    import io

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.error("PLT-034: python-pptx not available")
        return jsonify({"error": "PowerPoint export unavailable — python-pptx not installed"}), 500

    try:
        data = _get_portfolio_export_data()
    except Exception as exc:
        logger.warning("PLT-034: Portfolio PPTX data gathering failed: %s", exc)
        data = {
            "total_solutions": 0,
            "avg_completeness": 0,
            "solutions_by_status": {},
            "top10": [],
            "bottom10": [],
            "junction_coverage": {},
            "generated_at": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        }

    BRAND_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
    ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x6B, 0x72, 0x80)
    HEADER_BG = RGBColor(0xF1, 0xF5, 0xF9)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _add_title_footer(slide, footer_text="CONFIDENTIAL"):
        """Add a small footer to a slide."""
        left = Inches(0.5)
        top = Inches(7.0)
        width = Inches(12)
        height = Inches(0.4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(8)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    def _set_cell(table, row, col, text, bold=False, font_size=11, color=None, alignment=None):
        """Set a table cell's text and formatting."""
        cell = table.cell(row, col)
        cell.text = str(text)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold
            if color:
                paragraph.font.color.rgb = color
            if alignment:
                paragraph.alignment = alignment

    # ------------------------------------------------------------------
    # Slide 1: Title
    # ------------------------------------------------------------------
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide1 = prs.slides.add_slide(slide_layout)

    # Title text box
    left = Inches(1)
    top = Inches(2.5)
    txBox = slide1.shapes.add_textbox(left, top, Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Architecture Portfolio Summary"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle / date
    p2 = tf.add_paragraph()
    p2.text = data["generated_at"]
    p2.font.size = Pt(16)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER

    _add_title_footer(slide1)

    # ------------------------------------------------------------------
    # Slide 2: Key Metrics + Solutions by Status
    # ------------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Metrics"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    # KPI row
    kpis = [
        ("Total Solutions", str(data["total_solutions"])),
        ("Avg Completeness", f"{data['avg_completeness']}%"),
    ]
    for i, (label, value) in enumerate(kpis):
        left = Inches(1 + i * 4)
        txBox = slide2.shapes.add_textbox(left, Inches(1.3), Inches(3.5), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    # Status table
    status_items = sorted(data["solutions_by_status"].items()) if data["solutions_by_status"] else [("No data", 0)]
    rows = len(status_items) + 1
    tbl = slide2.shapes.add_table(rows, 2, Inches(1), Inches(3.2), Inches(6), Inches(0.4 * rows)).table
    tbl.columns[0].width = Inches(4)
    tbl.columns[1].width = Inches(2)
    _set_cell(tbl, 0, 0, "Governance Status", bold=True, font_size=12)
    _set_cell(tbl, 0, 1, "Count", bold=True, font_size=12)
    for i, (status, count) in enumerate(status_items, 1):
        label = status.replace("_", " ").title() if isinstance(status, str) else str(status)
        _set_cell(tbl, i, 0, label, font_size=11)
        _set_cell(tbl, i, 1, str(count), font_size=11)

    _add_title_footer(slide2)

    # ------------------------------------------------------------------
    # Slide 3: Top 10 Highest Completeness
    # ------------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Top 10 — Highest Completeness"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    top_items = data["top10"] if data["top10"] else [{"name": "No data available", "score": "-"}]
    rows = len(top_items) + 1
    tbl = slide3.shapes.add_table(rows, 3, Inches(1), Inches(1.3), Inches(10), Inches(0.4 * rows)).table
    tbl.columns[0].width = Inches(1)
    tbl.columns[1].width = Inches(7)
    tbl.columns[2].width = Inches(2)
    _set_cell(tbl, 0, 0, "#", bold=True, font_size=12)
    _set_cell(tbl, 0, 1, "Solution", bold=True, font_size=12)
    _set_cell(tbl, 0, 2, "Score", bold=True, font_size=12)
    for i, s in enumerate(top_items, 1):
        _set_cell(tbl, i, 0, str(i), font_size=11)
        _set_cell(tbl, i, 1, s["name"], font_size=11)
        score_text = f"{s['score']}%" if isinstance(s["score"], (int, float)) else str(s["score"])
        _set_cell(tbl, i, 2, score_text, font_size=11)

    _add_title_footer(slide3)

    # ------------------------------------------------------------------
    # Slide 4: Top 10 Gaps (Lowest Completeness)
    # ------------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Top 10 Gaps — Lowest Completeness"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    bottom_items = data["bottom10"] if data["bottom10"] else [{"name": "No data available", "score": "-"}]
    rows = len(bottom_items) + 1
    tbl = slide4.shapes.add_table(rows, 3, Inches(1), Inches(1.3), Inches(10), Inches(0.4 * rows)).table
    tbl.columns[0].width = Inches(1)
    tbl.columns[1].width = Inches(7)
    tbl.columns[2].width = Inches(2)
    _set_cell(tbl, 0, 0, "#", bold=True, font_size=12)
    _set_cell(tbl, 0, 1, "Solution", bold=True, font_size=12)
    _set_cell(tbl, 0, 2, "Score", bold=True, font_size=12)
    for i, s in enumerate(bottom_items, 1):
        _set_cell(tbl, i, 0, str(i), font_size=11)
        _set_cell(tbl, i, 1, s["name"], font_size=11)
        score_text = f"{s['score']}%" if isinstance(s["score"], (int, float)) else str(s["score"])
        _set_cell(tbl, i, 2, score_text, font_size=11)

    _add_title_footer(slide4)

    # ------------------------------------------------------------------
    # Slide 5: Junction Coverage
    # ------------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Junction Coverage (14 Types)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    jc = data["junction_coverage"]
    jc_items = list(jc.items()) if jc else [("No data", {"with": 0, "without": 0})]
    rows = len(jc_items) + 1
    tbl = slide5.shapes.add_table(rows, 3, Inches(1), Inches(1.3), Inches(10), Inches(0.35 * rows)).table
    tbl.columns[0].width = Inches(5)
    tbl.columns[1].width = Inches(2.5)
    tbl.columns[2].width = Inches(2.5)
    _set_cell(tbl, 0, 0, "Junction Type", bold=True, font_size=12)
    _set_cell(tbl, 0, 1, "Solutions With", bold=True, font_size=12)
    _set_cell(tbl, 0, 2, "Solutions Without", bold=True, font_size=12)
    for i, (jname, jdata) in enumerate(jc_items, 1):
        label = jname.replace("_", " ").title()
        _set_cell(tbl, i, 0, label, font_size=11)
        _set_cell(tbl, i, 1, str(jdata["with"]), font_size=11)
        _set_cell(tbl, i, 2, str(jdata["without"]), font_size=11)

    _add_title_footer(slide5)

    # ------------------------------------------------------------------
    # Serialize and return
    # ------------------------------------------------------------------
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    from flask import send_file
    filename = f"portfolio_summary_{datetime.utcnow().strftime('%Y%m%d')}.pptx"
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename,
    )


# ---------------------------------------------------------------------------
# AI Confidence Calibration (ARC-I02)
# ---------------------------------------------------------------------------
@admin_bp_v2.route("/ai-confidence")
@login_required
@admin_required
@timed_route
def ai_confidence_calibration():
    """AI confidence calibration dashboard — tracks whether AI confidence
    labels match actual architect acceptance rates."""
    from app.models.ai_suggestion import AISuggestion
    from sqlalchemy import func, case

    results = db.session.query(
        case(
            (AISuggestion.confidence >= 0.8, 'High'),
            (AISuggestion.confidence >= 0.5, 'Medium'),
            else_='Low'
        ).label('tier'),
        func.count().label('total'),
        func.sum(case((AISuggestion.architect_verdict == 'accepted', 1), else_=0)).label('accepted'),
        func.sum(case((AISuggestion.architect_verdict == 'modified', 1), else_=0)).label('modified'),
        func.sum(case((AISuggestion.architect_verdict == 'rejected', 1), else_=0)).label('rejected'),
        func.sum(case((AISuggestion.architect_verdict.is_(None), 1), else_=0)).label('no_verdict'),
    ).group_by('tier').all()

    # Ensure all three tiers are present, even with zero rows
    tier_order = ['High', 'Medium', 'Low']
    tier_map = {}
    for row in results:
        accepted = row.accepted or 0
        modified = row.modified or 0
        rejected = row.rejected or 0
        no_verdict = row.no_verdict or 0
        total_with_verdict = accepted + modified + rejected
        rate = round(accepted / total_with_verdict * 100) if total_with_verdict > 0 else 0
        tier_map[row.tier] = {
            'name': row.tier,
            'total': row.total,
            'accepted': accepted,
            'modified': modified,
            'rejected': rejected,
            'no_verdict': no_verdict,
            'acceptance_rate': rate,
        }

    calibration_data = []
    for tier_name in tier_order:
        if tier_name in tier_map:
            calibration_data.append(tier_map[tier_name])
        else:
            calibration_data.append({
                'name': tier_name,
                'total': 0,
                'accepted': 0,
                'modified': 0,
                'rejected': 0,
                'no_verdict': 0,
                'acceptance_rate': 0,
            })

    return render_template(
        'admin/ai_confidence_calibration.html',
        calibration_data=calibration_data,
    )


# ============================================================================
# Solution AI Prompt Management
# ============================================================================

# Registry of default prompts — keys match class attributes on SolutionAIOrchestrator
_SOLUTION_PROMPT_DEFAULTS = None


def _get_prompt_defaults():
    """Lazy-load default prompts from the orchestrator class attributes."""
    global _SOLUTION_PROMPT_DEFAULTS
    if _SOLUTION_PROMPT_DEFAULTS is not None:
        return _SOLUTION_PROMPT_DEFAULTS

    from app.modules.solutions_strategic.v2.services.solution_ai_orchestrator import (
        SolutionAIOrchestrator,
    )

    _SOLUTION_PROMPT_DEFAULTS = {
        "draft_architecture": {
            "name": "Draft Architecture (Phase A)",
            "description": "Generates the complete TOGAF ADM Motivation Layer with all ArchiMate 3.2 element types — stakeholders, drivers, assessments, goals, outcomes, principles, requirements, constraints, values, and risks.",
            "prompt_text": SolutionAIOrchestrator.DRAFT_ARCHITECTURE_PROMPT,
            "category": "Phase A",
            "variables": "solution_name, solution_type, business_domain, complexity, adm_phase, problem_statement, current_state, budget_range, timeline_months, compliance_needs, key_stakeholders, industry_context, technology_preferences, org_context",
        },
        "architecture_variants": {
            "name": "Architecture Variants",
            "description": "Produces 3 alternative architecture options: cost-optimized, timeline-optimized, and risk-balanced, each with entities and trade-offs.",
            "prompt_text": SolutionAIOrchestrator.ARCHITECTURE_VARIANTS_PROMPT,
            "category": "Phase A",
            "variables": "solution_name, business_domain, problem_statement",
        },
        "strategy_specialist": {
            "name": "Strategy Layer (Phase B)",
            "description": "ArchiMate Strategy specialist — generates courses of action, value streams, capability gap analysis, and resources from motivation layer and selected capabilities.",
            "prompt_text": SolutionAIOrchestrator.STRATEGY_SPECIALIST_PROMPT,
            "category": "Phase B",
            "variables": "solution_name, business_domain, problem_statement, stakeholders_json, drivers_json, assessments_json, goals_json, outcomes_json, principles_json, requirements_json, constraints_json, values_json, capability_count, capabilities_json, nfr_checklist",
        },
        "business_specialist": {
            "name": "Business Layer (Phase C)",
            "description": "ArchiMate Business specialist — generates business actors, processes, services, roles, objects, and events from strategy and motivation layers.",
            "prompt_text": SolutionAIOrchestrator.BUSINESS_SPECIALIST_PROMPT,
            "category": "Phase C",
            "variables": "solution_name, business_domain, stakeholders_json, drivers_json, goals_json, principles_json, constraints_json, capabilities_json, courses_of_action_json, existing_business_json, nfr_checklist, advanced_business_schema",
        },
        "application_specialist": {
            "name": "Application Layer (Phase C)",
            "description": "ArchiMate Application specialist — generates application components (matching existing portfolio apps), services, data objects, and interfaces.",
            "prompt_text": SolutionAIOrchestrator.APPLICATION_SPECIALIST_PROMPT,
            "category": "Phase C",
            "variables": "solution_name, business_domain, requirements_json, principles_json, constraints_json, capabilities_json, business_services_json, app_count, existing_apps_json, nfr_checklist",
        },
        "technology_specialist": {
            "name": "Technology Layer (Phase D)",
            "description": "ArchiMate Technology specialist — generates nodes, system software, technology services, networks, and deployment artifacts.",
            "prompt_text": SolutionAIOrchestrator.TECHNOLOGY_SPECIALIST_PROMPT,
            "category": "Phase D",
            "variables": "solution_name, business_domain, tech_preferences, constraints_json, principles_json, app_components_json, existing_infra_json, nfr_checklist",
        },
        "implementation_specialist": {
            "name": "Implementation & Migration (Phase F)",
            "description": "TOGAF Phase F Migration Planning — generates plateaus, gaps, work packages, deliverables, and implementation milestones.",
            "prompt_text": SolutionAIOrchestrator.IMPLEMENTATION_SPECIALIST_PROMPT,
            "category": "Phase F",
            "variables": "solution_name, business_domain, timeline_constraint, budget_constraint, goals_json, stakeholders_json, principles_json, gaps_json, capabilities_json, plateaus_json, nfr_checklist",
        },
        "capability_suggestion": {
            "name": "Capability Suggestions (Step 2)",
            "description": "Maps business problems to APQC capabilities — suggests which capabilities from the catalog a solution should address.",
            "prompt_text": _get_capability_suggestion_default(),
            "category": "Step 2",
            "variables": "solution_description, solution_type, business_domain, capability_catalog",
        },
    }
    return _SOLUTION_PROMPT_DEFAULTS


def _get_capability_suggestion_default():
    """Get the default capability suggestion prompt from SolutionAIService."""
    try:
        from app.modules.ai_chat.services.solution_ai_service import SolutionAIService
        return SolutionAIService.CAPABILITY_SUGGESTION_PROMPT
    except Exception:
        return "(Could not load default prompt)"


def _override_key(prompt_key):
    return f"solution_prompt_{prompt_key}"


@admin_bp_v2.route("/solution-prompts")
@timed_route
@login_required
@admin_required
def solution_prompts_page():
    """Render the solution AI prompt management page."""
    return render_template("admin/solution_prompts.html")


@admin_bp_v2.route("/solution-prompts/data")
@login_required
@admin_required
def solution_prompts_data():
    """JSON API: return all solution prompt configs merged with DB overrides."""
    defaults = _get_prompt_defaults()
    prompts = []

    for key, config in defaults.items():
        override_name = _override_key(key)
        override = AIPromptTemplate.query.filter_by(name=override_name).first()

        prompts.append({
            "key": key,
            "name": config["name"],
            "description": config["description"],
            "category": config["category"],
            "variables": config["variables"],
            "default_prompt": config["prompt_text"],
            "current_prompt": override.system_prompt if override else config["prompt_text"],
            "has_override": override is not None,
        })

    return jsonify({"prompts": prompts})


@admin_bp_v2.route("/solution-prompts/<prompt_key>/update", methods=["POST"])
@login_required
@admin_required
@audit_log("update_solution_prompt")
def solution_prompt_update(prompt_key):
    """Save a custom override for a solution prompt."""
    defaults = _get_prompt_defaults()
    if prompt_key not in defaults:
        return jsonify({"error": f"Unknown prompt: {prompt_key}"}), 404

    payload = request.get_json(silent=True) or {}
    prompt_text = (payload.get("prompt_text") or "").strip()

    if not prompt_text:
        return jsonify({"error": "Prompt text cannot be empty"}), 400

    override_name = _override_key(prompt_key)
    override = AIPromptTemplate.query.filter_by(name=override_name).first()

    if not override:
        override = AIPromptTemplate(
            name=override_name,
            description=defaults[prompt_key]["description"],
            system_prompt=prompt_text,
            user_prompt_template="",
            category="solution_prompt",
        )
        db.session.add(override)
    else:
        override.system_prompt = prompt_text
        override.updated_at = datetime.utcnow()

    try:
        db.session.commit()
        logger.info("Solution prompt override saved for %s by user %s", prompt_key, current_user.id)
    except Exception:
        db.session.rollback()
        logger.exception("Failed to save solution prompt override for %s", prompt_key)
        return jsonify({"error": "Database error saving override"}), 500

    config = defaults[prompt_key]
    return jsonify({
        "success": True,
        "prompt": {
            "key": prompt_key,
            "name": config["name"],
            "description": config["description"],
            "category": config["category"],
            "variables": config["variables"],
            "default_prompt": config["prompt_text"],
            "current_prompt": override.system_prompt,
            "has_override": True,
        },
    })


@admin_bp_v2.route("/solution-prompts/<prompt_key>/reset", methods=["POST"])
@login_required
@admin_required
@audit_log("reset_solution_prompt")
def solution_prompt_reset(prompt_key):
    """Remove custom override, reverting to hardcoded default."""
    defaults = _get_prompt_defaults()
    if prompt_key not in defaults:
        return jsonify({"error": f"Unknown prompt: {prompt_key}"}), 404

    override_name = _override_key(prompt_key)
    override = AIPromptTemplate.query.filter_by(name=override_name).first()

    if override:
        try:
            db.session.delete(override)
            db.session.commit()
            logger.info("Solution prompt override reset for %s by user %s", prompt_key, current_user.id)
        except Exception:
            db.session.rollback()
            logger.exception("Failed to reset solution prompt for %s", prompt_key)
            return jsonify({"error": "Database error resetting prompt"}), 500

    config = defaults[prompt_key]
    return jsonify({
        "success": True,
        "prompt": {
            "key": prompt_key,
            "name": config["name"],
            "description": config["description"],
            "category": config["category"],
            "variables": config["variables"],
            "default_prompt": config["prompt_text"],
            "current_prompt": config["prompt_text"],
            "has_override": False,
        },
    })


# ─── Vendor Pricing Import ──────────────────────────────────────────────
@admin_bp_v2.route("/vendor-pricing/import", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
def vendor_pricing_import():
    """Admin page for importing vendor pricing from CSV/Excel or contracts."""
    staging_items = []
    upload_error = None

    if request.method == "POST":
        file = request.files.get("file")
        if not file or file.filename == "":
            flash("No file selected", "error")
            return render_template("admin/vendor_pricing_import.html",
                                   staging_items=[], upload_error="No file selected")

        filename = secure_filename(file.filename)
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if ext in ("csv", "xlsx", "xls"):
            # Path A: Structured upload — parse with ImportFileParser
            try:
                from app.services.import_core.file_parser import ImportFileParser
                parser = ImportFileParser()
                rows, errors = parser.parse_file(file, filename)
                if errors:
                    upload_error = "; ".join(errors[:5])
                else:
                    required = {"vendor_name", "product_name", "tier", "annual_cost"}
                    validation = parser.validate_file_structure(
                        rows, list(rows[0].keys()) if rows else [],
                        required
                    )
                    if not validation.get("valid"):
                        upload_error = "; ".join(validation.get("errors", [])[:5])
                    else:
                        staging_items = rows
            except Exception as e:
                upload_error = f"Parse error: {str(e)}"

        elif ext in ("pdf", "docx", "txt", "md"):
            # Path B: Document ingestion — LLM extraction
            try:
                from app.modules.vendors.services.pricing_extraction_service import PricingExtractionService
                svc = PricingExtractionService()
                result = svc.extract_from_file(file)
                if result["success"]:
                    staging_items = result["items"]
                else:
                    upload_error = result["error"]
            except Exception as e:
                upload_error = f"Extraction error: {str(e)}"
        else:
            upload_error = f"Unsupported file type: .{ext}. Use CSV, XLSX, PDF, or DOCX."

    return render_template("admin/vendor_pricing_import.html",
                           staging_items=staging_items, upload_error=upload_error)


@admin_bp_v2.route("/vendor-pricing/confirm", methods=["POST"])
@timed_route
@login_required
@admin_required
def vendor_pricing_confirm():
    """Confirm staged pricing items — write to VendorProductPricing as contract_verified."""
    from difflib import SequenceMatcher
    from datetime import datetime, timezone
    from app.models.vendor.vendor_organization import (
        VendorOrganization, VendorProduct, VendorProductPricing,
    )

    data = request.get_json()
    if not data or "items" not in data:
        return jsonify({"error": "No items to confirm"}), 400

    created = 0
    updated = 0
    errors = []

    for item in data["items"]:
        vendor_name = item.get("vendor_name", "").strip()
        product_name = item.get("product_name", "").strip()
        if not vendor_name or not product_name:
            errors.append("Missing vendor/product name in item")
            continue

        # Exact match first, then fuzzy match vendor
        vendor = VendorOrganization.query.filter(
            db.func.lower(VendorOrganization.name) == vendor_name.lower()
        ).first()
        if not vendor:
            all_vendors = VendorOrganization.query.all()
            best_match = None
            best_score = 0
            for v in all_vendors:
                score = SequenceMatcher(None, vendor_name.lower(), v.name.lower()).ratio()
                if score > best_score:
                    best_score = score
                    best_match = v
            if best_match and best_score >= 0.7:
                vendor = best_match
            else:
                errors.append(f"Vendor not found: {vendor_name}")
                continue

        # Exact match first, then fuzzy match product
        product = VendorProduct.query.filter(
            VendorProduct.vendor_organization_id == vendor.id,
            db.func.lower(VendorProduct.name) == product_name.lower()
        ).first()
        if not product:
            products = VendorProduct.query.filter_by(vendor_organization_id=vendor.id).all()
            best_match = None
            best_score = 0
            for p in products:
                score = SequenceMatcher(None, product_name.lower(), p.name.lower()).ratio()
                if score > best_score:
                    best_score = score
                    best_match = p
            if best_match and best_score >= 0.6:
                product = best_match
            else:
                errors.append(f"Product not found: {vendor_name} / {product_name}")
                continue

        tier = item.get("tier", "Standard")
        annual_cost = item.get("annual_cost")

        # Check for existing pricing row (same product + tier)
        existing = VendorProductPricing.query.filter_by(
            product_id=product.id, tier_name=tier
        ).first()

        if existing:
            existing.list_price_annual = annual_cost
            existing.data_source_type = "contract_verified"
            existing.typical_discount_percent = item.get("discount_percent")
            existing.contract_term_months = item.get("contract_term_months", 12)
            existing.setup_fee = item.get("setup_fee")
            existing.implementation_fee = item.get("implementation_fee")
            existing.source = "Contract Import"
            existing.last_verified_at = datetime.now(timezone.utc)
            existing.organization_id = 1
            if item.get("expiry_date"):
                try:
                    from datetime import date
                    existing.expiry_date = date.fromisoformat(item["expiry_date"])
                except (ValueError, TypeError):
                    logger.exception("Failed to operation")
                    pass
            updated += 1
        else:
            pricing = VendorProductPricing(
                product_id=product.id,
                pricing_model=item.get("unit_type", "per_user"),
                tier_name=tier,
                list_price_annual=annual_cost,
                min_users=1,
                max_users=item.get("unit_quantity"),
                currency=item.get("currency", "USD"),
                source="Contract Import",
                billing_frequency="annual",
                contract_term_months=item.get("contract_term_months", 12),
                setup_fee=item.get("setup_fee"),
                implementation_fee=item.get("implementation_fee"),
                typical_discount_percent=item.get("discount_percent"),
                data_source_type="contract_verified",
                last_verified_at=datetime.now(timezone.utc),
                organization_id=1,
                created_by_id=current_user.id,
            )
            if item.get("expiry_date"):
                try:
                    from datetime import date
                    pricing.expiry_date = date.fromisoformat(item["expiry_date"])
                except (ValueError, TypeError):
                    logger.exception("Failed to operation")
                    pass
            db.session.add(pricing)
            created += 1

    db.session.commit()
    return jsonify({
        "success": True,
        "created": created,
        "updated": updated,
        "errors": errors,
    }), 200


# ─── Vendor Pricing Analytics ──────────────────────────────────────────────
@admin_bp_v2.route("/pricing-analytics")
@timed_route
@login_required
@admin_required
def pricing_analytics():
    """Vendor pricing analytics dashboard."""
    from app.modules.vendors.services.confidence_engine import ConfidenceEngine
    engine = ConfidenceEngine()
    analytics = engine.get_analytics()
    return render_template("admin/pricing_analytics.html", analytics=analytics)


# ============================================================================
# Governance Gates Management (GOV-03)
# ============================================================================


@admin_bp_v2.route("/governance-gates")
@timed_route
@login_required
@admin_required
def governance_gates():
    """Governance gates configuration page."""
    from app.modules.solutions_strategic.v2.services.governance_gate_service import (
        DEFAULT_GATES,
    )

    return render_template(
        "admin/governance_gates.html",
        default_gates=DEFAULT_GATES,
    )


@admin_bp_v2.route("/api/governance-gates", methods=["GET"])
@timed_route
@login_required
@admin_required
def governance_gates_list():
    """List all governance gates from DB."""
    from app.models.governance_gates import GovernanceGate

    try:
        gates = GovernanceGate.query.order_by(GovernanceGate.gate_name).all()
        return jsonify({"success": True, "gates": [g.to_dict() for g in gates]})
    except Exception as e:
        logger.error("Error listing governance gates: %s", e, exc_info=True)
        return jsonify({"success": False, "error": "Failed to load gates"}), 500


@admin_bp_v2.route("/api/governance-gates", methods=["POST"])
@timed_route
@login_required
@admin_required
@audit_log("admin_governance_gate_create")
def governance_gates_create():
    """Create a new governance gate."""
    from app.models.governance_gates import GovernanceGate

    data = request.get_json()
    if not data or not data.get("gate_name"):
        return jsonify({"success": False, "error": "gate_name is required"}), 400

    gate_name = data["gate_name"].strip()
    if GovernanceGate.query.filter_by(gate_name=gate_name).first():
        return jsonify({"success": False, "error": f"Gate '{gate_name}' already exists"}), 409

    try:
        gate = GovernanceGate(
            gate_name=gate_name,
            description=data.get("description", ""),
            required_sections=data.get("required_sections", []),
            min_completeness=int(data.get("min_completeness", 60)),
            required_decisions_count=int(data.get("required_decisions_count", 0)),
            require_risk_mitigations=bool(data.get("require_risk_mitigations", False)),
            enabled=bool(data.get("enabled", True)),
        )
        db.session.add(gate)
        db.session.commit()
        return jsonify({"success": True, "gate": gate.to_dict()}), 201
    except Exception as e:
        db.session.rollback()
        logger.error("Error creating governance gate: %s", e, exc_info=True)
        return jsonify({"success": False, "error": "Failed to create gate"}), 500


@admin_bp_v2.route("/api/governance-gates/<int:gate_id>", methods=["PUT"])
@timed_route
@login_required
@admin_required
@audit_log("admin_governance_gate_update")
def governance_gates_update(gate_id):
    """Update an existing governance gate."""
    from app.models.governance_gates import GovernanceGate

    gate = GovernanceGate.query.get_or_404(gate_id)
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "error": "No data provided"}), 400

    try:
        if "gate_name" in data:
            new_name = data["gate_name"].strip()
            existing = GovernanceGate.query.filter_by(gate_name=new_name).first()
            if existing and existing.id != gate_id:
                return jsonify({"success": False, "error": f"Gate '{new_name}' already exists"}), 409
            gate.gate_name = new_name
        if "description" in data:
            gate.description = data["description"]
        if "required_sections" in data:
            gate.required_sections = data["required_sections"]
        if "min_completeness" in data:
            gate.min_completeness = int(data["min_completeness"])
        if "required_decisions_count" in data:
            gate.required_decisions_count = int(data["required_decisions_count"])
        if "require_risk_mitigations" in data:
            gate.require_risk_mitigations = bool(data["require_risk_mitigations"])
        if "enabled" in data:
            gate.enabled = bool(data["enabled"])

        db.session.commit()
        return jsonify({"success": True, "gate": gate.to_dict()})
    except Exception as e:
        db.session.rollback()
        logger.error("Error updating governance gate %s: %s", gate_id, e, exc_info=True)
        return jsonify({"success": False, "error": "Failed to update gate"}), 500


@admin_bp_v2.route("/api/governance-gates/<int:gate_id>", methods=["DELETE"])
@timed_route
@login_required
@admin_required
@audit_log("admin_governance_gate_delete")
def governance_gates_delete(gate_id):
    """Soft-delete a governance gate by disabling it."""
    from app.models.governance_gates import GovernanceGate

    gate = GovernanceGate.query.get_or_404(gate_id)
    try:
        gate.enabled = False
        db.session.commit()
        return jsonify({"success": True, "message": f"Gate '{gate.gate_name}' disabled"})
    except Exception as e:
        db.session.rollback()
        logger.error("Error disabling governance gate %s: %s", gate_id, e, exc_info=True)
        return jsonify({"success": False, "error": "Failed to disable gate"}), 500


# ─── Organization Management (multi-tenancy) ───────────────────────────


@admin_bp_v2.route("/organizations")
@timed_route
@login_required
@platform_admin_required
def organizations_list():
    """List all organizations with user counts."""
    orgs = Organization.query.order_by(Organization.name).all()
    org_data = []
    for org in orgs:
        user_count = User.query.filter_by(organization_id=org.id).count()
        org_data.append({"org": org, "user_count": user_count})
    return render_template("admin/organizations/list.html", organizations=org_data)


@admin_bp_v2.route("/organizations/new", methods=["GET", "POST"])
@timed_route
@login_required
@platform_admin_required
def organization_create():
    """Create a new organization."""
    if request.method == "POST":
        name = request.form.get("name", "").strip()
        slug = request.form.get("slug", "").strip() or name.lower().replace(" ", "-")
        plan = request.form.get("plan", "free")
        max_users = int(request.form.get("max_users") or 10)

        if not name:
            flash("Organization name is required.", "error")
            return render_template("admin/organizations/form.html", org=None)

        if Organization.query.filter_by(slug=slug).first():
            flash(f'An organization with slug "{slug}" already exists.', "error")
            return render_template("admin/organizations/form.html", org=None)

        org = Organization(name=name, slug=slug, plan=plan, max_users=max_users)
        db.session.add(org)
        db.session.commit()
        flash(f'Organization "{name}" created.', "success")
        return redirect(url_for("admin.organizations_list"))

    return render_template("admin/organizations/form.html", org=None)


@admin_bp_v2.route("/organizations/<int:org_id>")
@timed_route
@login_required
@platform_admin_required
def organization_detail(org_id):
    """View organization details and its users."""
    org = Organization.query.get_or_404(org_id)
    users = User.query.filter_by(organization_id=org.id).all()
    return render_template("admin/organizations/detail.html", org=org, users=users)


@admin_bp_v2.route("/organizations/<int:org_id>/edit", methods=["GET", "POST"])
@timed_route
@login_required
@platform_admin_required
def organization_edit(org_id):
    """Edit an existing organization."""
    org = Organization.query.get_or_404(org_id)

    if request.method == "POST":
        org.name = request.form.get("name", "").strip() or org.name
        new_slug = request.form.get("slug", "").strip()
        if new_slug and new_slug != org.slug:
            existing = Organization.query.filter_by(slug=new_slug).first()
            if existing and existing.id != org.id:
                flash(f'Slug "{new_slug}" is already taken.', "error")
                return render_template("admin/organizations/form.html", org=org)
            org.slug = new_slug
        org.plan = request.form.get("plan", org.plan)
        org.max_users = int(request.form.get("max_users") or org.max_users)
        db.session.commit()
        flash(f'Organization "{org.name}" updated.', "success")
        return redirect(url_for("admin.organization_detail", org_id=org.id))

    return render_template("admin/organizations/form.html", org=org)


@admin_bp_v2.route("/organizations/<int:org_id>/toggle", methods=["POST"])
@timed_route
@login_required
@platform_admin_required
def organization_toggle(org_id):
    """Activate or deactivate an organization."""
    org = Organization.query.get_or_404(org_id)
    org.is_active = not org.is_active
    db.session.commit()
    status = "activated" if org.is_active else "deactivated"
    flash(f'Organization "{org.name}" {status}.', "success")
    return redirect(url_for("admin.organization_detail", org_id=org.id))


@admin_bp_v2.route(
    "/organizations/<int:org_id>/users/<int:user_id>/toggle-admin",
    methods=["POST"],
)
@timed_route
@login_required
@platform_admin_required
def toggle_org_admin(org_id, user_id):
    """Toggle org-admin status for a user within an organization."""
    user = User.query.get_or_404(user_id)
    if user.organization_id != org_id:
        flash("User does not belong to this organization.", "error")
        return redirect(url_for("admin.organization_detail", org_id=org_id))
    user.is_org_admin = not user.is_org_admin
    db.session.commit()
    role_label = "granted" if user.is_org_admin else "revoked"
    flash(f'Org-admin role {role_label} for {user.full_name() or user.email}.', "success")
    return redirect(url_for("admin.organization_detail", org_id=org_id))


@admin_bp_v2.route("/organizations/<int:org_id>/delete", methods=["POST"])
@timed_route
@login_required
@platform_admin_required
def organization_delete(org_id):
    """Delete an organization. Moves all users to the Default org first."""
    org = Organization.query.get_or_404(org_id)

    if org.slug == "default":
        flash("Cannot delete the Default organization.", "error")
        return redirect(url_for("admin.organization_detail", org_id=org_id))

    # Move all users to Default org before deleting
    default_org = Organization.query.filter_by(slug="default").first()
    if not default_org:
        flash("Cannot delete — no Default organization to reassign users.", "error")
        return redirect(url_for("admin.organization_detail", org_id=org_id))

    moved = User.query.filter_by(organization_id=org.id).update(
        {"organization_id": default_org.id, "is_org_admin": False},
        synchronize_session=False,
    )

    org_name = org.name
    db.session.delete(org)
    db.session.commit()
    flash(f'Organization "{org_name}" deleted. {moved} user(s) moved to Default.', "success")
    return redirect(url_for("admin.organizations_list"))


@admin_bp_v2.route(
    "/organizations/<int:org_id>/users/<int:user_id>/remove",
    methods=["POST"],
)
@timed_route
@login_required
@platform_admin_required
def remove_user_from_org(org_id, user_id):
    """Remove a user from this organization (moves to Default org)."""
    user = User.query.get_or_404(user_id)
    if user.organization_id != org_id:
        flash("User does not belong to this organization.", "error")
        return redirect(url_for("admin.organization_detail", org_id=org_id))

    default_org = Organization.query.filter_by(slug="default").first()
    if not default_org or default_org.id == org_id:
        flash("Cannot remove user — no other organization to move them to.", "error")
        return redirect(url_for("admin.organization_detail", org_id=org_id))

    user.organization_id = default_org.id
    user.is_org_admin = False
    db.session.commit()
    flash(f'{user.full_name() or user.email} moved to Default organization.', "success")
    return redirect(url_for("admin.organization_detail", org_id=org_id))


# ---------------------------------------------------------------------------
# Power Platform CoE Integration
# ---------------------------------------------------------------------------

_PP_PROVIDER = "power_platform_coe"
_PP_LABEL = "default"

# ── Salesforce Org Discovery (PROG-003) — mirrors the Power Platform block ──
_SF_PROVIDER = "salesforce_org"
_SF_LABEL = "default"


def _sf_settings_row():
    """Return the APISettings row for Salesforce org discovery, or None."""
    from app.models.models import APISettings
    return APISettings.query.filter_by(
        provider=_SF_PROVIDER, key_label=_SF_LABEL
    ).first()


@admin_bp_v2.route("/integrations/salesforce", methods=["GET"])
@timed_route
@login_required
def salesforce_integration():
    """GET /admin/integrations/salesforce — org discovery configuration UI."""
    row = _sf_settings_row()
    sf_config = {}
    if row:
        sf_config = {
            "instance_url": row.jira_url or "",
            "client_id": row.jira_email or "",
            "configured": bool(row.jira_url),
        }
    # Programme picker for import targeting
    from app.modules.solutions_strategic.v2.services.programme_governance_service import (
        ProgrammeGovernanceService,
    )
    programmes = ProgrammeGovernanceService.list_programmes()
    return render_template(
        "admin/integrations_salesforce.html", sf_config=sf_config, programmes=programmes
    )


@admin_bp_v2.route("/integrations/salesforce/save", methods=["POST"])
@timed_route
@login_required
def salesforce_save_credentials():
    """POST /admin/integrations/salesforce/save — persist credentials."""
    from app.modules.solutions_strategic.v2.services.salesforce_discovery_service import (
        SalesforceDiscoveryService,
    )
    data = request.get_json() or request.form
    instance_url = (data.get("instance_url") or "").strip()
    client_id = (data.get("client_id") or "").strip()
    client_secret = (data.get("client_secret") or "").strip()
    if not instance_url or not client_id:
        return jsonify({"status": "error", "error": "instance_url and client_id are required"}), 400
    SalesforceDiscoveryService.save_settings(instance_url, client_id, client_secret)
    return jsonify({"success": True, "status": "ok", "message": "Credentials saved."})


@admin_bp_v2.route("/integrations/salesforce/test", methods=["POST"])
@timed_route
@login_required
def salesforce_test_connection():
    """POST /admin/integrations/salesforce/test — test credentials."""
    from app.modules.solutions_strategic.v2.services.salesforce_discovery_service import (
        SalesforceDiscoveryService,
    )
    data = request.get_json() or request.form
    instance_url = (data.get("instance_url") or "").strip()
    client_id = (data.get("client_id") or "").strip()
    client_secret = (data.get("client_secret") or "").strip()
    if not all([instance_url, client_id, client_secret]):
        return jsonify({"status": "error", "error": "instance_url, client_id, and client_secret are required"}), 400
    result = SalesforceDiscoveryService.test_connection(instance_url, client_id, client_secret)
    return jsonify(result), 200 if result.get("status") == "ok" else 400


@admin_bp_v2.route("/integrations/salesforce/discover", methods=["POST"])
@timed_route
@login_required
def salesforce_discover():
    """POST /admin/integrations/salesforce/discover — list org apps + packages."""
    from app.models.application_portfolio import ApplicationComponent
    from app.modules.solutions_strategic.v2.services.salesforce_discovery_service import (
        SalesforceDiscoveryService,
    )
    row = _sf_settings_row()
    if not row or not row.jira_url:
        return jsonify({"error": "Salesforce credentials not configured. Save credentials first."}), 400
    apps = SalesforceDiscoveryService.discover_apps(
        row.jira_url or "", row.jira_email or "", row.api_key or ""
    )
    linked_ids = {
        r.source_identifier
        for r in ApplicationComponent.query.filter(
            ApplicationComponent.source_identifier.isnot(None)
        ).with_entities(ApplicationComponent.source_identifier).all()
    }
    for app in apps:
        app["status"] = "linked" if app["id"] in linked_ids else "unregistered"
    return jsonify({
        "apps": apps,
        "total": len(apps),
        "ungoverned": sum(1 for a in apps if not a.get("owner_email")),
        "linked": sum(1 for a in apps if a["status"] == "linked"),
    })


@admin_bp_v2.route("/integrations/salesforce/import", methods=["POST"])
@timed_route
@login_required
def salesforce_import():
    """POST /admin/integrations/salesforce/import — import selected app_ids
    with optional programme_initiative_id baseline linkage."""
    from app.modules.solutions_strategic.v2.services.salesforce_discovery_service import (
        SalesforceDiscoveryService,
    )
    data = request.get_json() or request.form
    app_ids = data.get("app_ids") or []
    if isinstance(app_ids, str):
        import json as _json
        try:
            app_ids = _json.loads(app_ids)
        except Exception:
            app_ids = [app_ids]
    if not app_ids:
        return jsonify({"error": "app_ids is required and must be non-empty"}), 400
    row = _sf_settings_row()
    if not row or not row.jira_url:
        return jsonify({"error": "Salesforce credentials not configured."}), 400
    discovered = SalesforceDiscoveryService.discover_apps(
        row.jira_url or "", row.jira_email or "", row.api_key or ""
    )
    result = SalesforceDiscoveryService.import_apps(
        app_ids,
        discovered,
        user_id=current_user.id,
        programme_initiative_id=data.get("programme_initiative_id"),
    )
    return jsonify(result)


def _pp_settings_row():
    """Return the APISettings row for Power Platform CoE, or None."""
    from app.models.models import APISettings
    return APISettings.query.filter_by(
        provider=_PP_PROVIDER, key_label=_PP_LABEL
    ).first()


@admin_bp_v2.route("/integrations/power-platform", methods=["GET"])
@timed_route
@login_required
def power_platform_integration():
    """GET /admin/integrations/power-platform — CoE configuration and discovery UI."""
    row = _pp_settings_row()
    config = {}
    if row:
        config = {
            "tenant_id": row.jira_url or "",
            "client_id": row.jira_email or "",
            "env_url": row.custom_endpoint_url or "",
            "configured": bool(row.jira_url),
        }
    return render_template("admin/integrations_power_platform.html", config=config)


@admin_bp_v2.route("/integrations/power-platform/save", methods=["POST"])
@timed_route
@login_required
def power_platform_save_credentials():
    """POST /admin/integrations/power-platform/save — persist credentials to api_settings."""
    from app.models.models import APISettings
    data = request.get_json() or request.form
    tenant_id = (data.get("tenant_id") or "").strip()
    client_id = (data.get("client_id") or "").strip()
    client_secret = (data.get("client_secret") or "").strip()
    env_url = (data.get("env_url") or "").strip()
    row = _pp_settings_row()
    if not row:
        row = APISettings(provider=_PP_PROVIDER, key_label=_PP_LABEL)
        db.session.add(row)
    row.jira_url = tenant_id
    row.jira_email = client_id
    if client_secret:
        row.api_key = client_secret
    row.custom_endpoint_url = env_url
    db.session.commit()
    return jsonify({"success": True, "status": "ok", "message": "Credentials saved."})


@admin_bp_v2.route("/integrations/power-platform/test", methods=["POST"])
@timed_route
@login_required
def power_platform_test_connection():
    """POST /admin/integrations/power-platform/test — test credentials."""
    from app.modules.solutions_strategic.v2.services.power_platform_coe_service import (
        PowerPlatformCoeService,
    )
    data = request.get_json() or request.form
    tenant_id = (data.get("tenant_id") or "").strip()
    client_id = (data.get("client_id") or "").strip()
    client_secret = (data.get("client_secret") or "").strip()
    env_url = (data.get("env_url") or "").strip()
    if not all([tenant_id, client_id, client_secret]):
        return jsonify({"status": "error", "error": "tenant_id, client_id, and client_secret are required"}), 400
    result = PowerPlatformCoeService.test_connection(tenant_id, client_id, client_secret, env_url)
    http_status = 200 if result.get("status") == "ok" else 400
    return jsonify(result), http_status


@admin_bp_v2.route("/integrations/power-platform/discover", methods=["POST"])
@timed_route
@login_required
def power_platform_discover():
    """POST /admin/integrations/power-platform/discover — trigger discovery, return app list."""
    from app.modules.solutions_strategic.v2.services.power_platform_coe_service import (
        PowerPlatformCoeService,
    )
    row = _pp_settings_row()
    if not row or not row.jira_url:
        return jsonify({"error": "Power Platform credentials not configured. Save credentials first."}), 400
    tenant_id = row.jira_url or ""
    client_id = row.jira_email or ""
    client_secret = row.api_key or ""
    apps = PowerPlatformCoeService.discover_apps(tenant_id, client_id, client_secret)

    # Annotate with ARCHIE link status (same as v1)
    from app.models.application_portfolio import ApplicationComponent
    linked_ids = {
        r.source_identifier
        for r in ApplicationComponent.query.filter(
            ApplicationComponent.source_identifier.isnot(None)
        ).with_entities(ApplicationComponent.source_identifier).all()
    }
    for app in apps:
        app["status"] = "linked" if app["id"] in linked_ids else "unregistered"

    ungoverned = sum(1 for a in apps if not a.get("owner_email"))
    linked = sum(1 for a in apps if a["status"] == "linked")

    return jsonify({"apps": apps, "total": len(apps), "ungoverned": ungoverned, "linked": linked})


@admin_bp_v2.route("/integrations/power-platform/import", methods=["POST"])
@timed_route
@login_required
def power_platform_import():
    """POST /admin/integrations/power-platform/import — import selected app_ids."""
    from app.modules.solutions_strategic.v2.services.power_platform_coe_service import (
        PowerPlatformCoeService,
    )
    data = request.get_json() or request.form
    app_ids = data.get("app_ids") or []
    if isinstance(app_ids, str):
        import json as _json
        try:
            app_ids = _json.loads(app_ids)
        except Exception:
            app_ids = [app_ids]
    if not app_ids:
        return jsonify({"error": "app_ids is required and must be non-empty"}), 400
    row = _pp_settings_row()
    if not row or not row.jira_url:
        return jsonify({"error": "Power Platform credentials not configured."}), 400
    tenant_id = row.jira_url or ""
    client_id = row.jira_email or ""
    client_secret = row.api_key or ""
    discovered = PowerPlatformCoeService.discover_apps(tenant_id, client_id, client_secret)
    from flask_login import current_user
    result = PowerPlatformCoeService.import_apps(app_ids, discovered, current_user.id)
    return jsonify(result)


# ---------------------------------------------------------------------------
# Slack AI Architect Integration
# ---------------------------------------------------------------------------

@admin_bp_v2.route("/integrations/slack", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
def slack_integration():
    """GET/POST /admin/integrations/slack — Slack bot configuration."""
    from app.services.slack_architect_service import SlackArchitectService

    if request.method == "POST":
        bot_token = (request.form.get("bot_token") or "").strip()
        signing_secret = (request.form.get("signing_secret") or "").strip()
        monitored_channels = (request.form.get("monitored_channels") or "").strip()
        portfolio_scan = request.form.get("portfolio_scan") == "1"
        mention_response = request.form.get("mention_response") == "1"
        SlackArchitectService.save_config(
            bot_token, signing_secret, monitored_channels, portfolio_scan, mention_response
        )
        flash("Slack configuration saved.", "success")
        return redirect(url_for("admin.slack_integration"))

    config = SlackArchitectService.get_config()
    return render_template("admin/integrations_slack.html", config=config)


@admin_bp_v2.route("/integrations/slack/test", methods=["POST"])
@timed_route
@login_required
@admin_required
def slack_test_connection():
    """POST /admin/integrations/slack/test — test bot token."""
    from app.services.slack_architect_service import SlackArchitectService

    data = request.get_json() or request.form
    bot_token = (data.get("bot_token") or "").strip()
    if not bot_token:
        cfg = SlackArchitectService.get_config()
        bot_token = cfg.get("bot_token", "")
    if not bot_token:
        return jsonify({"status": "error", "error": "Bot token is required."}), 400
    result = SlackArchitectService.test_connection(bot_token)
    return jsonify(result), 200 if result.get("status") == "ok" else 400


# ---------------------------------------------------------------------------
# Teams Meeting Intelligence
# ---------------------------------------------------------------------------

@admin_bp_v2.route("/integrations/teams-meetings", methods=["GET", "POST"])
@timed_route
@login_required
@admin_required
def teams_meetings_integration():
    """GET/POST /admin/integrations/teams-meetings — meeting transcript config."""
    from app.services.teams_meeting_service import TeamsMeetingService

    if request.method == "POST":
        notification_url = (request.form.get("notification_url") or "").strip()
        transcript_analysis = request.form.get("transcript_analysis") == "1"
        signal_creation = request.form.get("signal_creation") == "1"
        TeamsMeetingService.save_config(notification_url, transcript_analysis, signal_creation)
        flash("Teams Meetings configuration saved.", "success")
        return redirect(url_for("admin.teams_meetings_integration"))

    config = TeamsMeetingService.get_config()
    return render_template("admin/integrations_teams_meetings.html", config=config)


@admin_bp_v2.route("/ai-corrections", methods=["GET"])
@timed_route
@login_required
@admin_required
def ai_corrections():
    """Batch review and correct AI-generated ArchiMate elements."""
    from app.models.application_portfolio import ApplicationComponent

    try:
        from app.models.models import ArchiMateElement
    except ImportError:
        ArchiMateElement = None  # type: ignore[assignment,misc]

    from app.modules.architecture.services.feedback_learning_service import ExtractionFeedback

    # Get apps that have AI-generated ArchiMate elements
    apps_with_elements = (
        ApplicationComponent.query
        .filter(ApplicationComponent.archimate_element_id.isnot(None))
        .order_by(ApplicationComponent.updated_at.desc())
        .limit(50)
        .all()
    )

    rows = []
    for app in apps_with_elements:
        elem = ArchiMateElement.query.get(app.archimate_element_id) if ArchiMateElement else None
        if elem:
            rows.append({
                "app_id": app.id,
                "app_name": app.name,
                "element_id": elem.id,
                "element_name": elem.name,
                "element_type": elem.type,
                "element_layer": elem.layer,
                "completeness": app.completeness_score,
            })

    # Count existing corrections
    correction_count = ExtractionFeedback.query.count()

    return render_template(
        "admin/ai_corrections.html",
        rows=rows,
        correction_count=correction_count,
        total_with_elements=len(rows),
    )


@admin_bp_v2.route("/sap-clean-core", methods=["GET"])
@timed_route
@login_required
@admin_required
def sap_clean_core_dashboard():
    """GET /admin/sap-clean-core — portfolio-level SAP clean-core compliance overview."""
    from app.services.sap_clean_core_service import SAPCleanCoreService
    scan = SAPCleanCoreService.quick_scan_portfolio(limit=30)
    return jsonify(scan)


@admin_bp_v2.route("/sap-clean-core/<int:solution_id>", methods=["GET"])
@timed_route
@login_required
@admin_required
def sap_clean_core_solution(solution_id):
    """GET /admin/sap-clean-core/<id> — validate a single solution."""
    from app.services.sap_clean_core_service import SAPCleanCoreService
    result = SAPCleanCoreService.validate_solution(solution_id)
    return jsonify(result)


@admin_bp_v2.route("/integrations/teams-meetings/subscribe", methods=["POST"])
@timed_route
@login_required
@admin_required
def teams_meetings_subscribe():
    """POST /admin/integrations/teams-meetings/subscribe — register Graph subscription."""
    from app.services.teams_meeting_service import TeamsMeetingService

    data = request.get_json() or request.form
    notification_url = (data.get("notification_url") or "").strip()
    if not notification_url:
        cfg = TeamsMeetingService.get_config()
        notification_url = cfg.get("notification_url", "")
    if not notification_url:
        return jsonify({"status": "error", "error": "notification_url is required."}), 400
    result = TeamsMeetingService.subscribe_to_call_records(notification_url)
    return jsonify(result), 200 if result.get("status") == "ok" else 400
